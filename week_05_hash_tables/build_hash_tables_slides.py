"""
Generate a university-grade PowerPoint deck for Week 05: Hash Tables.

Output: hash_tables.pptx (16:9 widescreen)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# -------- Theme --------
NAVY        = RGBColor(0x0B, 0x2A, 0x4A)
TEAL        = RGBColor(0x12, 0x7A, 0x8A)
GOLD        = RGBColor(0xE0, 0xA8, 0x1F)
LIGHT_GREY  = RGBColor(0xF2, 0xF4, 0xF7)
DARK_GREY   = RGBColor(0x33, 0x3D, 0x4A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xC0, 0x39, 0x2B)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
PURPLE      = RGBColor(0x6A, 0x1B, 0x9A)
CODE_BG     = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG     = RGBColor(0xE6, 0xE6, 0xE6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ---------- Generic helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=18, color=DARK_GREY,
                bullet_color=TEAL, line_spacing=1.25, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        rb = p.add_run(); rb.text = "▸  "
        rb.font.name = font; rb.font.size = Pt(size)
        rb.font.bold = True; rb.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run(); r1.text = head
            r1.font.name = font; r1.font.size = Pt(size)
            r1.font.bold = True; r1.font.color.rgb = NAVY
            r2 = p.add_run(); r2.text = tail
            r2.font.name = font; r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = item
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb

def add_header(slide, title, subtitle=None, page_num=None, total=None):
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    add_rect(slide, Inches(0.35), 0, SLIDE_W - Inches(0.35), Inches(0.08), GOLD)
    add_text(slide, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.78), Inches(11.5), Inches(0.45),
                 subtitle, size=15, italic=True, color=TEAL)
    add_rect(slide, Inches(0.6), Inches(1.22), Inches(2.0), Inches(0.04), GOLD)
    add_text(slide, Inches(0.6), SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
             "Data Structures & Algorithms · Week 05 · Hash Tables",
             size=10, color=DARK_GREY, italic=True)
    if page_num is not None and total is not None:
        add_text(slide, SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.45),
                 Inches(1.2), Inches(0.3),
                 f"Slide {page_num} / {total}", size=10, color=DARK_GREY,
                 align=PP_ALIGN.RIGHT)

def add_code_block(slide, x, y, w, h, code, size=14):
    add_rect(slide, x, y, w, h, CODE_BG)
    add_rect(slide, x, y, Inches(0.08), h, GOLD)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.25), h - Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.line_spacing = 1.15
        r = p.add_run(); r.text = ln if ln else " "
        r.font.name = "Consolas"; r.font.size = Pt(size)
        r.font.color.rgb = CODE_FG

# ---------- Slides ----------
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(5.5), SLIDE_W, Inches(0.12), GOLD)
    add_rect(s, 0, Inches(5.7), SLIDE_W, Inches(0.04), TEAL)
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.9), Inches(1.0),
                              Inches(2.4), Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb = GOLD
    chip.line.fill.background(); chip.shadow.inherit = False
    add_text(s, Inches(0.9), Inches(1.0), Inches(2.4), Inches(0.55),
             "WEEK 05", size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(1.4),
             "Hash Tables & Dictionaries", size=58, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.7),
             "Achieving O(1) Lookups through Mathematical Hashing",
             size=24, color=GOLD, italic=True)
    add_text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.6),
             "A University-Grade Lecture in Data Structures & Algorithms",
             size=18, color=LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(6.0), Inches(8), Inches(0.4),
             "Course   ·   CSC: Data Structures & Algorithms",
             size=14, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.4), Inches(8), Inches(0.4),
             "Module   ·   Week 05 of 12",
             size=14, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.8), Inches(8), Inches(0.4),
             "Format   ·   Lecture · Demonstration · Practical Lab",
             size=14, color=WHITE)

def slide_outline(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Lecture Outline",
               "The magic behind O(1) performance", n, total)
    items = [
        ("1. The Need for Speed.  ", "Why arrays and linked lists are not enough for fast lookups."),
        ("2. Direct Address Tables.  ", "The intuition of using keys as array indices."),
        ("3. Hash Functions.  ", "Mapping arbitrary data (strings, objects) to integer indices."),
        ("4. The Collision Problem.  ", "What happens when two keys hash to the same bucket?"),
        ("5. Collision Resolution: Chaining.  ", "Linked lists inside arrays."),
        ("6. Collision Resolution: Open Addressing.  ", "Linear probing, quadratic probing, and double hashing."),
        ("7. The Load Factor.  ", "When and why to resize your hash table."),
        ("8. Hash Tables in Practice.  ", "Python dicts, Java HashMaps, and real-world applications.")
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(12), Inches(5.5),
                items, size=18, line_spacing=1.32)

def slide_motivation(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "1. The Need for Speed",
               "The limitations of sequential search", n, total)
               
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
             "Searching for a value in traditional linear data structures takes O(N) time.", size=20, bold=True, color=NAVY)
             
    add_text(s, Inches(0.9), Inches(2.2), Inches(5.5), Inches(4.5),
             "Array:\n"
             "Finding an element by value requires checking every element until you find it (Linear Search).\n\n"
             "Linked List:\n"
             "Requires traversing pointers from the head until the node is found.\n\n"
             "If we have 1 billion records, O(N) is too slow for real-time systems (e.g., Database Lookups, Caching, DNS Resolution).",
             size=16, color=DARK_GREY)
             
    add_rect(s, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.0), LIGHT_GREY)
    add_text(s, Inches(7.0), Inches(2.3), Inches(5.5), Inches(0.5), "The Goal: O(1) Time", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.2), Inches(2.9), Inches(5.1), Inches(3.0),
             "We want to instantly know exactly where our data is stored in memory, without scanning.\n\n"
             "To do this, we need a mathematical function that takes our 'Key' and directly computes its 'Index' in an array.",
             size=16, color=DARK_GREY)

def slide_hash_function(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "2. The Hash Function",
               "Translating Keys into Array Indices", n, total)
               
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.0),
             "A Hash Function takes an arbitrary input (like a string) and outputs a fixed-size integer. "
             "We then use the modulo operator (%) to constrain this integer to the size of our array.",
             size=18, bold=True, color=NAVY)
             
    add_rect(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.0), TEAL)
    add_text(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(0.5),
             "Index = HashFunction(Key) % ArraySize", size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
             
    add_text(s, Inches(0.9), Inches(4.0), Inches(6.0), Inches(0.5), "Properties of a Good Hash Function:", size=18, bold=True, color=NAVY)
    items = [
        "Deterministic: Same input always yields the same output.",
        "Fast to Compute: Shouldn't take O(N) time itself.",
        "Uniform Distribution: Should spread keys evenly across the array to avoid clusters."
    ]
    add_bullets(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(2.5), items, size=16, line_spacing=1.3)

def slide_collisions(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "3. The Collision Problem",
               "What happens when two keys share an index?", n, total)
               
    add_text(s, Inches(0.9), Inches(1.5), Inches(6.0), Inches(2.0),
             "Because the number of possible keys (e.g., all possible strings) is infinite, but our array size is finite, "
             "collisions are mathematically inevitable (Pigeonhole Principle).\n\n"
             "A collision occurs when:\n"
             "Hash(Key1) == Hash(Key2)",
             size=16, color=DARK_GREY)
             
    # Visualization
    add_rect(s, Inches(7.5), Inches(1.5), Inches(5.0), Inches(4.5), LIGHT_GREY)
    add_text(s, Inches(7.5), Inches(1.6), Inches(5.0), Inches(0.5), "Collision Example", size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
    
    add_text(s, Inches(7.7), Inches(2.5), Inches(2.0), Inches(0.5), '"Alice"', size=16, bold=True, color=NAVY)
    add_text(s, Inches(7.7), Inches(3.5), Inches(2.0), Inches(0.5), '"Bob"', size=16, bold=True, color=NAVY)
    
    add_rect(s, Inches(10.5), Inches(2.8), Inches(1.5), Inches(0.8), GOLD)
    add_text(s, Inches(10.5), Inches(2.8), Inches(1.5), Inches(0.8), "Index 4", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    
    # Arrows
    add_text(s, Inches(9.2), Inches(2.5), Inches(1.0), Inches(0.5), "----->", size=16)
    add_text(s, Inches(9.2), Inches(3.5), Inches(1.0), Inches(0.5), "--/-->", size=16)
    
    add_text(s, Inches(0.9), Inches(4.5), Inches(6.0), Inches(1.0),
             "We must have a strategy to resolve these collisions.\n"
             "Two primary strategies exist: Chaining and Open Addressing.", size=16, bold=True, color=NAVY)

def slide_chaining(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "4. Resolution: Separate Chaining",
               "Linked lists inside arrays", n, total)
               
    add_text(s, Inches(0.9), Inches(1.5), Inches(5.5), Inches(4.5),
             "Instead of storing the value directly in the array bucket, we store a pointer to a Linked List.\n\n"
             "Insertion:\n"
             "Hash the key, go to the index, and append the key-value pair to the linked list. O(1).\n\n"
             "Lookup:\n"
             "Hash the key, go to the index, and traverse the linked list to find the key. O(L) where L is the list length.\n\n"
             "Pros: Simple, handles infinite collisions easily.\n"
             "Cons: Memory overhead for pointers, poor cache locality.",
             size=15, color=DARK_GREY)
             
    add_rect(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(4.5), LIGHT_GREY)
    add_text(s, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5), "Chaining Visualization", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    
    # Array indices
    for i in range(4):
        add_rect(s, Inches(7.5), Inches(2.2 + i*0.8), Inches(0.8), Inches(0.6), GOLD)
        add_text(s, Inches(7.5), Inches(2.2 + i*0.8), Inches(0.8), Inches(0.6), f"[{i}]", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        
    # Lists
    add_rect(s, Inches(9.0), Inches(2.2), Inches(1.5), Inches(0.5), WHITE, line=DARK_GREY)
    add_text(s, Inches(9.0), Inches(2.2), Inches(1.5), Inches(0.5), '"Alice"', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    
    add_rect(s, Inches(9.0), Inches(3.8), Inches(1.5), Inches(0.5), WHITE, line=DARK_GREY)
    add_text(s, Inches(9.0), Inches(3.8), Inches(1.5), Inches(0.5), '"Bob"', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(11.0), Inches(3.8), Inches(1.5), Inches(0.5), WHITE, line=DARK_GREY)
    add_text(s, Inches(11.0), Inches(3.8), Inches(1.5), Inches(0.5), '"Charlie"', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def slide_open_addressing(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "5. Resolution: Open Addressing",
               "Finding the next available slot", n, total)
               
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.0),
             "In Open Addressing, all elements are stored directly in the array. No linked lists. "
             "If a collision occurs, we probe (search) for the next empty bucket.",
             size=16, bold=True, color=NAVY)
             
    add_text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.5), "Probing Strategies:", size=16, bold=True, color=TEAL)
    
    items = [
        ("Linear Probing: ", "Check index + 1, index + 2, index + 3... (Causes 'Primary Clustering')."),
        ("Quadratic Probing: ", "Check index + 1^2, index + 2^2, index + 3^2... (Avoids primary clustering)."),
        ("Double Hashing: ", "Use a second hash function to determine the step size. (Best distribution).")
    ]
    add_bullets(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(2.0), items, size=15)
    
    add_rect(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.5), LIGHT_GREY)
    add_text(s, Inches(1.1), Inches(5.3), Inches(11.0), Inches(1.3),
             "Important: When deleting a key in Open Addressing, you cannot just set the bucket to Null. "
             "This would break the probing chain for other elements! Instead, we place a special 'Tombstone' marker.",
             size=15, color=DARK_GREY, italic=True)

def slide_load_factor(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "6. Load Factor & Resizing",
               "Keeping the hash table fast", n, total)
               
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
             "Load Factor (α) = Number of Entries / Array Size", size=22, bold=True, color=RED)
             
    add_text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.5),
             "As the table fills up, collisions become more frequent, and O(1) operations degrade to O(N).\n\n"
             "To maintain O(1) performance, we must Resize (Rehash) the table when the Load Factor crosses a threshold (usually 0.7 or 70%).",
             size=16, color=DARK_GREY)
             
    add_rect(s, Inches(0.9), Inches(3.8), Inches(11.5), Inches(2.5), CODE_BG)
    add_text(s, Inches(1.1), Inches(4.0), Inches(11.0), Inches(0.5), "Rehashing Process:", size=16, bold=True, color=GOLD)
    items = [
        "Create a new array double the size of the old array.",
        "Iterate through the old array.",
        "Re-calculate the hash (because ArraySize changed!) and insert into new array.",
        "Takes O(N) time, but occurs infrequently (Amortized O(1))."
    ]
    # Quick custom bullet logic for code bg
    for i, it in enumerate(items):
        add_text(s, Inches(1.3), Inches(4.5 + i*0.4), Inches(10.5), Inches(0.4), f"• {it}", size=14, color=WHITE)

def slide_implementation(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "7. Simple Implementation (Chaining)",
               "Building a Hash Map from scratch", n, total)
               
    code = (
        "class HashTable:\n"
        "    def __init__(self, size=10):\n"
        "        self.size = size\n"
        "        self.table = [ [] for _ in range(self.size) ] # Array of lists\n"
        "\n"
        "    def _hash(self, key):\n"
        "        return hash(key) % self.size\n"
        "\n"
        "    def insert(self, key, value):\n"
        "        idx = self._hash(key)\n"
        "        for kvp in self.table[idx]:\n"
        "            if kvp[0] == key:\n"
        "                kvp[1] = value # Update existing\n"
        "                return\n"
        "        self.table[idx].append([key, value]) # Append to chain\n"
        "\n"
        "    def get(self, key):\n"
        "        idx = self._hash(key)\n"
        "        for kvp in self.table[idx]:\n"
        "            if kvp[0] == key: return kvp[1]\n"
        "        return None"
    )
    add_code_block(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.0), code, size=13)

def slide_summary(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Summary & Time Complexity",
               "The ultimate trade-off: Memory for Speed", n, total)
               
    # Table headers
    add_rect(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.6), NAVY)
    add_text(s, Inches(1.0), Inches(1.5), Inches(3.0), Inches(0.6), "Operation", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.5), Inches(1.5), Inches(3.5), Inches(0.6), "Average Case", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.5), Inches(1.5), Inches(3.5), Inches(0.6), "Worst Case", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    
    rows = [
        ("Search / Lookup", "O(1)", "O(N) (If all hash to same bucket)"),
        ("Insert", "O(1)", "O(N)"),
        ("Delete", "O(1)", "O(N)")
    ]
    
    for i, (op, avg, wst) in enumerate(rows):
        y = 2.1 + i*0.6
        bg = LIGHT_GREY if i%2==0 else WHITE
        add_rect(s, Inches(0.9), Inches(y), Inches(11.5), Inches(0.6), bg)
        add_text(s, Inches(1.0), Inches(y), Inches(3.0), Inches(0.6), op, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.5), Inches(y), Inches(3.5), Inches(0.6), avg, size=15, color=GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.5), Inches(y), Inches(3.5), Inches(0.6), wst, size=15, color=RED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        
    add_text(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.5), "Key Takeaways:", size=18, bold=True, color=NAVY)
    items = [
        "Hash Tables sacrifice memory (Array size must be > Number of Items) to gain O(1) Speed.",
        "Data is Unordered. If you need sorted keys, use a Binary Search Tree (O(log N)).",
        "Python dicts, Sets, and Caches use Hash Tables under the hood."
    ]
    add_bullets(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(2.0), items, size=15, line_spacing=1.3)

# --- Build ---
slides = [
    slide_title,
    slide_outline,
    slide_motivation,
    slide_hash_function,
    slide_collisions,
    slide_chaining,
    slide_open_addressing,
    slide_load_factor,
    slide_implementation,
    slide_summary
]

for i, f in enumerate(slides):
    if f == slide_title:
        f()
    else:
        f(i, len(slides)-1)

prs.save("hash_tables.pptx")
print("Saved hash_tables.pptx")
