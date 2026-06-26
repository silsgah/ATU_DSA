"""
Generate a university-grade PowerPoint deck for Week 07: Heaps & Priority Queues.

Output: heaps_dsa_final.pptx  (16:9 widescreen, same visual language as Week 06)

Run:
    python3 build_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ──────────────────────────────────────────────────────────────────────────────
# Colour palette  (same theme as Week 06)
# ──────────────────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0B, 0x2A, 0x4A)
TEAL       = RGBColor(0x12, 0x7A, 0x8A)
GOLD       = RGBColor(0xE0, 0xA8, 0x1F)
LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF7)
DARK_GREY  = RGBColor(0x33, 0x3D, 0x4A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
RED        = RGBColor(0xC0, 0x39, 0x2B)
GREEN      = RGBColor(0x27, 0x7B, 0x3E)
PURPLE     = RGBColor(0x6A, 0x1B, 0x9A)
CODE_BG    = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG    = RGBColor(0xE6, 0xE6, 0xE6)
ORANGE     = RGBColor(0xD4, 0x73, 0x14)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

TOTAL_SLIDES = 15   # updated after counting below


# ──────────────────────────────────────────────────────────────────────────────
# Generic drawing helpers
# ──────────────────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, line_color=None, line_width=Pt(1)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = line_width
    shp.shadow.inherit = False
    return shp


def add_rounded_rect(slide, x, y, w, h, fill, line_color=None, radius=Pt(8)):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    # adjust corner radius
    adj = shp.adjustments
    if adj:
        adj[0] = 0.05
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=DARK_GREY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.05)
    tf.margin_right  = Inches(0.05)
    tf.margin_top    = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name  = font
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_paragraph(tf, text, *, size=17, bold=False, color=DARK_GREY,
                  font="Calibri", italic=False, align=PP_ALIGN.LEFT,
                  space_before=Pt(0), space_after=Pt(5)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.space_after  = space_after
    r = p.add_run()
    r.text = text
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p


def add_bullets(slide, x, y, w, h, items, *, size=17, color=DARK_GREY,
                bullet_color=TEAL, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.05)
    tf.margin_right  = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment   = PP_ALIGN.LEFT
        p.line_spacing = 1.3
        p.space_after  = Pt(5)
        rb = p.add_run()
        rb.text = "▸  "
        rb.font.name  = font
        rb.font.size  = Pt(size)
        rb.font.bold  = True
        rb.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run(); r1.text = head
            r1.font.name = font; r1.font.size = Pt(size)
            r1.font.bold = True; r1.font.color.rgb = NAVY
            r2 = p.add_run(); r2.text = tail
            r2.font.name = font; r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = item
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def add_code_block(slide, x, y, w, h, code_lines, *, size=13):
    add_rect(slide, x, y, w, h, CODE_BG)
    add_rect(slide, x, y, Inches(0.07), h, GOLD)
    tb = slide.shapes.add_textbox(
        x + Inches(0.15), y + Inches(0.1),
        w - Inches(0.25), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    for i, ln in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment   = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = ln if ln else " "
        r.font.name  = "Consolas"
        r.font.size  = Pt(size)
        r.font.color.rgb = CODE_FG


def add_edge(slide, x1, y1, x2, y2, color=DARK_GREY, weight=1.75):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    return conn


def add_tree_node(slide, cx, cy, value, *, r=Inches(0.38),
                  fill=WHITE, border=NAVY, text_color=NAVY,
                  size=15, highlight=False):
    if highlight:
        fill = GOLD
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(2)
    shp.shadow.inherit = False
    add_text(slide, cx - r, cy - r, r * 2, r * 2,
             str(value), size=size, bold=True, color=text_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def draw_heap_tree(slide, positions, edges, *, node_r=Inches(0.38),
                   fill=WHITE, border=NAVY, text_color=NAVY,
                   highlight_nodes=None, font_size=15):
    """Draw a labelled heap tree from a positions dict and edges list."""
    highlight_nodes = highlight_nodes or set()
    for (a, b) in edges:
        x1, y1 = positions[a]
        x2, y2 = positions[b]
        add_edge(slide, x1, y1 + node_r, x2, y2 - node_r,
                 color=DARK_GREY, weight=1.75)
    for label, (cx, cy) in positions.items():
        hl = label in highlight_nodes
        add_tree_node(slide, cx, cy, label, r=node_r,
                      fill=fill, border=border, text_color=text_color,
                      size=font_size, highlight=hl)


def add_array_cells(slide, x, y, values, *, cell_w=Inches(0.55),
                    cell_h=Inches(0.55), highlight_indices=None,
                    root_index=0, font_size=15):
    """Draw an array visualisation with index labels below each cell."""
    highlight_indices = set(highlight_indices or [])
    for i, val in enumerate(values):
        cx = x + i * (cell_w + Inches(0.04))
        if i == root_index:
            bg = TEAL
            tc = WHITE
        elif i in highlight_indices:
            bg = GOLD
            tc = NAVY
        else:
            bg = LIGHT_GREY
            tc = NAVY
        add_rect(slide, cx, y, cell_w, cell_h, bg,
                 line_color=NAVY, line_width=Pt(1))
        add_text(slide, cx, y, cell_w, cell_h,
                 str(val), size=font_size, bold=True, color=tc,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # index label
        add_text(slide, cx, y + cell_h + Inches(0.04),
                 cell_w, Inches(0.25),
                 f"[{i}]", size=10, color=DARK_GREY,
                 align=PP_ALIGN.CENTER)


def add_header(slide, title, subtitle=None, page_num=None, total=None):
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    add_rect(slide, Inches(0.35), 0,
             SLIDE_W - Inches(0.35), Inches(0.08), GOLD)
    add_text(slide, Inches(0.6), Inches(0.18),
             Inches(11.5), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.78),
                 Inches(11.5), Inches(0.45),
                 subtitle, size=15, italic=True, color=TEAL)
    add_rect(slide, Inches(0.6), Inches(1.22),
             Inches(2.0), Inches(0.04), GOLD)
    add_text(slide, Inches(0.6), SLIDE_H - Inches(0.45),
             Inches(9), Inches(0.3),
             "Data Structures & Algorithms  ·  Week 07  ·  Heaps & Priority Queues",
             size=10, color=DARK_GREY, italic=True)
    if page_num is not None and total is not None:
        add_text(slide, SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.45),
                 Inches(1.2), Inches(0.3),
                 f"Slide {page_num} / {total}",
                 size=10, color=DARK_GREY, align=PP_ALIGN.RIGHT)


# ──────────────────────────────────────────────────────────────────────────────
# Individual slides
# ──────────────────────────────────────────────────────────────────────────────

def slide_cover():
    s = prs.slides.add_slide(BLANK)
    # Full navy background
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # Gold accent bar (left)
    add_rect(s, 0, 0, Inches(0.5), SLIDE_H, GOLD)
    # Gold accent bar (bottom)
    add_rect(s, 0, SLIDE_H - Inches(0.1), SLIDE_W, Inches(0.1), GOLD)

    # Week chip
    add_rounded_rect(s, Inches(0.8), Inches(1.2),
                     Inches(3.5), Inches(0.5),
                     TEAL, line_color=None)
    add_text(s, Inches(0.8), Inches(1.2), Inches(3.5), Inches(0.5),
             "📅  Week 07  ·  Data Structures & Algorithms",
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Main title
    add_text(s, Inches(0.8), Inches(1.95), Inches(9.5), Inches(1.6),
             "Heaps & Priority Queues",
             size=52, bold=True, color=WHITE)

    # Gold underline
    add_rect(s, Inches(0.8), Inches(3.55), Inches(5.5), Inches(0.07), GOLD)

    # Tagline
    add_text(s, Inches(0.8), Inches(3.75), Inches(8.5), Inches(0.9),
             "Efficiently locating the minimum or maximum element — "
             "from the heap property to real-world scheduling algorithms.",
             size=17, color=LIGHT_GREY, italic=True)

    # Topic pills
    topics = [
        ("Binary Heap", TEAL),
        ("Min-Heap / Max-Heap", GREEN),
        ("Priority Queue", GOLD),
        ("Sift Operations", ORANGE),
        ("K-th Largest Element", PURPLE),
    ]
    px = Inches(0.8)
    py = Inches(4.85)
    for label, col in topics:
        w = Inches(len(label) * 0.14 + 0.6)
        add_rounded_rect(s, px, py, w, Inches(0.42), col)
        add_text(s, px, py, w, Inches(0.42),
                 label, size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        px += w + Inches(0.18)

    # Institution
    add_text(s, Inches(0.8), SLIDE_H - Inches(0.9),
             Inches(6), Inches(0.35),
             "ATU  ·  Computer Science  ·  2024 – 2025",
             size=13, color=LIGHT_GREY, italic=True)


def slide_agenda(page):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GREY)
    add_header(s, "Agenda", "What we cover this week", page, TOTAL_SLIDES)

    agenda = [
        "1.  What is a Binary Heap?",
        "2.  The Heap Property — Min-Heap vs. Max-Heap",
        "3.  Storing a Heap in a Plain Array",
        "4.  Sift-Up  (restoring the heap after insertion)",
        "5.  Sift-Down  (restoring the heap after extraction)",
        "6.  Worked Insert & Extract Examples",
        "7.  Priority Queues and Real-World Applications",
        "8.  Python's heapq Module",
        "9.  K-th Largest Element Problem",
        "10. Floyd's O(N) Linear-Time Heap Construction",
    ]

    tb = slide.shapes.add_textbox if False else None  # noqa
    bx = Inches(0.8)
    by = Inches(1.5)
    bw = Inches(11.5)
    tb = s.shapes.add_textbox(bx, by, bw, Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(agenda):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.3
        p.space_after = Pt(4)
        rb = p.add_run(); rb.text = "▸  "
        rb.font.name = "Calibri"; rb.font.size = Pt(18)
        rb.font.bold = True; rb.font.color.rgb = TEAL
        r = p.add_run(); r.text = line
        r.font.name = "Calibri"; r.font.size = Pt(18)
        r.font.color.rgb = DARK_GREY


def slide_what_is_heap(page):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GREY)
    add_header(s, "What Is a Binary Heap?", "Core concept", page, TOTAL_SLIDES)

    # Left panel — description
    lx, ly, lw = Inches(0.6), Inches(1.5), Inches(6.4)
    add_text(s, lx, ly, lw, Inches(0.5),
             "Definition",
             size=19, bold=True, color=NAVY)
    add_rect(s, lx, ly + Inches(0.48), Inches(1.6), Inches(0.04), TEAL)

    desc = (
        "A Binary Heap is a special kind of complete binary tree designed "
        "to answer one specific question very efficiently:\n\n"
        "\u201cWhat is the absolute largest or smallest element in this dataset?\u201d\n\n"
        "It answers that question in O(1) time. Inserting or removing that "
        "top element while maintaining the heap structure takes O(log N) time."
    )
    add_text(s, lx, ly + Inches(0.65), lw, Inches(2.0),
             desc, size=16, color=DARK_GREY)

    add_text(s, lx, ly + Inches(2.85), lw, Inches(0.4),
             "Structural Requirements",
             size=19, bold=True, color=NAVY)
    add_rect(s, lx, ly + Inches(3.28), Inches(1.6), Inches(0.04), TEAL)

    reqs = [
        ("Complete binary tree: ", "every level is fully filled except possibly the last, "
         "which is filled from left to right."),
        ("Heap Property: ", "a strict ordering rule between every parent and its children "
         "(no ordering between siblings)."),
        ("No pointers needed: ", "because the tree is complete, it is stored in a plain array."),
    ]
    tb = s.shapes.add_textbox(lx, ly + Inches(3.45), lw, Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (head, tail) in enumerate(reqs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.3; p.space_after = Pt(6)
        rb = p.add_run(); rb.text = "▸  "
        rb.font.name = "Calibri"; rb.font.size = Pt(16)
        rb.font.bold = True; rb.font.color.rgb = TEAL
        r1 = p.add_run(); r1.text = head
        r1.font.name = "Calibri"; r1.font.size = Pt(16)
        r1.font.bold = True; r1.font.color.rgb = NAVY
        r2 = p.add_run(); r2.text = tail
        r2.font.name = "Calibri"; r2.font.size = Pt(16)
        r2.font.color.rgb = DARK_GREY

    # Right panel — complexity table
    rx = Inches(7.4)
    add_rect(s, rx, Inches(1.5), Inches(5.4), Inches(4.5), WHITE,
             line_color=NAVY, line_width=Pt(1))
    add_rect(s, rx, Inches(1.5), Inches(5.4), Inches(0.55), NAVY)
    add_text(s, rx, Inches(1.5), Inches(5.4), Inches(0.55),
             "Time Complexity Summary",
             size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ("Peek minimum / maximum", "O(1)"),
        ("Insert an element",       "O(log N)"),
        ("Remove minimum / maximum", "O(log N)"),
        ("Search for a value",       "O(N)"),
        ("Build heap from N items",  "O(N)"),
    ]
    for i, (op, tc) in enumerate(rows):
        ry = Inches(2.15) + i * Inches(0.72)
        bg = LIGHT_GREY if i % 2 == 0 else WHITE
        add_rect(s, rx, ry, Inches(5.4), Inches(0.72), bg)
        add_text(s, rx + Inches(0.12), ry + Inches(0.1),
                 Inches(3.6), Inches(0.52),
                 op, size=15, color=DARK_GREY, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, rx + Inches(3.8), ry + Inches(0.12),
                 Inches(1.35), Inches(0.48), GOLD)
        add_text(s, rx + Inches(3.8), ry + Inches(0.12),
                 Inches(1.35), Inches(0.48),
                 tc, size=15, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_heap_property(page):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GREY)
    add_header(s, "The Heap Property", "The rule every heap must obey", page, TOTAL_SLIDES)

    for col, (label, subtitle, node_color, border_color, values, root_val) in enumerate([
        ("Min-Heap",
         "Parent ≤ Both Children  →  Root = Minimum",
         RGBColor(0x12, 0xA0, 0x9A), TEAL,
         {1: (Inches(3.5), Inches(2.3)),
          3: (Inches(2.3), Inches(3.5)),
          5: (Inches(4.7), Inches(3.5)),
          9: (Inches(1.5), Inches(4.7)),
          7: (Inches(3.1), Inches(4.7)),
          6: (Inches(4.1), Inches(4.7)),
          8: (Inches(5.5), Inches(4.7))},
         1),
        ("Max-Heap",
         "Parent ≥ Both Children  →  Root = Maximum",
         RGBColor(0xC0, 0x39, 0x2B), RED,
         {9: (Inches(9.6),  Inches(2.3)),
          7: (Inches(8.4),  Inches(3.5)),
          5: (Inches(10.8), Inches(3.5)),
          1: (Inches(7.6),  Inches(4.7)),
          3: (Inches(9.2),  Inches(4.7)),
          4: (Inches(10.2), Inches(4.7)),
          2: (Inches(11.6), Inches(4.7))},
         9),
    ]):
        # Card background
        cx = Inches(0.55 + col * 6.7)
        add_rect(s, cx, Inches(1.45), Inches(6.1), Inches(5.7),
                 WHITE, line_color=border_color, line_width=Pt(2))
        add_rect(s, cx, Inches(1.45), Inches(6.1), Inches(0.65), border_color)
        add_text(s, cx, Inches(1.45), Inches(6.1), Inches(0.65),
                 label, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx, Inches(2.1), Inches(6.1), Inches(0.4),
                 subtitle, size=13, italic=True, color=border_color,
                 align=PP_ALIGN.CENTER)

        # Draw edges first
        edge_pairs = [
            (root_val, list(values.keys())[1]),
            (root_val, list(values.keys())[2]),
            (list(values.keys())[1], list(values.keys())[3]),
            (list(values.keys())[1], list(values.keys())[4]),
            (list(values.keys())[2], list(values.keys())[5]),
            (list(values.keys())[2], list(values.keys())[6]),
        ]
        nr = Inches(0.38)
        for (a, b) in edge_pairs:
            x1, y1 = values[a]
            x2, y2 = values[b]
            add_edge(s, x1, y1 + nr, x2, y2 - nr,
                     color=border_color, weight=1.75)

        # Draw nodes
        for val, (nx, ny) in values.items():
            is_root = (val == root_val)
            fill = border_color if is_root else WHITE
            tc   = WHITE if is_root else DARK_GREY
            add_tree_node(s, nx, ny, val, r=nr,
                          fill=fill, border=border_color,
                          text_color=tc, size=16)

    # Key distinction callout at bottom
    add_rect(s, Inches(0.55), Inches(6.9), Inches(12.3), Inches(0.48),
             RGBColor(0xFF, 0xF3, 0xCD), line_color=GOLD, line_width=Pt(1.5))
    add_text(s, Inches(0.7), Inches(6.9), Inches(12.1), Inches(0.48),
             "Key distinction from Binary Search Trees:  In a heap, there is NO ordering between sibling nodes. "
             "A left child is not necessarily smaller than a right child.",
             size=13, italic=True, color=DARK_GREY,
             anchor=MSO_ANCHOR.MIDDLE)


def slide_array_representation(page):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GREY)
    add_header(s, "Array Representation",
               "Storing the tree in a plain one-dimensional array", page, TOTAL_SLIDES)

    # Left: explanation
    lx = Inches(0.6)
    add_text(s, lx, Inches(1.55), Inches(5.8), Inches(0.45),
             "Why an Array?", size=19, bold=True, color=NAVY)
    add_rect(s, lx, Inches(1.98), Inches(1.8), Inches(0.04), TEAL)
    add_text(s, lx, Inches(2.08), Inches(5.8), Inches(1.1),
             "Because a heap is always a complete binary tree, its nodes have no "
             "gaps. We can store the entire tree level by level in a standard array — "
             "no pointer or node class is required.",
             size=15, color=DARK_GREY)

    add_text(s, lx, Inches(3.25), Inches(5.8), Inches(0.45),
             "Index Formulas  (0-indexed array)", size=19, bold=True, color=NAVY)
    add_rect(s, lx, Inches(3.68), Inches(1.8), Inches(0.04), TEAL)

    formulas = [
        ("Parent of node at index i:", "(i − 1) // 2"),
        ("Left child of node at i:",   "2 × i + 1"),
        ("Right child of node at i:",  "2 × i + 2"),
    ]
    fy = Inches(3.82)
    for label, formula in formulas:
        add_text(s, lx, fy, Inches(2.8), Inches(0.42),
                 label, size=14, color=DARK_GREY, italic=True)
        add_rect(s, lx + Inches(2.85), fy + Inches(0.02),
                 Inches(2.8), Inches(0.38), TEAL)
        add_text(s, lx + Inches(2.85), fy + Inches(0.02),
                 Inches(2.8), Inches(0.38),
                 formula, size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font="Consolas")
        fy += Inches(0.58)

    add_text(s, lx, Inches(5.62), Inches(5.8), Inches(0.38),
             "The root is always stored at index 0.",
             size=14, color=TEAL, italic=True, bold=True)

    # Right: tree diagram
    rx_base = Inches(7.1)
    add_text(s, rx_base, Inches(1.55), Inches(5.8), Inches(0.4),
             "Example:  Min-Heap  [1, 3, 5, 9, 7, 6, 8]",
             size=17, bold=True, color=NAVY)
    add_rect(s, rx_base, Inches(1.93), Inches(2.0), Inches(0.04), TEAL)

    # Tree positions
    nr = Inches(0.36)
    positions = {
        "1": (Inches(10.1), Inches(2.6)),
        "3": (Inches(9.0),  Inches(3.7)),
        "5": (Inches(11.2), Inches(3.7)),
        "9": (Inches(8.3),  Inches(4.8)),
        "7": (Inches(9.7),  Inches(4.8)),
        "6": (Inches(10.7), Inches(4.8)),
        "8": (Inches(12.1), Inches(4.8)),
    }
    edges = [("1","3"),("1","5"),("3","9"),("3","7"),("5","6"),("5","8")]
    idx_labels = {
        "1":"[0]","3":"[1]","5":"[2]","9":"[3]","7":"[4]","6":"[5]","8":"[6]"
    }
    for (a,b) in edges:
        x1,y1 = positions[a]; x2,y2 = positions[b]
        add_edge(s, x1, y1+nr, x2, y2-nr, color=TEAL, weight=1.75)
    for val,(cx,cy) in positions.items():
        is_root = val == "1"
        fill = TEAL if is_root else WHITE
        tc   = WHITE if is_root else NAVY
        add_tree_node(s, cx, cy, val, r=nr, fill=fill,
                      border=TEAL, text_color=tc, size=15)
        add_text(s, cx - nr, cy - nr*2.0, nr*2, Inches(0.28),
                 idx_labels[val], size=10, color=DARK_GREY,
                 align=PP_ALIGN.CENTER)

    # Array visualisation
    add_array_cells(s, Inches(7.2), Inches(5.9),
                    [1, 3, 5, 9, 7, 6, 8],
                    cell_w=Inches(0.62), cell_h=Inches(0.55),
                    root_index=0, font_size=15)


def slide_sift_up(page):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GREY)
    add_header(s, "Sift-Up (Bubble Up)",
               "Restoring the heap property after inserting a new element", page, TOTAL_SLIDES)

    # Left: steps
    lx = Inches(0.6)
    add_text(s, lx, Inches(1.55), Inches(5.8), Inches(0.4),
             "Algorithm", size=19, bold=True, color=NAVY)
    add_rect(s, lx, Inches(1.93), Inches(1.6), Inches(0.04), TEAL)

    steps = [
        ("Step 1:", "Append the new element to the end of the array."),
        ("Step 2:", "Find its parent using index  (i − 1) // 2."),
        ("Step 3:", "If the new element is smaller than its parent (Min-Heap), swap them."),
        ("Step 4:", "Move the focus index up to the parent's position."),
        ("Step 5:", "Repeat Steps 2–4 until the element is larger than its parent, "
                   "or it reaches the root (index 0)."),
    ]
    tb = s.shapes.add_textbox(lx, Inches(2.1), Inches(5.8), Inches(3.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (head, tail) in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.3; p.space_after = Pt(6)
        rb = p.add_run(); rb.text = "▸  "
        rb.font.name="Calibri"; rb.font.size=Pt(16)
        rb.font.bold=True; rb.font.color.rgb=TEAL
        r1 = p.add_run(); r1.text = head + " "
        r1.font.name="Calibri"; r1.font.size=Pt(16)
        r1.font.bold=True; r1.font.color.rgb=NAVY
        r2 = p.add_run(); r2.text = tail
        r2.font.name="Calibri"; r2.font.size=Pt(16)
        r2.font.color.rgb=DARK_GREY

    # Complexity callout
    add_rect(s, lx, Inches(5.7), Inches(5.8), Inches(0.65),
             RGBColor(0xFF, 0xF3, 0xCD), line_color=GOLD, line_width=Pt(1.5))
    add_text(s, lx + Inches(0.1), Inches(5.72), Inches(5.6), Inches(0.62),
             "In the worst case, the element travels from the bottom level to the root: "
             "at most  log N  swaps.   Time complexity:  O(log N)",
             size=14, color=DARK_GREY, anchor=MSO_ANCHOR.MIDDLE)

    # Right: code
    rx = Inches(6.8)
    add_text(s, rx, Inches(1.55), Inches(6.1), Inches(0.4),
             "Python Implementation", size=19, bold=True, color=NAVY)
    add_rect(s, rx, Inches(1.93), Inches(1.8), Inches(0.04), TEAL)

    code = [
        "def _sift_up(self, index):",
        "    # While we have a parent AND",
        "    # the parent is LARGER than the current element:",
        "    while self.has_parent(index):",
        "        parent_i = self.get_parent_index(index)",
        "",
        "        if self.heap[parent_i] > self.heap[index]:",
        "            # Swap current element with its parent",
        "            self.heap[parent_i], self.heap[index] = \\",
        "                self.heap[index], self.heap[parent_i]",
        "            index = parent_i   # move focus upward",
        "        else:",
        "            break  # heap property is satisfied",
    ]
    add_code_block(s, rx, Inches(2.1), Inches(6.1), Inches(4.2), code, size=13)


def slide_sift_down(page):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GREY)
    add_header(s, "Sift-Down (Sink Down)",
               "Restoring the heap property after removing the root element", page, TOTAL_SLIDES)

    lx = Inches(0.6)
    add_text(s, lx, Inches(1.55), Inches(5.8), Inches(0.4),
             "Algorithm", size=19, bold=True, color=NAVY)
    add_rect(s, lx, Inches(1.93), Inches(1.6), Inches(0.04), TEAL)

    steps = [
        ("Step 1:", "Record the root value — this is the element being returned."),
        ("Step 2:", "Move the last array element to position [0] (the root). "
                   "Remove the last position."),
        ("Step 3:", "Compare the element with its left and right children."),
        ("Step 4:", "Identify the smaller of the two children (for a Min-Heap)."),
        ("Step 5:", "If the element is larger than the smaller child, swap them."),
        ("Step 6:", "Repeat Steps 3–5, moving downward until the element is smaller "
                   "than both children, or becomes a leaf."),
    ]
    tb = s.shapes.add_textbox(lx, Inches(2.1), Inches(5.8), Inches(3.8))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (head, tail) in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25; p.space_after = Pt(5)
        rb = p.add_run(); rb.text = "▸  "
        rb.font.name="Calibri"; rb.font.size=Pt(15)
        rb.font.bold=True; rb.font.color.rgb=TEAL
        r1 = p.add_run(); r1.text = head + " "
        r1.font.name="Calibri"; r1.font.size=Pt(15)
        r1.font.bold=True; r1.font.color.rgb=NAVY
        r2 = p.add_run(); r2.text = tail
        r2.font.name="Calibri"; r2.font.size=Pt(15)
        r2.font.color.rgb=DARK_GREY

    add_rect(s, lx, Inches(5.98), Inches(5.8), Inches(0.55),
             RGBColor(0xFF, 0xF3, 0xCD), line_color=GOLD, line_width=Pt(1.5))
    add_text(s, lx + Inches(0.1), Inches(6.0), Inches(5.6), Inches(0.52),
             "At most log N levels to traverse downward.   Time complexity:  O(log N)",
             size=14, color=DARK_GREY, anchor=MSO_ANCHOR.MIDDLE)

    rx = Inches(6.8)
    add_text(s, rx, Inches(1.55), Inches(6.1), Inches(0.4),
             "Python Implementation", size=19, bold=True, color=NAVY)
    add_rect(s, rx, Inches(1.93), Inches(1.8), Inches(0.04), TEAL)

    code = [
        "def _sift_down(self, index):",
        "    # Continue while there is at least a left child",
        "    while self.has_left_child(index):",
        "        left  = self.get_left_child_index(index)",
        "        right = self.get_right_child_index(index)",
        "",
        "        # Assume left child is the smaller candidate",
        "        smaller = left",
        "",
        "        # If right child exists and is smaller, use it",
        "        if (self.has_right_child(index) and",
        "                self.heap[right] < self.heap[left]):",
        "            smaller = right",
        "",
        "        # If current element is already smaller, stop",
        "        if self.heap[index] < self.heap[smaller]:",
        "            break",
        "",
        "        # Swap and continue downward",
        "        self.heap[index], self.heap[smaller] = (",
        "            self.heap[smaller], self.heap[index]",
        "        )",
        "        index = smaller",
    ]
    add_code_block(s, rx, Inches(2.1), Inches(6.1), Inches(4.85), code, size=11.5)


def slide_insert_example(page):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GREY)
    add_header(s, "Worked Example — Insert",
               "Inserting the value 2 into the Min-Heap  [1, 3, 5, 9, 7, 6, 8]",
               page, TOTAL_SLIDES)

    add_text(s, Inches(0.6), Inches(1.55), Inches(12.3), Inches(0.38),
             "Initial heap array:", size=16, color=DARK_GREY)
    add_array_cells(s, Inches(0.6), Inches(1.92),
                    [1, 3, 5, 9, 7, 6, 8],
                    cell_w=Inches(0.6), cell_h=Inches(0.55),
                    root_index=0)

    steps = [
        ("Step 1 — Append",
         "Append 2 at the end (index 7)  →  [1, 3, 5, 9, 7, 6, 8, 2]",
         [1, 3, 5, 9, 7, 6, 8, 2], {7}),
        ("Step 2 — Compare with parent at index 3 (value 9)",
         "2 < 9  →  Swap indices 7 and 3  →  [1, 3, 5, 2, 7, 6, 8, 9]",
         [1, 3, 5, 2, 7, 6, 8, 9], {3}),
        ("Step 3 — Compare with parent at index 1 (value 3)",
         "2 < 3  →  Swap indices 3 and 1  →  [1, 2, 5, 3, 7, 6, 8, 9]",
         [1, 2, 5, 3, 7, 6, 8, 9], {1}),
        ("Step 4 — Compare with parent at index 0 (value 1)",
         "2 > 1  →  Stop. The heap property is satisfied. ✓",
         [1, 2, 5, 3, 7, 6, 8, 9], set()),
    ]

    for i, (heading, desc, arr, hl) in enumerate(steps):
        y = Inches(2.8) + i * Inches(1.08)
        add_rect(s, Inches(0.6), y, Inches(12.3), Inches(1.0),
                 WHITE, line_color=TEAL, line_width=Pt(1))
        add_rect(s, Inches(0.6), y, Inches(0.06), Inches(1.0), TEAL)
        add_text(s, Inches(0.8), y + Inches(0.06),
                 Inches(6.5), Inches(0.38),
                 heading, size=14, bold=True, color=NAVY)
        add_text(s, Inches(0.8), y + Inches(0.43),
                 Inches(5.8), Inches(0.38),
                 desc, size=13, color=DARK_GREY)
        add_array_cells(s, Inches(7.1), y + Inches(0.18),
                        arr,
                        cell_w=Inches(0.5), cell_h=Inches(0.46),
                        highlight_indices=hl, root_index=0,
                        font_size=13)


def slide_extract_example(page):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GREY)
    add_header(s, "Worked Example — Extract Minimum",
               "Removing the minimum (root) from  [1, 3, 5, 9, 7, 6, 8]",
               page, TOTAL_SLIDES)

    add_text(s, Inches(0.6), Inches(1.55), Inches(12.3), Inches(0.38),
             "Initial heap array:", size=16, color=DARK_GREY)
    add_array_cells(s, Inches(0.6), Inches(1.92),
                    [1, 3, 5, 9, 7, 6, 8],
                    cell_w=Inches(0.6), cell_h=Inches(0.55),
                    root_index=0)

    steps = [
        ("Step 1 — Record and Replace",
         "Record root value 1 (the result). Move last element (8) to root.",
         [8, 3, 5, 9, 7, 6], {0}),
        ("Step 2 — Smaller child of 8 is 3 (left, index 1)",
         "8 > 3  →  Swap indices 0 and 1  →  [3, 8, 5, 9, 7, 6]",
         [3, 8, 5, 9, 7, 6], {1}),
        ("Step 3 — Smaller child of 8 is 7 (right, index 4)",
         "8 > 7  →  Swap indices 1 and 4  →  [3, 7, 5, 9, 8, 6]",
         [3, 7, 5, 9, 8, 6], {4}),
        ("Step 4 — Node 8 has no children",
         "Stop. Min-Heap property is fully restored. Minimum returned: 1  ✓",
         [3, 7, 5, 9, 8, 6], set()),
    ]

    for i, (heading, desc, arr, hl) in enumerate(steps):
        y = Inches(2.8) + i * Inches(1.08)
        add_rect(s, Inches(0.6), y, Inches(12.3), Inches(1.0),
                 WHITE, line_color=RED, line_width=Pt(1))
        add_rect(s, Inches(0.6), y, Inches(0.06), Inches(1.0), RED)
        add_text(s, Inches(0.8), y + Inches(0.06),
                 Inches(6.5), Inches(0.38),
                 heading, size=14, bold=True, color=NAVY)
        add_text(s, Inches(0.8), y + Inches(0.43),
                 Inches(5.8), Inches(0.38),
                 desc, size=13, color=DARK_GREY)
        add_array_cells(s, Inches(7.1), y + Inches(0.18),
                        arr,
                        cell_w=Inches(0.56), cell_h=Inches(0.46),
                        highlight_indices=hl, root_index=0,
                        font_size=13)


def slide_priority_queue(page):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GREY)
    add_header(s, "Priority Queues",
               "The abstract data structure that a heap implements", page, TOTAL_SLIDES)

    # What is a PQ
    lx = Inches(0.6)
    add_text(s, lx, Inches(1.55), Inches(12.3), Inches(0.45),
             "A Priority Queue is an abstract data structure: elements leave in order of priority, "
             "not in order of arrival. A Binary Heap is the standard implementation.",
             size=16, color=DARK_GREY)

    # Comparison table
    add_text(s, lx, Inches(2.2), Inches(5.8), Inches(0.4),
             "Regular Queue vs. Priority Queue", size=19, bold=True, color=NAVY)
    add_rect(s, lx, Inches(2.58), Inches(1.8), Inches(0.04), TEAL)

    table_data = [
        ("Regular Queue (FIFO)",          "Priority Queue"),
        ("First inserted leaves first.",   "Highest-priority element leaves first."),
        ("No concept of priority.",        "Each element carries a priority value."),
        ("Dequeue:  O(1)",                 "Dequeue:  O(log N)"),
        ("Implemented with array/list.",   "Implemented with a Binary Heap."),
    ]
    for row_i, (left, right) in enumerate(table_data):
        ry = Inches(2.7) + row_i * Inches(0.6)
        bg  = NAVY        if row_i == 0 else (LIGHT_GREY if row_i % 2 == 1 else WHITE)
        tc  = WHITE       if row_i == 0 else DARK_GREY
        fsz = 14          if row_i == 0 else 13
        bld = (row_i == 0)
        add_rect(s, lx,               ry, Inches(2.85), Inches(0.6), bg,
                 line_color=DARK_GREY if row_i > 0 else None, line_width=Pt(0.5))
        add_rect(s, lx + Inches(2.85), ry, Inches(2.85), Inches(0.6), bg,
                 line_color=DARK_GREY if row_i > 0 else None, line_width=Pt(0.5))
        add_text(s, lx + Inches(0.1),               ry + Inches(0.08),
                 Inches(2.65), Inches(0.44),
                 left, size=fsz, bold=bld, color=tc, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, lx + Inches(2.85) + Inches(0.1), ry + Inches(0.08),
                 Inches(2.65), Inches(0.44),
                 right, size=fsz, bold=bld, color=tc, anchor=MSO_ANCHOR.MIDDLE)

    # Applications
    rx = Inches(7.2)
    add_text(s, rx, Inches(2.2), Inches(5.8), Inches(0.4),
             "Real-World Applications", size=19, bold=True, color=NAVY)
    add_rect(s, rx, Inches(2.58), Inches(1.8), Inches(0.04), TEAL)

    apps = [
        (TEAL,   "OS Process Scheduling",
         "The CPU runs the highest-priority process first. High-priority processes "
         "can pre-empt lower-priority ones."),
        (GREEN,  "Network Packet Routing",
         "Routers prioritise urgent packets (voice calls, video streams) over bulk "
         "data transfers using a priority queue."),
        (PURPLE, "Dijkstra's Shortest Path",
         "A Min-Heap always selects the nearest unvisited node first — core to GPS "
         "navigation and routing software."),
    ]
    for i, (col, title, desc) in enumerate(apps):
        ay = Inches(2.75) + i * Inches(1.5)
        add_rect(s, rx, ay, Inches(5.8), Inches(1.35),
                 WHITE, line_color=col, line_width=Pt(2))
        add_rect(s, rx, ay, Inches(0.08), Inches(1.35), col)
        add_text(s, rx + Inches(0.2), ay + Inches(0.1),
                 Inches(5.4), Inches(0.42),
                 title, size=16, bold=True, color=col)
        add_text(s, rx + Inches(0.2), ay + Inches(0.52),
                 Inches(5.4), Inches(0.72),
                 desc, size=13, color=DARK_GREY)


def slide_heapq(page):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GREY)
    add_header(s, "Python's heapq Module",
               "The built-in Min-Heap operating on a plain Python list", page, TOTAL_SLIDES)

    # Table of functions
    lx = Inches(0.6)
    add_text(s, lx, Inches(1.55), Inches(5.8), Inches(0.4),
             "Core Functions", size=19, bold=True, color=NAVY)
    add_rect(s, lx, Inches(1.93), Inches(1.8), Inches(0.04), TEAL)

    funcs = [
        ("Function",              "Description",                           "Time"),
        ("heappush(h, x)",        "Insert x and restore heap property",    "O(log N)"),
        ("heappop(h)",            "Remove and return the smallest element", "O(log N)"),
        ("heapify(h)",            "Convert a list into a heap in-place",    "O(N)"),
        ("h[0]",                  "Read smallest element without removing", "O(1)"),
        ("heappushpop(h, x)",     "Push then pop in one efficient step",    "O(log N)"),
    ]
    for row_i, (fn, desc, tc_val) in enumerate(funcs):
        ry = Inches(2.07) + row_i * Inches(0.62)
        bg = NAVY if row_i == 0 else (LIGHT_GREY if row_i % 2 == 1 else WHITE)
        fc = WHITE if row_i == 0 else DARK_GREY
        bld = (row_i == 0)
        fn_font = "Consolas" if row_i > 0 else "Calibri"
        add_rect(s, lx,              ry, Inches(1.8), Inches(0.62), bg,
                 line_color=DARK_GREY if row_i > 0 else None, line_width=Pt(0.5))
        add_rect(s, lx+Inches(1.8),  ry, Inches(3.0), Inches(0.62), bg,
                 line_color=DARK_GREY if row_i > 0 else None, line_width=Pt(0.5))
        add_rect(s, lx+Inches(4.8),  ry, Inches(1.2), Inches(0.62), bg,
                 line_color=DARK_GREY if row_i > 0 else None, line_width=Pt(0.5))
        add_text(s, lx+Inches(0.06),     ry+Inches(0.08), Inches(1.68), Inches(0.46),
                 fn, size=13, bold=bld, color=fc,
                 font=fn_font, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, lx+Inches(1.86),     ry+Inches(0.08), Inches(2.88), Inches(0.46),
                 desc, size=13, bold=bld, color=fc, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, lx+Inches(4.86),     ry+Inches(0.08), Inches(1.08), Inches(0.46),
                 tc_val, size=13, bold=bld, color=fc,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Important note
    add_rect(s, lx, Inches(5.95), Inches(5.95), Inches(0.7),
             RGBColor(0xFF, 0xF3, 0xCD), line_color=GOLD, line_width=Pt(1.5))
    add_text(s, lx+Inches(0.12), Inches(5.95), Inches(5.71), Inches(0.7),
             "Important: Python's heapq only provides a Min-Heap. To simulate a Max-Heap, "
             "insert negated values (−x) and negate the result on extraction.",
             size=13, italic=True, color=DARK_GREY, anchor=MSO_ANCHOR.MIDDLE)

    # Code example
    rx = Inches(7.1)
    add_text(s, rx, Inches(1.55), Inches(5.8), Inches(0.4),
             "Usage Example", size=19, bold=True, color=NAVY)
    add_rect(s, rx, Inches(1.93), Inches(1.8), Inches(0.04), TEAL)

    code = [
        "import heapq",
        "",
        "numbers = [5, 1, 9, 3, 2]",
        "",
        "# Build a Min-Heap in place — O(N)",
        "heapq.heapify(numbers)",
        "# numbers is now: [1, 2, 9, 3, 5]",
        "",
        "# Read the minimum without removing it",
        "minimum = numbers[0]        # -> 1",
        "",
        "# Remove and return the minimum — O(log N)",
        "val = heapq.heappop(numbers)   # -> 1",
        "",
        "# Insert a new element — O(log N)",
        "heapq.heappush(numbers, 4)",
        "",
        "# ---- Simulating a Max-Heap ----",
        "max_heap = []",
        "for x in [5, 1, 9]:",
        "    heapq.heappush(max_heap, -x)",
        "largest = -heapq.heappop(max_heap)   # -> 9",
    ]
    add_code_block(s, rx, Inches(2.1), Inches(5.9), Inches(5.5), code, size=12)


def slide_kth_largest(page):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GREY)
    add_header(s, "K-th Largest Element",
               "A classic problem solved efficiently using a Min-Heap", page, TOTAL_SLIDES)

    # Problem statement
    add_rect(s, Inches(0.6), Inches(1.55), Inches(12.3), Inches(0.75),
             RGBColor(0xE8, 0xF4, 0xFD), line_color=TEAL, line_width=Pt(1.5))
    add_text(s, Inches(0.75), Inches(1.6), Inches(12.0), Inches(0.65),
             "Problem:  Given an unsorted list of N numbers, find the K-th largest element.   "
             "Constraint: the solution must run in O(N log K) time, not O(N log N).",
             size=15, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

    lx = Inches(0.6)
    add_text(s, lx, Inches(2.5), Inches(5.8), Inches(0.4),
             "Strategy: Min-Heap Capped at Size K", size=18, bold=True, color=NAVY)
    add_rect(s, lx, Inches(2.88), Inches(2.0), Inches(0.04), TEAL)

    steps = [
        ("Step 1:", "Maintain a Min-Heap of at most K elements."),
        ("Step 2:", "For each element in the list, push it onto the heap."),
        ("Step 3:", "If the heap grows beyond size K, remove the smallest "
                   "element (the root)."),
        ("Step 4:", "After all elements are processed, the root of the heap "
                   "is the K-th largest element overall."),
    ]
    tb = s.shapes.add_textbox(lx, Inches(3.04), Inches(5.8), Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (head, tail) in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing=1.3; p.space_after=Pt(6)
        rb = p.add_run(); rb.text="▸  "
        rb.font.name="Calibri"; rb.font.size=Pt(15)
        rb.font.bold=True; rb.font.color.rgb=TEAL
        r1 = p.add_run(); r1.text=head+" "
        r1.font.name="Calibri"; r1.font.size=Pt(15)
        r1.font.bold=True; r1.font.color.rgb=NAVY
        r2 = p.add_run(); r2.text=tail
        r2.font.name="Calibri"; r2.font.size=Pt(15)
        r2.font.color.rgb=DARK_GREY

    # Why not sort?
    add_rect(s, lx, Inches(5.78), Inches(5.8), Inches(0.85),
             RGBColor(0xFF, 0xF3, 0xCD), line_color=GOLD, line_width=Pt(1.5))
    add_text(s, lx+Inches(0.12), Inches(5.82), Inches(5.56), Inches(0.78),
             "Why not sort the array?  Sorting takes O(N log N) and loads all N elements "
             "into memory. The heap approach uses only O(K) space and runs in O(N log K) — "
             "significantly faster when K is much smaller than N.",
             size=13, color=DARK_GREY, anchor=MSO_ANCHOR.MIDDLE)

    # Code
    rx = Inches(7.0)
    add_text(s, rx, Inches(2.5), Inches(5.9), Inches(0.4),
             "Python Solution", size=18, bold=True, color=NAVY)
    add_rect(s, rx, Inches(2.88), Inches(1.8), Inches(0.04), TEAL)

    code = [
        "import heapq",
        "",
        "def find_kth_largest(nums: list, k: int) -> int:",
        "    min_heap = []",
        "",
        "    for num in nums:",
        "        heapq.heappush(min_heap, num)",
        "",
        "        # If heap exceeds K elements,",
        "        # discard the smallest one.",
        "        # This keeps only the K largest seen so far.",
        "        if len(min_heap) > k:",
        "            heapq.heappop(min_heap)",
        "",
        "    # The root is the K-th largest overall.",
        "    return min_heap[0]",
        "",
        "# Example",
        "nums = [3, 2, 1, 5, 6, 4]",
        "k = 2",
        "print(find_kth_largest(nums, k))   # -> 5",
    ]
    add_code_block(s, rx, Inches(3.04), Inches(5.9), Inches(3.82), code, size=12)

    # Complexity box
    add_rect(s, rx, Inches(7.0), Inches(5.9), Inches(0.2), TEAL)   # spacer
    for i, (label, val) in enumerate([("Time: ", "O(N log K)"), ("Space: ", "O(K)")]):
        bx = rx + i * Inches(2.9)
        add_rect(s, bx, Inches(7.02), Inches(2.6), Inches(0.38),
                 WHITE, line_color=TEAL, line_width=Pt(1))
        add_text(s, bx+Inches(0.08), Inches(7.02), Inches(1.0), Inches(0.38),
                 label, size=13, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, bx+Inches(1.0), Inches(7.02), Inches(1.5), Inches(0.38),
                 val, size=13, bold=True, color=NAVY,
                 font="Consolas", anchor=MSO_ANCHOR.MIDDLE)


def slide_floyd(page):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GREY)
    add_header(s, "Floyd's O(N) Heap Construction",
               "Building a valid heap from an unsorted array in linear time", page, TOTAL_SLIDES)

    add_text(s, Inches(0.6), Inches(1.55), Inches(12.3), Inches(0.38),
             "There are two approaches to building a heap from an unsorted array of N elements:",
             size=16, color=DARK_GREY)

    # Approach 1
    add_rect(s, Inches(0.6), Inches(2.05), Inches(5.85), Inches(2.85),
             WHITE, line_color=RED, line_width=Pt(2))
    add_rect(s, Inches(0.6), Inches(2.05), Inches(5.85), Inches(0.55), RED)
    add_text(s, Inches(0.6), Inches(2.05), Inches(5.85), Inches(0.55),
             "Approach 1  —  Insert One by One",
             size=17, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.75), Inches(2.7), Inches(5.55), Inches(1.6),
             "Call push() (Sift-Up) N times, once per element.\n\n"
             "Each insertion costs O(log N).\n"
             "Total cost for N elements:  O(N log N)\n\n"
             "This is the straightforward but slower approach.",
             size=15, color=DARK_GREY)

    # Approach 2
    add_rect(s, Inches(7.0), Inches(2.05), Inches(5.85), Inches(2.85),
             WHITE, line_color=GREEN, line_width=Pt(2))
    add_rect(s, Inches(7.0), Inches(2.05), Inches(5.85), Inches(0.55), GREEN)
    add_text(s, Inches(7.0), Inches(2.05), Inches(5.85), Inches(0.55),
             "Approach 2  —  Floyd's Bottom-Up Method",
             size=17, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.15), Inches(2.7), Inches(5.55), Inches(0.8),
             "Start from the last non-leaf node at index (N // 2) − 1 and call "
             "Sift-Down on every node, moving up toward the root.",
             size=15, color=DARK_GREY)

    code = [
        "n = len(arr)",
        "# From last non-leaf to root:",
        "for i in range(n // 2 - 1, -1, -1):",
        "    sift_down(arr, i, n)",
    ]
    add_code_block(s, Inches(7.15), Inches(3.6), Inches(5.55), Inches(1.1), code, size=13)

    add_text(s, Inches(7.15), Inches(4.78), Inches(5.55), Inches(0.38),
             "Total cost:  O(N)  — Python's heapq.heapify() uses this method.",
             size=14, bold=True, color=GREEN)

    # Explanation callout
    add_rect(s, Inches(0.6), Inches(5.05), Inches(12.3), Inches(1.2),
             RGBColor(0xE8, 0xF8, 0xF1), line_color=GREEN, line_width=Pt(1.5))
    add_text(s, Inches(0.6), Inches(5.05), Inches(1.3), Inches(1.2),
             "Why O(N)?", size=15, bold=True, color=GREEN,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(1.9), Inches(5.15), Inches(0.04), Inches(1.0), GREEN)
    add_text(s, Inches(2.05), Inches(5.1), Inches(10.7), Inches(1.1),
             "Approximately half of all nodes in a complete binary tree are leaf nodes that "
             "require no sifting at all. Only a small number of nodes exist near the root where "
             "the most work is performed. When the total work across all levels is summed, it "
             "converges to O(N) — proven rigorously using geometric series analysis.",
             size=14, color=DARK_GREY)

    # Comparison table
    add_rect(s, Inches(0.6), Inches(6.4), Inches(12.3), Inches(0.55), NAVY)
    for i, (col1, col2, col3) in enumerate([
        ("Approach",            "Build Time",  "Extra Space"),
        ("Insert one by one",   "O(N log N)", "O(1)"),
        ("Floyd's bottom-up",   "O(N)",        "O(1)"),
    ]):
        ry = Inches(6.4) + i * Inches(0.54)
        bg = NAVY if i == 0 else (LIGHT_GREY if i % 2 == 1 else WHITE)
        fc = WHITE if i == 0 else DARK_GREY
        for j, txt in enumerate([col1, col2, col3]):
            bx = Inches(0.6) + j * Inches(4.1)
            add_rect(s, bx, ry, Inches(4.1), Inches(0.54), bg,
                     line_color=DARK_GREY if i > 0 else None, line_width=Pt(0.5))
            add_text(s, bx+Inches(0.1), ry+Inches(0.05), Inches(3.9), Inches(0.44),
                     txt, size=14, bold=(i==0), color=fc,
                     anchor=MSO_ANCHOR.MIDDLE)


def slide_summary(page):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, 0, Inches(0.5), SLIDE_H, GOLD)
    add_rect(s, 0, SLIDE_H - Inches(0.1), SLIDE_W, Inches(0.1), GOLD)

    add_text(s, Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.65),
             "Week 07 Summary", size=36, bold=True, color=WHITE)
    add_rect(s, Inches(0.8), Inches(1.25), Inches(4.0), Inches(0.06), GOLD)

    summary = [
        ("Binary Heap", "A complete binary tree specialised for fast retrieval of the minimum or maximum."),
        ("Heap Property", "Every parent is ≤ both children (Min-Heap) or ≥ both children (Max-Heap)."),
        ("Array Storage", "Because the tree is complete, it maps directly to a 1-D array — no pointers needed."),
        ("Sift-Up", "Called after insertion. Moves the new element upward until the property is restored.  O(log N)"),
        ("Sift-Down", "Called after extraction. Moves the displaced element downward until the property is restored.  O(log N)"),
        ("Priority Queue", "Abstract structure using a heap — elements leave in priority order, not arrival order."),
        ("heapq Module", "Python's built-in Min-Heap on a list: heappush, heappop, heapify."),
        ("K-th Largest", "Maintain a Min-Heap of size K while scanning N elements.  Time: O(N log K),  Space: O(K)"),
        ("Floyd's Method", "Build a heap bottom-up by calling Sift-Down from the last non-leaf to the root.  O(N)"),
    ]

    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(12.0), Inches(5.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (head, tail) in enumerate(summary):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.3; p.space_after = Pt(5)
        rb = p.add_run(); rb.text = "▸  "
        rb.font.name="Calibri"; rb.font.size=Pt(16)
        rb.font.bold=True; rb.font.color.rgb=GOLD
        r1 = p.add_run(); r1.text = head + ":  "
        r1.font.name="Calibri"; r1.font.size=Pt(16)
        r1.font.bold=True; r1.font.color.rgb=GOLD
        r2 = p.add_run(); r2.text = tail
        r2.font.name="Calibri"; r2.font.size=Pt(16)
        r2.font.color.rgb=LIGHT_GREY

    add_text(s, Inches(0.8), SLIDE_H - Inches(0.65),
             Inches(11.5), Inches(0.38),
             "Data Structures & Algorithms  ·  Week 07  ·  ATU Computer Science  ·  2024 – 2025",
             size=12, color=RGBColor(0xA0, 0xB0, 0xC0), italic=True)


# ──────────────────────────────────────────────────────────────────────────────
# Build all slides
# ──────────────────────────────────────────────────────────────────────────────

TOTAL_SLIDES = 15

slide_cover()                      # 1
slide_agenda(2)                    # 2
slide_what_is_heap(3)              # 3
slide_heap_property(4)             # 4
slide_array_representation(5)     # 5
slide_sift_up(6)                   # 6
slide_sift_down(7)                 # 7
slide_insert_example(8)           # 8
slide_extract_example(9)          # 9
slide_priority_queue(10)          # 10
slide_heapq(11)                   # 11
slide_kth_largest(12)             # 12
slide_floyd(13)                   # 13
slide_summary(14)                 # 14

# Slide 15 — blank "Questions?" slide
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, 0, Inches(0.5), SLIDE_H, GOLD)
add_rect(s, 0, SLIDE_H - Inches(0.1), SLIDE_W, Inches(0.1), GOLD)
add_text(s, Inches(1.5), Inches(2.6), Inches(10.0), Inches(1.4),
         "Questions?", size=72, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(4.3), Inches(10.0), Inches(0.6),
         "Data Structures & Algorithms  ·  Week 07  ·  Heaps & Priority Queues",
         size=18, italic=True, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

prs.save("heaps_dsa_final.pptx")
print("✅  Saved:  heaps_dsa_final.pptx")
