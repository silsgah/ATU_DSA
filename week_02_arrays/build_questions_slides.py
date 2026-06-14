import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# -------- Theme --------
NAVY        = RGBColor(0x0B, 0x2A, 0x4A)
TEAL        = RGBColor(0x12, 0x7A, 0x8A)
GOLD        = RGBColor(0xE0, 0xA8, 0x1F)
LIGHT_GREY  = RGBColor(0xF2, 0xF4, 0xF7)
DARK_GREY   = RGBColor(0x33, 0x3D, 0x4A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xC0, 0x39, 0x2B)
PURPLE      = RGBColor(0x6A, 0x1B, 0x9A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_GREY,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
              italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=18, color=DARK_GREY,
                bullet_color=TEAL, line_spacing=1.25, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        rb = p.add_run(); rb.text = "▸  "
        rb.font.name = font; rb.font.size = Pt(size)
        rb.font.bold = True; rb.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run(); r1.text = head
            r1.font.name = font; r1.font.size = Pt(size)
            r1.font.bold = True; r1.font.color.rgb = NAVY
            r2 = p.add_run(); r2.text = tail
            r2.font.name = font; r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = item
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb

def add_header(slide, title, subtitle=None, page_num=None, total=None):
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    add_rect(slide, Inches(0.35), 0, SLIDE_W - Inches(0.35), Inches(0.08), GOLD)
    add_text(slide, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.78), Inches(11.5), Inches(0.45),
                 subtitle, size=15, italic=True, color=TEAL)
    add_rect(slide, Inches(0.6), Inches(1.22), Inches(2.0), Inches(0.04), GOLD)
    add_text(slide, Inches(0.6), SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
             "Data Structures & Algorithms · Week 02 · Practice Questions",
             size=10, color=DARK_GREY, italic=True)
    if page_num is not None and total is not None:
        add_text(slide, SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.45),
                 Inches(1.2), Inches(0.3),
                 f"Slide {page_num} / {total}", size=10, color=DARK_GREY,
                 align=PP_ALIGN.RIGHT)

# ---------- Slides ----------

def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(5.5), SLIDE_W, Inches(0.12), GOLD)
    add_rect(s, 0, Inches(5.7), SLIDE_W, Inches(0.04), TEAL)
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.9), Inches(1.0),
                               Inches(2.8), Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb = GOLD
    chip.line.fill.background(); chip.shadow.inherit = False
    add_text(s, Inches(0.9), Inches(1.0), Inches(2.8), Inches(0.55),
             "WEEK 02 - WORKBOOK", size=16, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(1.4),
             "Arrays, Strings, & Pointers", size=58, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.7),
             "Self-Study Questions, Visual Exercises, and Modeling Scenarios",
             size=24, color=GOLD, italic=True)
    add_text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.6),
             "A Supplementary Workbook for Data Structures & Algorithms Students",
             size=18, color=LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(6.0), Inches(8), Inches(0.4),
             "Topic    ·   Dynamic Arrays & Two-Pointer Pattern",
             size=14, color=WHITE)

def slide_conceptual_1(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Part 1: Conceptual Questions (1/2)", "Understanding memory allocations and resizes", n, total)
    
    items = [
        ("1. Dynamic Array Growth: ", "Physical computer memory is allocated in contiguous blocks. Since arrays must be stored sequentially, explain how dynamic structures (like Python lists) grow. What steps occur when capacity is exceeded?"),
        ("2. Amortized Cost: ", "If we append N elements one by one to a dynamic array that starts with a capacity of 1 and doubles when full: how many resizing operations occur in total? Why is the time complexity considered O(1) amortized?"),
        ("3. Memory Overhead: ", "Suppose a dynamic array grows by doubling capacity. If it resizes from 1 million cells to 2 million, what is the maximum number of elements actually present immediately after resize? How much memory is 'wasted' temporarily?")
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.0), items, size=15, line_spacing=1.3)

def slide_conceptual_2(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Part 1: Conceptual Questions (2/2)", "String immutability and pointer movements", n, total)
    
    items = [
        ("1. String Immutability: ", "In languages like Python and Java, strings are immutable. If you concatenate characters to a string inside a loop N times (e.g., s += char), what is the total time complexity? Why is it better to use a character list instead?"),
        ("2. Two-Pointer Mechanics: ", "Contrast the two primary variations of the Two-Pointer pattern:\n(a) Boundary Pointers: Left and right moving towards the center (e.g. palindrome check).\n(b) Sliding Window: Fast and slow pointers starting from the same end. Describe a scenario suited for each."),
        ("3. Memory Pointers in Arrays: ", "An array lookup (A[i]) runs in O(1) time. Explain the mathematical formula that calculates the exact memory address of A[i] under the hood using the array base address, index i, and data type size.")
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.0), items, size=15, line_spacing=1.3)

def slide_trace_twosum(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Part 2: Hand-Tracing Two Pointers", "Simulating Two Sum II on a sorted array", n, total)
    
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.8),
             "Exercise: Trace two_sum_sorted([1, 3, 5, 8, 10, 15], 13) using the two-pointer approach.\nWrite down the pointers, sum, and decisions at each loop iteration.",
             size=15, italic=True, color=DARK_GREY)
             
    # Table template for tracing
    add_rect(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.55), NAVY)
    add_text(s, Inches(0.9), Inches(2.5), Inches(1.5), Inches(0.4), "Step", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.4), Inches(2.5), Inches(2.5), Inches(0.4), "Left Pointer (Idx / Val)", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(4.9), Inches(2.5), Inches(2.5), Inches(0.4), "Right Pointer (Idx / Val)", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.4), Inches(2.5), Inches(1.8), Inches(0.4), "Current Sum", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.2), Inches(2.5), Inches(3.2), Inches(0.4), "Decision / Pointer Move", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    for i in range(6):
        y_off = Inches(2.95 + i * 0.55)
        bg = WHITE if i % 2 == 0 else LIGHT_GREY
        add_rect(s, Inches(0.9), y_off, Inches(11.5), Inches(0.55), bg, line=DARK_GREY)
        if i == 0:
            add_text(s, Inches(0.9), y_off + Inches(0.12), Inches(1.5), Inches(0.4), "0 (Start)", size=12, color=DARK_GREY, align=PP_ALIGN.CENTER)
            add_text(s, Inches(2.4), y_off + Inches(0.12), Inches(2.5), Inches(0.4), "idx 0 (val = 1)", size=12, color=DARK_GREY, align=PP_ALIGN.CENTER)
            add_text(s, Inches(4.9), y_off + Inches(0.12), Inches(2.5), Inches(0.4), "idx 5 (val = 15)", size=12, color=DARK_GREY, align=PP_ALIGN.CENTER)
            add_text(s, Inches(7.4), y_off + Inches(0.12), Inches(1.8), Inches(0.4), "16", size=12, color=DARK_GREY, align=PP_ALIGN.CENTER)
            add_text(s, Inches(9.2), y_off + Inches(0.12), Inches(3.2), Inches(0.4), "16 > 13: Move Right (right -= 1)", size=12, color=DARK_GREY, align=PP_ALIGN.CENTER)
        else:
            add_text(s, Inches(0.9), y_off + Inches(0.12), Inches(1.5), Inches(0.4), str(i), size=12, color=DARK_GREY, align=PP_ALIGN.CENTER)

def slide_shifting(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Part 3: Performance & Shifting Analysis", "Analyzing element shifting operations on array insertions", n, total)
    
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.8),
             "Scenario: A task list application stores N elements in a contiguous dynamic array.",
             size=16, bold=True, color=NAVY)
             
    items = [
        ("1. Insertion Shifting: ", "If a user inserts a task at the front (index 0) of the array, how many elements must be shifted in memory? What is the time complexity of this operation?"),
        ("2. Deletion Shifting: ", "Compare the time complexity of: (a) popping an element from the very end of the array, and (b) deleting an element from index 0. Explain the difference in operations."),
        ("3. Memory Shrinking logic: ", "Our custom DynamicArray implementation shrinks the capacity to half when count <= capacity // 4. Why don't we shrink the capacity immediately when count < capacity // 2? (Hint: Think about alternating append and pop calls right at the boundary).")
    ]
    add_bullets(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(4.3), items, size=15, line_spacing=1.3)

def slide_extensions(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Part 4: Coding Extensions", "Advanced array challenges for study", n, total)
    
    items = [
        ("1. Merge Sorted Arrays In-Place: ", "Given two sorted integer arrays nums1 and nums2, merge nums2 into nums1 in-place as one sorted array. nums1 has a length of m + n (where m is elements in nums1 and n is elements in nums2).\nHint: Start filling the array from the back (right to left) using two pointers to avoid overwriting elements."),
        ("2. Sliding Window Maximum: ", "Given an array of integers nums and a sliding window size k, find the maximum element in each sliding window position.\nHint: Use a double-ended queue (deque) to store indices of useful elements in the window in decreasing order of value."),
        ("3. Move Zeroes In-Place: ", "Given an integer array nums, move all 0s to the end of it while maintaining the relative order of the non-zero elements. You must do this in-place without making a copy of the array.\nHint: Maintain a 'write' pointer for non-zero elements and fill the rest with zeroes.")
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.0), items, size=14, line_spacing=1.2)

# --- Build ---
slides = [
    slide_title,
    slide_conceptual_1,
    slide_conceptual_2,
    slide_trace_twosum,
    slide_shifting,
    slide_extensions
]

for i, f in enumerate(slides):
    if f == slide_title:
        f()
    else:
        f(i, len(slides)-1)

prs.save("questions.pptx")
print("Saved questions.pptx")
