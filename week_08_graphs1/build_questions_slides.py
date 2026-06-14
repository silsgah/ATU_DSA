import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# -------- Theme --------
NAVY        = RGBColor(0x0B, 0x2A, 0x4A)
TEAL        = RGBColor(0x12, 0x7A, 0x8A)
GOLD        = RGBColor(0xE0, 0xA8, 0x1F)
LIGHT_GREY  = RGBColor(0xF2, 0xF4, 0xF7)
DARK_GREY   = RGBColor(0x33, 0x3D, 0x4A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xC0, 0x39, 0x2B)
PURPLE      = RGBColor(0x6A, 0x1B, 0x9A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_GREY,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
              italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=18, color=DARK_GREY,
                bullet_color=TEAL, line_spacing=1.25, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        rb = p.add_run(); rb.text = "▸  "
        rb.font.name = font; rb.font.size = Pt(size)
        rb.font.bold = True; rb.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run(); r1.text = head
            r1.font.name = font; r1.font.size = Pt(size)
            r1.font.bold = True; r1.font.color.rgb = NAVY
            r2 = p.add_run(); r2.text = tail
            r2.font.name = font; r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = item
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb

def add_header(slide, title, subtitle=None, page_num=None, total=None):
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    add_rect(slide, Inches(0.35), 0, SLIDE_W - Inches(0.35), Inches(0.08), GOLD)
    add_text(slide, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.78), Inches(11.5), Inches(0.45),
                 subtitle, size=15, italic=True, color=TEAL)
    add_rect(slide, Inches(0.6), Inches(1.22), Inches(2.0), Inches(0.04), GOLD)
    add_text(slide, Inches(0.6), SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
             "Data Structures & Algorithms · Week 08 · Practice Questions",
             size=10, color=DARK_GREY, italic=True)
    if page_num is not None and total is not None:
        add_text(slide, SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.45),
                 Inches(1.2), Inches(0.3),
                 f"Slide {page_num} / {total}", size=10, color=DARK_GREY,
                 align=PP_ALIGN.RIGHT)

def draw_node(slide, cx, cy, label, color=WHITE, border_col=NAVY, r=Inches(0.4)):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.color.rgb = border_col; shp.line.width = Pt(2)
    shp.shadow.inherit = False
    add_text(slide, cx - r, cy - r, r * 2, r * 2,
             str(label), size=16, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def draw_edge(slide, x1, y1, x2, y2, color=DARK_GREY, weight=1.75):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    return conn

# ---------- Slides ----------

def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(5.5), SLIDE_W, Inches(0.12), GOLD)
    add_rect(s, 0, Inches(5.7), SLIDE_W, Inches(0.04), TEAL)
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.9), Inches(1.0),
                               Inches(2.8), Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb = GOLD
    chip.line.fill.background(); chip.shadow.inherit = False
    add_text(s, Inches(0.9), Inches(1.0), Inches(2.8), Inches(0.55),
             "WEEK 08 - WORKBOOK", size=16, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(1.4),
             "Graphs I: Practice Problems", size=58, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.7),
             "Self-Study Questions, Visual Exercises, and Modeling Scenarios",
             size=24, color=GOLD, italic=True)
    add_text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.6),
             "A Supplementary Workbook for Data Structures & Algorithms Students",
             size=18, color=LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(6.0), Inches(8), Inches(0.4),
             "Topic    ·   Graph Fundamentals & Traversals",
             size=14, color=WHITE)

def slide_conceptual_1(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Part 1: Conceptual Questions (1/2)", "Space complexity trade-offs in representations", n, total)
    
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.8),
             "Scenario: A social network has 10 million users (V = 10 million). On average, each user has 150 friends (E = 1.5 billion edges total).",
             size=16, bold=True, color=NAVY)
             
    items = [
        ("1. Adjacency Matrix Space: ", "If you use a V x V boolean grid (using 1 byte per cell), how much RAM is required? (Hint: V x V bytes = 10^14 bytes). What is this in Terabytes?"),
        ("2. Adjacency List Space: ", "If you store it as an Adjacency List (where each user ID and list reference node pointer takes 8 bytes), how much memory will it require? (Hint: Space complexity is O(V + E))."),
        ("3. Representation Choice: ", "Based on your calculations above, which data structure is appropriate for this real-world application? Explain the practical engineering implications.")
    ]
    add_bullets(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(4.3), items, size=15, line_spacing=1.3)

def slide_conceptual_2(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Part 1: Conceptual Questions (2/2)", "Understanding traversal algorithms & failure modes", n, total)
    
    items = [
        ("1. BFS vs DFS Shortest Path: ", "Explain why BFS is guaranteed to find the shortest path (minimum edge count) between two vertices in an unweighted graph, whereas DFS is not. Can you sketch a simple 4-node graph where DFS finds a sub-optimal path?"),
        ("2. The Visited Set Failure: ", "Describe exactly what happens if you run a BFS or DFS on a cyclic graph but forget to track visited nodes. Which one will fail with a 'Stack Overflow' error first, and which one will consume all system memory?"),
        ("3. Recursion vs Iteration: ", "Recursive DFS uses the Call Stack. If a graph is structured as a single long line of V vertices, what is the maximum recursion depth, and what problem does this present in languages like Python?")
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.0), items, size=15, line_spacing=1.3)

def slide_trace_bfs(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Part 2: Tracing Breadth-First Search (BFS)", "Hand-trace the queue execution step-by-step", n, total)
    
    add_text(s, Inches(0.9), Inches(1.5), Inches(6.0), Inches(0.8),
             "Exercise: Trace BFS starting at node 'A'.\nAssume neighbors are processed in alphabetical order.",
             size=14, italic=True, color=DARK_GREY)
             
    # Table template for tracing
    add_rect(s, Inches(0.9), Inches(2.3), Inches(6.0), Inches(0.45), NAVY)
    add_text(s, Inches(0.9), Inches(2.35), Inches(1.2), Inches(0.3), "Step", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.1), Inches(2.35), Inches(1.2), Inches(0.3), "Pop", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.3), Inches(2.35), Inches(1.8), Inches(0.3), "Queue Status", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(5.1), Inches(2.35), Inches(1.8), Inches(0.3), "Visited Set", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    for i in range(7):
        y_off = Inches(2.75 + i * 0.5)
        bg = WHITE if i % 2 == 0 else LIGHT_GREY
        add_rect(s, Inches(0.9), y_off, Inches(6.0), Inches(0.5), bg, line=DARK_GREY)
        # initial state hints
        if i == 0:
            add_text(s, Inches(0.9), y_off + Inches(0.1), Inches(1.2), Inches(0.3), "0 (Start)", size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)
            add_text(s, Inches(2.1), y_off + Inches(0.1), Inches(1.2), Inches(0.3), "-", size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)
            add_text(s, Inches(3.3), y_off + Inches(0.1), Inches(1.8), Inches(0.3), "['A']", size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)
            add_text(s, Inches(5.1), y_off + Inches(0.1), Inches(1.8), Inches(0.3), "{'A'}", size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)
        else:
            add_text(s, Inches(0.9), y_off + Inches(0.1), Inches(1.2), Inches(0.3), str(i), size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)

    # Right: Graph Visual
    add_rect(s, Inches(7.5), Inches(1.5), Inches(5.0), Inches(5.0), LIGHT_GREY)
    add_text(s, Inches(7.5), Inches(1.6), Inches(5.0), Inches(0.5), "Target Graph", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    
    pos = {
        'A': (Inches(10.0), Inches(2.3)),
        'B': (Inches(8.5), Inches(3.6)),
        'C': (Inches(11.5), Inches(3.6)),
        'D': (Inches(8.5), Inches(5.1)),
        'E': (Inches(11.5), Inches(5.1)),
        'F': (Inches(10.0), Inches(4.3)),
        'G': (Inches(10.0), Inches(5.8))
    }
    edges = [('A','B'), ('A','C'), ('B','D'), ('B','E'), ('C','F'), ('E','F'), ('F','G')]
    
    for (u, v) in edges:
        x1, y1 = pos[u]; x2, y2 = pos[v]
        draw_edge(s, x1, y1, x2, y2, color=TEAL, weight=2)
        
    for label, (cx, cy) in pos.items():
        draw_node(s, cx, cy, label)

def slide_trace_dfs(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Part 2: Tracing Depth-First Search (DFS)", "Hand-trace the recursion call stack step-by-step", n, total)
    
    add_text(s, Inches(0.9), Inches(1.5), Inches(6.0), Inches(0.8),
             "Exercise: Trace DFS starting at node 'A'.\nAssume neighbors are visited in alphabetical order.",
             size=14, italic=True, color=DARK_GREY)
             
    # Table template for tracing
    add_rect(s, Inches(0.9), Inches(2.3), Inches(6.0), Inches(0.45), PURPLE)
    add_text(s, Inches(0.9), Inches(2.35), Inches(1.2), Inches(0.3), "Step", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.1), Inches(2.35), Inches(1.2), Inches(0.3), "Action", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.3), Inches(2.35), Inches(1.8), Inches(0.3), "Call Stack Trace", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(5.1), Inches(2.35), Inches(1.8), Inches(0.3), "Visited Set", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    for i in range(7):
        y_off = Inches(2.75 + i * 0.5)
        bg = WHITE if i % 2 == 0 else LIGHT_GREY
        add_rect(s, Inches(0.9), y_off, Inches(6.0), Inches(0.5), bg, line=DARK_GREY)
        if i == 0:
            add_text(s, Inches(0.9), y_off + Inches(0.1), Inches(1.2), Inches(0.3), "0 (Start)", size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)
            add_text(s, Inches(2.1), y_off + Inches(0.1), Inches(1.2), Inches(0.3), "Call dfs(A)", size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)
            add_text(s, Inches(3.3), y_off + Inches(0.1), Inches(1.8), Inches(0.3), "[dfs(A)]", size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)
            add_text(s, Inches(5.1), y_off + Inches(0.1), Inches(1.8), Inches(0.3), "{'A'}", size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)
        else:
            add_text(s, Inches(0.9), y_off + Inches(0.1), Inches(1.2), Inches(0.3), str(i), size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)

    # Right: Graph Visual (identical to BFS trace for direct comparison)
    add_rect(s, Inches(7.5), Inches(1.5), Inches(5.0), Inches(5.0), LIGHT_GREY)
    add_text(s, Inches(7.5), Inches(1.6), Inches(5.0), Inches(0.5), "Target Graph", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    
    pos = {
        'A': (Inches(10.0), Inches(2.3)),
        'B': (Inches(8.5), Inches(3.6)),
        'C': (Inches(11.5), Inches(3.6)),
        'D': (Inches(8.5), Inches(5.1)),
        'E': (Inches(11.5), Inches(5.1)),
        'F': (Inches(10.0), Inches(4.3)),
        'G': (Inches(10.0), Inches(5.8))
    }
    edges = [('A','B'), ('A','C'), ('B','D'), ('B','E'), ('C','F'), ('E','F'), ('F','G')]
    
    for (u, v) in edges:
        x1, y1 = pos[u]; x2, y2 = pos[v]
        draw_edge(s, x1, y1, x2, y2, color=GOLD, weight=2)
        
    for label, (cx, cy) in pos.items():
        draw_node(s, cx, cy, label)

def slide_modeling(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Part 3: Real-World Modeling Challenges", "How do we represent system logic as a graph schema?", n, total)
    
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.8),
             "Exercise: For each scenario, specify: (a) Vertices, (b) Edges, (c) Directed/Undirected, (d) Weighted/Unweighted.",
             size=15, italic=True, color=DARK_GREY)
             
    items = [
        ("1. Google Maps Navigation: ", "Finding the fastest route between two cities.\n(Vertices = Cities/Intersections, Edges = Highways/Roads, Directed, Weighted by travel time/distance)"),
        ("2. LinkedIn Connections: ", "Finding mutual connections and suggesting new friends.\n(Vertices = Users, Edges = Connections, Undirected, Unweighted)"),
        ("3. Academic Prerequisites: ", "Determining which classes a student must take in order.\n(Vertices = Courses, Edges = Prerequisites, Directed (one-way prerequisite), Unweighted)"),
        ("4. Web Search Crawler: ", "Finding and index-ranking pages by traversing links.\n(Vertices = Webpages, Edges = Hyperlinks on pages, Directed, Unweighted)")
    ]
    add_bullets(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(4.5), items, size=14, line_spacing=1.2)

def slide_extensions(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Part 4: Coding Extensions", "Extend the code in practice.py to master traversals", n, total)
    
    items = [
        ("1. Full Path Reconstruction: ", "Modify shortest_path(graph, start, target) so that instead of just returning the distance integer, it returns the actual list of nodes visited along the optimal path (e.g., ['A', 'C', 'F', 'G']).\nHint: Maintain a parent dictionary {child: parent} during traversal."),
        ("2. Cycle Detection (DFS): ", "Write a function has_cycle(graph) that returns True if an undirected graph contains a loop, and False otherwise.\nHint: During DFS, if you visit a node that is already marked as visited and is not the immediate parent of the current node, you have found a cycle."),
        ("3. Connected Components Count: ", "Write a function count_components(graph) that returns the total count of isolated graph networks.\nHint: Loop through all vertices and trigger a DFS/BFS whenever you hit an unvisited node, incrementing the component count each time.")
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.0), items, size=14, line_spacing=1.2)

# --- Build ---
slides = [
    slide_title,
    slide_conceptual_1,
    slide_conceptual_2,
    slide_trace_bfs,
    slide_trace_dfs,
    slide_modeling,
    slide_extensions
]

for i, f in enumerate(slides):
    if f == slide_title:
        f()
    else:
        f(i, len(slides)-1)

prs.save("questions.pptx")
print("Saved questions.pptx")
