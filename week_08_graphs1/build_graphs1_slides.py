"""
Generate a university-grade PowerPoint deck for Week 08: Graphs I (Fundamentals).

Output: graphs1.pptx (16:9 widescreen)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# -------- Theme --------
NAVY        = RGBColor(0x0B, 0x2A, 0x4A)
TEAL        = RGBColor(0x12, 0x7A, 0x8A)
GOLD        = RGBColor(0xE0, 0xA8, 0x1F)
LIGHT_GREY  = RGBColor(0xF2, 0xF4, 0xF7)
DARK_GREY   = RGBColor(0x33, 0x3D, 0x4A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xC0, 0x39, 0x2B)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
PURPLE      = RGBColor(0x6A, 0x1B, 0x9A)
CODE_BG     = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG     = RGBColor(0xE6, 0xE6, 0xE6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ---------- Generic helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=18, color=DARK_GREY,
                bullet_color=TEAL, line_spacing=1.25, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        rb = p.add_run(); rb.text = "▸  "
        rb.font.name = font; rb.font.size = Pt(size)
        rb.font.bold = True; rb.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run(); r1.text = head
            r1.font.name = font; r1.font.size = Pt(size)
            r1.font.bold = True; r1.font.color.rgb = NAVY
            r2 = p.add_run(); r2.text = tail
            r2.font.name = font; r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = item
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb

def add_header(slide, title, subtitle=None, page_num=None, total=None):
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    add_rect(slide, Inches(0.35), 0, SLIDE_W - Inches(0.35), Inches(0.08), GOLD)
    add_text(slide, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.78), Inches(11.5), Inches(0.45),
                 subtitle, size=15, italic=True, color=TEAL)
    add_rect(slide, Inches(0.6), Inches(1.22), Inches(2.0), Inches(0.04), GOLD)
    add_text(slide, Inches(0.6), SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
             "Data Structures & Algorithms · Week 08 · Graphs I",
             size=10, color=DARK_GREY, italic=True)
    if page_num is not None and total is not None:
        add_text(slide, SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.45),
                 Inches(1.2), Inches(0.3),
                 f"Slide {page_num} / {total}", size=10, color=DARK_GREY,
                 align=PP_ALIGN.RIGHT)

def add_code_block(slide, x, y, w, h, code, size=14):
    add_rect(slide, x, y, w, h, CODE_BG)
    add_rect(slide, x, y, Inches(0.08), h, GOLD)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.25), h - Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, ln in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.line_spacing = 1.15
        r = p.add_run(); r.text = ln if ln else " "
        r.font.name = "Consolas"; r.font.size = Pt(size)
        r.font.color.rgb = CODE_FG

# ---------- Node Drawing Helpers ----------
def draw_node(slide, cx, cy, label, color=WHITE, border_col=NAVY, r=Inches(0.4)):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.color.rgb = border_col; shp.line.width = Pt(2)
    shp.shadow.inherit = False
    add_text(slide, cx - r, cy - r, r * 2, r * 2,
             str(label), size=16, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def draw_edge(slide, x1, y1, x2, y2, color=DARK_GREY, weight=1.75, head=False):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if head:
        ln = conn.line._get_or_add_ln()
        etree.SubElement(ln, qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    return conn

# ---------- Slides ----------
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(5.5), SLIDE_W, Inches(0.12), GOLD)
    add_rect(s, 0, Inches(5.7), SLIDE_W, Inches(0.04), TEAL)
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.9), Inches(1.0),
                              Inches(2.4), Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb = GOLD
    chip.line.fill.background(); chip.shadow.inherit = False
    add_text(s, Inches(0.9), Inches(1.0), Inches(2.4), Inches(0.55),
             "WEEK 08", size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(1.4),
             "Graphs I: Fundamentals", size=58, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.7),
             "Representation, Traversal, and Network Logic",
             size=24, color=GOLD, italic=True)
    add_text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.6),
             "A University-Grade Lecture in Data Structures & Algorithms",
             size=18, color=LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(6.0), Inches(8), Inches(0.4),
             "Course   ·   CSC: Data Structures & Algorithms",
             size=14, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.4), Inches(8), Inches(0.4),
             "Module   ·   Week 08 of 12",
             size=14, color=WHITE)

def slide_outline(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Lecture Outline",
               "The foundation of modern networks", n, total)
    items = [
        ("1. What is a Graph?  ", "From Trees to free-form networks."),
        ("2. Graph Terminology.  ", "Vertices, Edges, Directed vs Undirected, Cycles."),
        ("3. Representing Graphs in Code.  ", "Adjacency Lists vs Adjacency Matrices."),
        ("4. Breadth-First Search (BFS).  ", "Concept, pseudocode, and execution tracing."),
        ("5. Depth-First Search (DFS).  ", "Concept, recursion, and call stack simulation.")
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(12), Inches(5.5),
                items, size=18, line_spacing=1.32)

def slide_intro(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "1. What is a Graph?",
               "The generalization of Trees and Linked Lists", n, total)
               
    add_text(s, Inches(0.9), Inches(1.5), Inches(6.0), Inches(4.0),
             "A Tree is just a highly restricted, directed, acyclic Graph!\n\n"
             "When we remove the strict hierarchical parent-child rules, we get a Graph:\n"
             "• Any node can connect to any other node.\n"
             "• There can be loops/cycles.\n"
             "• There is no inherent 'Root' node.\n\n"
             "Mathematically: G = (V, E) where V is a set of Vertices and E is a set of Edges.",
             size=16, color=DARK_GREY)
             
    # Right: small visual of a graph
    add_rect(s, Inches(7.5), Inches(1.5), Inches(5.0), Inches(4.5), LIGHT_GREY)
    
    pos = {
        'A': (Inches(8.5), Inches(2.5)),
        'B': (Inches(11.0), Inches(2.5)),
        'C': (Inches(8.5), Inches(4.5)),
        'D': (Inches(11.0), Inches(4.5)),
        'E': (Inches(9.75), Inches(3.5))
    }
    edges = [('A','B'), ('A','C'), ('B','E'), ('C','E'), ('C','D'), ('E','D'), ('B','D')]
    
    for (u, v) in edges:
        x1, y1 = pos[u]
        x2, y2 = pos[v]
        draw_edge(s, x1, y1, x2, y2, color=TEAL, weight=2)
        
    for label, (cx, cy) in pos.items():
        draw_node(s, cx, cy, label)

def slide_terminology(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "2. Graph Terminology",
               "The vocabulary of network data", n, total)
               
    items = [
        ("Vertex (Node): ", "The entities in the network (e.g., a person, a city, a webpage)."),
        ("Edge (Link): ", "The connection between two vertices."),
        ("Directed Graph: ", "Edges are one-way streets (e.g., following someone on Twitter)."),
        ("Undirected Graph: ", "Edges are two-way roads (e.g., mutual friends on Facebook)."),
        ("Weighted Graph: ", "Edges have a cost or distance associated with them (e.g., Google Maps travel time)."),
        ("Cyclic Graph: ", "Contains at least one path that loops back to the same vertex."),
        ("Acyclic Graph: ", "Impossible to walk in a continuous loop.")
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.5), items, size=16, line_spacing=1.3)

def slide_representation(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "3. Representing Graphs in Code",
               "How do we store this complex structure?", n, total)
               
    # Left: Adj List
    add_rect(s, Inches(0.9), Inches(1.5), Inches(5.5), Inches(4.5), LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(1.6), Inches(5.5), Inches(0.5), "Adjacency List", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.1), Inches(2.2), Inches(5.1), Inches(1.5),
             "A Hash Map (dictionary) where Key = Vertex, Value = List of Neighbors.\n"
             "• Space: O(V + E)\n"
             "• Ideal for sparse graphs (most real-world graphs).", size=14, color=DARK_GREY)
    
    code_list = (
        "adj_list = {\n"
        "  'A': ['B', 'C'],\n"
        "  'B': ['A', 'D', 'E'],\n"
        "  'C': ['A', 'F']\n"
        "}"
    )
    add_code_block(s, Inches(1.1), Inches(3.5), Inches(5.1), Inches(2.0), code_list, size=14)
    
    # Right: Adj Matrix
    add_rect(s, Inches(6.9), Inches(1.5), Inches(5.5), Inches(4.5), LIGHT_GREY)
    add_text(s, Inches(6.9), Inches(1.6), Inches(5.5), Inches(0.5), "Adjacency Matrix", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.1), Inches(2.2), Inches(5.1), Inches(1.5),
             "A 2D array of size V x V where matrix[i][j] = 1 if an edge exists, else 0.\n"
             "• Space: O(V²)\n"
             "• Fast edge lookups O(1), but wastes massive RAM for sparse networks.", size=14, color=DARK_GREY)
             
    code_mat = (
        "# A  B  C  D\n"
        "[[0, 1, 1, 0], # A\n"
        " [1, 0, 0, 1], # B\n"
        " [1, 0, 0, 1], # C\n"
        " [0, 1, 1, 0]] # D"
    )
    add_code_block(s, Inches(7.1), Inches(3.5), Inches(5.1), Inches(2.0), code_mat, size=14)

def slide_bfs_concept(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "4. Breadth-First Search (BFS): Concept",
               "Exploring level by level like a ripple in a pond", n, total)
    
    items = [
        ("The Level-Order Approach: ", "BFS starts at a source node and visits all of its direct neighbors first (Level 1), then all of their neighbors (Level 2), and so on."),
        ("FIFO Queue Mechanism: ", "It uses a Queue to keep track of nodes to explore next. First In, First Out ensures we explore closer nodes before further ones."),
        ("Shortest Path Guarantee: ", "In an unweighted graph, the first time BFS reaches a node, it has found the shortest path (minimum edge count) to that node."),
        ("Visited Set: ", "Crucial to keep track of visited vertices to avoid infinite loops and re-processing nodes in cyclic graphs."),
        ("Complexity: ", "Time: O(V + E) as every vertex is visited and every edge is checked. Space: O(V) for the queue and visited set.")
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(6.0), Inches(5.0), items, size=15, line_spacing=1.2)
    
    # Right: Visual Graph with levels
    add_rect(s, Inches(7.5), Inches(1.5), Inches(5.0), Inches(5.0), LIGHT_GREY)
    add_text(s, Inches(7.5), Inches(1.6), Inches(5.0), Inches(0.5), "BFS Layer Traversal Visual", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    
    pos = {
        'A': (Inches(10.0), Inches(2.3)),   # Source (Level 0)
        'B': (Inches(8.5), Inches(3.6)),    # Level 1
        'C': (Inches(11.5), Inches(3.6)),   # Level 1
        'D': (Inches(8.5), Inches(5.1)),    # Level 2
        'E': (Inches(11.5), Inches(5.1))    # Level 2
    }
    edges = [('A','B'), ('A','C'), ('B','D'), ('C','E')]
    
    for (u, v) in edges:
        x1, y1 = pos[u]; x2, y2 = pos[v]
        draw_edge(s, x1, y1, x2, y2, color=TEAL, weight=2)
        
    for label, (cx, cy) in pos.items():
        col = GOLD if label == 'A' else (TEAL if label in ['B', 'C'] else WHITE)
        draw_node(s, cx, cy, label, color=col)
        
    add_text(s, Inches(7.5), Inches(5.8), Inches(5.0), Inches(0.6),
             "Level 0 (Gold)  ➔  Level 1 (Teal)  ➔  Level 2 (White)",
             size=13, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)

def slide_bfs_code(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "4. Breadth-First Search (BFS): Code & Trace",
               "Implementing and executing the BFS algorithm", n, total)
               
    # Left: Code Block
    code = (
        "from collections import deque\n\n"
        "def bfs(graph, start):\n"
        "    visited = {start}\n"
        "    q = deque([start])\n\n"
        "    while q:\n"
        "        node = q.popleft()\n"
        "        print(node)  # Process node\n\n"
        "        for neighbor in graph[node]:\n"
        "            if neighbor not in visited:\n"
        "                visited.add(neighbor)\n"
        "                q.append(neighbor)"
    )
    add_code_block(s, Inches(0.9), Inches(1.5), Inches(6.0), Inches(5.0), code, size=13)
    
    # Right: Dry-Run / Trace Table
    add_rect(s, Inches(7.5), Inches(1.5), Inches(5.0), Inches(5.0), LIGHT_GREY)
    add_text(s, Inches(7.5), Inches(1.6), Inches(5.0), Inches(0.5), "Step-by-Step Execution Trace", size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
    
    add_text(s, Inches(7.7), Inches(2.2), Inches(4.6), Inches(0.6),
             "Graph: A-B, A-C, B-D, C-E (Start at 'A')", size=13, bold=True, color=DARK_GREY)
             
    # Table headers
    add_rect(s, Inches(7.7), Inches(2.9), Inches(4.6), Inches(0.4), NAVY)
    add_text(s, Inches(7.7), Inches(2.95), Inches(1.0), Inches(0.3), "Pop Node", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.7), Inches(2.95), Inches(1.8), Inches(0.3), "Queue Status", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(10.5), Inches(2.95), Inches(1.8), Inches(0.3), "Visited Set", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    # Rows
    trace_data = [
        ("Initial", "['A']", "{'A'}"),
        ("A", "['B', 'C']", "{'A', 'B', 'C'}"),
        ("B", "['C', 'D']", "{'A', 'B', 'C', 'D'}"),
        ("C", "['D', 'E']", "{'A', 'B', 'C', 'D', 'E'}"),
        ("D", "['E']", "{'A', 'B', 'C', 'D', 'E'}"),
        ("E", "[]", "{'A', 'B', 'C', 'D', 'E'}")
    ]
    
    for idx, (pop, q_state, vis) in enumerate(trace_data):
        y_offset = Inches(3.4 + idx * 0.45)
        bg_col = WHITE if idx % 2 == 0 else LIGHT_GREY
        add_rect(s, Inches(7.7), y_offset, Inches(4.6), Inches(0.4), bg_col, line=DARK_GREY)
        add_text(s, Inches(7.7), y_offset + Inches(0.05), Inches(1.0), Inches(0.3), pop, size=11, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)
        add_text(s, Inches(8.7), y_offset + Inches(0.05), Inches(1.8), Inches(0.3), q_state, size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)
        add_text(s, Inches(10.5), y_offset + Inches(0.05), Inches(1.8), Inches(0.3), vis, size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)

def slide_dfs_concept(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "5. Depth-First Search (DFS): Concept",
               "Plunging deep down a path until hitting a dead end", n, total)
               
    items = [
        ("The Depth-First Approach: ", "DFS starts at a source node and follows a branch as far as possible before backtracking. It prioritizes depth over breadth."),
        ("LIFO Stack / Recursion: ", "DFS uses a Stack structure. This is naturally implemented via Recursion (the Call Stack), but can also be done iteratively with a custom Stack."),
        ("Cycle Detection: ", "Excellent for detecting cycles (loops) in graphs by tracking the path of the current recursive call."),
        ("Topological Sorting: ", "Crucial for ordering tasks with dependencies (e.g. build systems, class prerequisites) using finish times."),
        ("Complexity: ", "Time: O(V + E) as every vertex is visited and every edge is checked. Space: O(V) for the recursion call stack in the worst case.")
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(6.0), Inches(5.0), items, size=15, line_spacing=1.2)
    
    # Right: Visual Graph with levels
    add_rect(s, Inches(7.5), Inches(1.5), Inches(5.0), Inches(5.0), LIGHT_GREY)
    add_text(s, Inches(7.5), Inches(1.6), Inches(5.0), Inches(0.5), "DFS Traversal Order Visual", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    
    pos = {
        'A': (Inches(10.0), Inches(2.3)),
        'B': (Inches(8.5), Inches(3.6)),
        'C': (Inches(11.5), Inches(3.6)),
        'D': (Inches(8.5), Inches(5.1)),
        'E': (Inches(11.5), Inches(5.1))
    }
    edges = [('A','B'), ('B','D'), ('A','C'), ('C','E')]
    
    for (u, v) in edges:
        x1, y1 = pos[u]; x2, y2 = pos[v]
        draw_edge(s, x1, y1, x2, y2, color=GOLD, weight=2)
        
    order = ['A', 'B', 'D', 'C', 'E']
    for label, (cx, cy) in pos.items():
        col = PURPLE if label == 'A' else (GOLD if label in ['B', 'D'] else WHITE)
        draw_node(s, cx, cy, label, color=col)
        
    add_text(s, Inches(7.5), Inches(5.8), Inches(5.0), Inches(0.6),
             "First Branch (Purple ➔ Gold)  ➔  Second Branch (White)",
             size=13, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)

def slide_dfs_code(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "5. Depth-First Search (DFS): Code & Trace",
               "Implementing and executing the recursive DFS", n, total)
               
    # Left: Code Block
    code = (
        "def dfs(graph, node, visited=None):\n"
        "    if visited is None:\n"
        "        visited = set()\n\n"
        "    if node not in visited:\n"
        "        visited.add(node)\n"
        "        print(node)  # Process node\n\n"
        "        for neighbor in graph[node]:\n"
        "            if neighbor not in visited:\n"
        "                dfs(graph, neighbor, visited)\n"
        "    return visited"
    )
    add_code_block(s, Inches(0.9), Inches(1.5), Inches(6.0), Inches(5.0), code, size=13)
    
    # Right: Call Stack Simulation
    add_rect(s, Inches(7.5), Inches(1.5), Inches(5.0), Inches(5.0), LIGHT_GREY)
    add_text(s, Inches(7.5), Inches(1.6), Inches(5.0), Inches(0.5), "Recursion Call Stack Trace", size=16, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    
    add_text(s, Inches(7.7), Inches(2.2), Inches(4.6), Inches(0.6),
             "Graph: A-B, B-D, A-C, C-E (Start at 'A')", size=13, bold=True, color=DARK_GREY)
             
    # Table headers
    add_rect(s, Inches(7.7), Inches(2.9), Inches(4.6), Inches(0.4), NAVY)
    add_text(s, Inches(7.7), Inches(2.95), Inches(1.0), Inches(0.3), "Action", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.7), Inches(2.95), Inches(1.8), Inches(0.3), "Call Stack", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(10.5), Inches(2.95), Inches(1.8), Inches(0.3), "Visited Set", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    # Rows
    trace_data = [
        ("Call dfs(A)", "[dfs(A)]", "{'A'}"),
        ("Call dfs(B)", "[dfs(A), dfs(B)]", "{'A', 'B'}"),
        ("Call dfs(D)", "[dfs(A), dfs(B), dfs(D)]", "{'A', 'B', 'D'}"),
        ("Pop dfs(D)", "[dfs(A), dfs(B)]", "{'A', 'B', 'D'}"),
        ("Pop dfs(B)", "[dfs(A)]", "{'A', 'B', 'D'}"),
        ("Call dfs(C)", "[dfs(A), dfs(C)]", "{'A', 'B', 'D', 'C'}"),
        ("Call dfs(E)", "[dfs(A), dfs(C), dfs(E)]", "{'A', 'B', 'D', 'C', 'E'}")
    ]
    
    for idx, (action, stack_state, vis) in enumerate(trace_data):
        y_offset = Inches(3.4 + idx * 0.45)
        bg_col = WHITE if idx % 2 == 0 else LIGHT_GREY
        add_rect(s, Inches(7.7), y_offset, Inches(4.6), Inches(0.4), bg_col, line=DARK_GREY)
        add_text(s, Inches(7.7), y_offset + Inches(0.05), Inches(1.0), Inches(0.3), action, size=11, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)
        add_text(s, Inches(8.7), y_offset + Inches(0.05), Inches(1.8), Inches(0.3), stack_state, size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)
        add_text(s, Inches(10.5), y_offset + Inches(0.05), Inches(1.8), Inches(0.3), vis, size=11, color=DARK_GREY, align=PP_ALIGN.CENTER)

def slide_summary(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Summary",
               "Graphs open up the rest of Computer Science", n, total)
               
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.0),
             "Graphs are the most versatile data structure. Mastering BFS and DFS is mandatory for complex problem solving.",
             size=18, bold=True, color=NAVY)
             
    items = [
        "Use Adjacency Lists almost always, unless the graph is extremely dense.",
        "Need the Shortest Path in an unweighted graph? Use BFS.",
        "Need to check if a path exists, or find cycles? DFS is usually cleaner to implement recursively.",
        "Always remember to track Visited nodes, otherwise cycles will cause infinite loops (Stack Overflow or Memory Exceeded)."
    ]
    add_bullets(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(3.0), items, size=16, line_spacing=1.4)
    
    add_rect(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.0), TEAL)
    add_text(s, Inches(1.1), Inches(5.6), Inches(11.0), Inches(0.8),
             "Next Week (Graphs II): We will introduce edge weights and explore Dijkstra's Algorithm for finding the shortest paths on real-world maps.",
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# --- Build ---
slides = [
    slide_title,
    slide_outline,
    slide_intro,
    slide_terminology,
    slide_representation,
    slide_bfs_concept,
    slide_bfs_code,
    slide_dfs_concept,
    slide_dfs_code,
    slide_summary
]

for i, f in enumerate(slides):
    if f == slide_title:
        f()
    else:
        f(i, len(slides)-1)

prs.save("graphs1.pptx")
print("Saved graphs1.pptx")
