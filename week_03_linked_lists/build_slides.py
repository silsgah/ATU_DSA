"""
Generate a university-grade PowerPoint deck for Week 03: Linked Lists.

Output: linked_lists_dsa_final.pptx (16:9 widescreen)
Author : DSA Curriculum
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree
import os

# -------- Theme --------
NAVY        = RGBColor(0x0B, 0x2A, 0x4A)   # deep navy (primary)
TEAL        = RGBColor(0x12, 0x7A, 0x8A)   # accent
GOLD        = RGBColor(0xE0, 0xA8, 0x1F)   # accent
LIGHT_GREY  = RGBColor(0xF2, 0xF4, 0xF7)   # background panels
DARK_GREY   = RGBColor(0x33, 0x3D, 0x4A)   # body text
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xC0, 0x39, 0x2B)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
CODE_BG     = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG     = RGBColor(0xE6, 0xE6, 0xE6)

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=18, color=DARK_GREY,
                bullet_color=TEAL, line_spacing=1.25, font="Calibri"):
    """items: list of strings or (label_bold, rest) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        # bullet glyph
        rb = p.add_run()
        rb.text = "▸  "
        rb.font.name = font
        rb.font.size = Pt(size)
        rb.font.bold = True
        rb.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run()
            r1.text = head
            r1.font.name = font
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = tail
            r2.font.name = font
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None, page_num=None, total=None):
    # left navy stripe
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    # top thin gold rule
    add_rect(slide, Inches(0.35), Inches(0.0), SLIDE_W - Inches(0.35), Inches(0.08), GOLD)
    # title
    add_text(slide, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.78), Inches(11.5), Inches(0.45),
                 subtitle, size=15, italic=True, color=TEAL)
    # divider
    add_rect(slide, Inches(0.6), Inches(1.22), Inches(2.0), Inches(0.04), GOLD)
    # footer brand
    add_text(slide, Inches(0.6), SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
             "Data Structures & Algorithms · Week 03 · Linked Lists",
             size=10, color=DARK_GREY, italic=True)
    if page_num is not None and total is not None:
        add_text(slide, SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.45),
                 Inches(1.2), Inches(0.3),
                 f"Slide {page_num} / {total}", size=10, color=DARK_GREY,
                 align=PP_ALIGN.RIGHT)


def add_code_block(slide, x, y, w, h, code, size=14):
    add_rect(slide, x, y, w, h, CODE_BG)
    # left accent
    add_rect(slide, x, y, Inches(0.08), h, GOLD)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.25), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = code.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = ln if ln else " "
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = CODE_FG


def add_node_box(slide, x, y, w, h, value, *, fill=WHITE, border=NAVY,
                 value_size=20, has_next=True, has_prev=False):
    """Draw a linked-list node: [data | next] (or [prev | data | next])."""
    # outer
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    outer.fill.solid()
    outer.fill.fore_color.rgb = fill
    outer.line.color.rgb = border
    outer.line.width = Pt(2)
    outer.shadow.inherit = False
    # split
    parts = (1 if has_prev else 0) + 1 + (1 if has_next else 0)
    cell_w = w / parts
    idx = 0
    if has_prev:
        # prev cell
        add_rect(slide, x + Emu(0), y, cell_w, h, LIGHT_GREY, line=border)
        add_text(slide, x, y, cell_w, h, "prev", size=12, color=DARK_GREY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        idx += 1
    # data cell
    add_text(slide, x + cell_w * idx, y, cell_w, h, str(value),
             size=value_size, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if has_next:
        idx2 = idx + 1
        add_rect(slide, x + cell_w * idx2, y, cell_w, h, LIGHT_GREY, line=border)
        add_text(slide, x + cell_w * idx2, y, cell_w, h, "next",
                 size=12, color=DARK_GREY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # vertical separator(s)
    for k in range(1, parts):
        sep_x = x + cell_w * k
        line = slide.shapes.add_connector(1, sep_x, y, sep_x, y + h)
        line.line.color.rgb = border
        line.line.width = Pt(1.25)


def add_arrow(slide, x1, y1, x2, y2, color=TEAL, weight=2.25):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    # add arrowhead via XML
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(
        ln, qn("a:tailEnd"),
        {"type": "triangle", "w": "med", "len": "med"}
    )
    return conn


def add_label(slide, x, y, w, h, text, *, color=DARK_GREY, size=12, bold=False,
              align=PP_ALIGN.CENTER):
    return add_text(slide, x, y, w, h, text, size=size, bold=bold,
                    color=color, align=align, anchor=MSO_ANCHOR.MIDDLE)


# ---------- Slide builders ----------
def slide_title():
    s = prs.slides.add_slide(BLANK)
    # full background
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # diagonal accent
    add_rect(s, 0, Inches(5.5), SLIDE_W, Inches(0.12), GOLD)
    add_rect(s, 0, Inches(5.7), SLIDE_W, Inches(0.04), TEAL)
    # week chip
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.9), Inches(1.0),
                              Inches(2.4), Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb = GOLD
    chip.line.fill.background(); chip.shadow.inherit = False
    add_text(s, Inches(0.9), Inches(1.0), Inches(2.4), Inches(0.55),
             "WEEK 03", size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # title
    add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(1.4),
             "Linked Lists", size=64, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.7),
             "Nodes, Pointers, and Foundational Algorithms",
             size=26, color=GOLD, italic=True)
    # caption
    add_text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.6),
             "A University-Grade Lecture in Data Structures & Algorithms",
             size=18, color=LIGHT_GREY)
    # bottom block
    add_text(s, Inches(0.9), Inches(6.0), Inches(8), Inches(0.4),
             "Course   ·   CSC: Data Structures & Algorithms",
             size=14, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.4), Inches(8), Inches(0.4),
             "Module   ·   Week 03 of 12",
             size=14, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.8), Inches(8), Inches(0.4),
             "Format   ·   Lecture · Demonstration · Practical Lab",
             size=14, color=WHITE)


def slide_outline(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Lecture Outline",
               "What we will cover today and how it builds your engineering toolkit",
               n, total)
    items = [
        ("1. Motivation.  ", "Why linked structures exist; problems with arrays."),
        ("2. The Node.  ", "Memory model, fields, and intuition."),
        ("3. Singly · Doubly · Circular.  ", "Three canonical variants."),
        ("4. Core Operations.  ", "Append, prepend, insert, delete, search."),
        ("5. Complexity.  ", "Big-O analysis vs. arrays."),
        ("6. Implementation.  ", "Production-quality Python walkthrough."),
        ("7. Reversal Algorithm.  ", "Pointer juggling without losing the chain."),
        ("8. Floyd's Tortoise & Hare.  ", "Cycle detection in O(1) memory."),
        ("9. Pitfalls & Best Practice.  ", "Where students lose marks—and points."),
        ("10. Applications & Lab.  ", "Real systems and the practice assignment."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(12), Inches(5.5),
                items, size=18, line_spacing=1.35)


def slide_objectives(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Learning Outcomes",
               "By the end of this lecture you will be able to…", n, total)
    objs = [
        "Explain how linked memory differs from contiguous memory and articulate the engineering trade-offs.",
        "Define and implement a Node class and a generic Singly Linked List.",
        "Distinguish singly, doubly, and circular variants and choose the right one.",
        "Reason precisely about asymptotic complexity for each operation.",
        "Implement in-place reversal and cycle detection (Floyd's algorithm) from first principles.",
        "Identify common pointer pitfalls and describe defensive coding patterns to avoid them.",
        "Map linked-list techniques onto real-world systems: undo stacks, LRU caches, music queues, allocators.",
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(12), Inches(5.5),
                objs, size=18, line_spacing=1.35)


def slide_motivation(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "1.  Motivation: Why Not Just Use Arrays?",
               "Every data structure is born from a constraint",
               n, total)
    # left text
    add_text(s, Inches(0.9), Inches(1.45), Inches(6.0), Inches(0.5),
             "Arrays are wonderful — until they aren't.",
             size=20, bold=True, color=NAVY)
    items = [
        "Fixed capacity: insertions in the middle force O(N) shifting.",
        "Resizing copies the entire buffer to a larger block.",
        "Memory pre-allocated even for unused slots.",
        "Concatenation requires reallocation.",
    ]
    add_bullets(s, Inches(0.9), Inches(2.05), Inches(6.0), Inches(3.6),
                items, size=16, line_spacing=1.3)
    add_text(s, Inches(0.9), Inches(5.6), Inches(6.0), Inches(0.5),
             "Linked Lists trade direct indexing for cheap structural change.",
             size=15, italic=True, color=TEAL)
    # right panel: array vs linked picture
    add_rect(s, Inches(7.3), Inches(1.45), Inches(5.4), Inches(5.3),
             LIGHT_GREY)
    add_text(s, Inches(7.3), Inches(1.5), Inches(5.4), Inches(0.4),
             "Conceptual Comparison", size=14, bold=True,
             color=NAVY, align=PP_ALIGN.CENTER)
    # array row
    add_text(s, Inches(7.5), Inches(2.0), Inches(2), Inches(0.3),
             "Array (contiguous)", size=12, bold=True, color=DARK_GREY)
    cell_w = Inches(0.65); x0 = Inches(7.5); y0 = Inches(2.35)
    for i, v in enumerate([10, 20, 30, 40, 50, 60]):
        add_rect(s, x0 + cell_w * i, y0, cell_w, Inches(0.6), WHITE, line=NAVY)
        add_text(s, x0 + cell_w * i, y0, cell_w, Inches(0.6), str(v),
                 size=14, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # linked row
    add_text(s, Inches(7.5), Inches(3.4), Inches(3), Inches(0.3),
             "Linked List (scattered)", size=12, bold=True, color=DARK_GREY)
    nx = Inches(7.5); ny = Inches(3.75); nw = Inches(1.05); nh = Inches(0.6)
    gap = Inches(0.25)
    vals = [10, 20, 30, 40]
    for i, v in enumerate(vals):
        add_node_box(s, nx + (nw + gap) * i, ny, nw, nh, v, value_size=14)
    # arrows between
    for i in range(len(vals) - 1):
        ax1 = nx + (nw + gap) * i + nw
        ax2 = nx + (nw + gap) * (i + 1)
        add_arrow(s, ax1, ny + nh / 2, ax2, ny + nh / 2)
    # null
    last_x = nx + (nw + gap) * (len(vals) - 1) + nw + Inches(0.05)
    add_text(s, last_x, ny, Inches(0.7), nh, "∅",
             size=22, bold=True, color=RED,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # caption
    add_text(s, Inches(7.4), Inches(5.0), Inches(5.2), Inches(1.5),
             "Arrays answer: \"give me element 5\".\n"
             "Linked lists answer: \"insert here, now\".\n"
             "The right tool depends on which question dominates your workload.",
             size=12, italic=True, color=DARK_GREY)


def slide_node_concept(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "2.  The Node",
               "The atom of every linked structure", n, total)
    # left text
    add_text(s, Inches(0.9), Inches(1.45), Inches(6.0), Inches(0.5),
             "Definition", size=22, bold=True, color=NAVY)
    add_text(s, Inches(0.9), Inches(1.95), Inches(6.0), Inches(1.2),
             "A Node is a small object that owns a payload and a reference "
             "(pointer) to one or more other Nodes. The list itself is just "
             "the head node — every other node is reached through pointers.",
             size=15, color=DARK_GREY)
    add_text(s, Inches(0.9), Inches(3.3), Inches(6.0), Inches(0.45),
             "Two fields, one consequence:",
             size=18, bold=True, color=TEAL)
    items = [
        ("data.  ", "The value stored at this position."),
        ("next.  ", "A reference to the following node, or None if last."),
        ("→  ", "Logical order is encoded by the pointers, not by memory layout."),
    ]
    add_bullets(s, Inches(0.9), Inches(3.75), Inches(6.0), Inches(2.8),
                items, size=16, line_spacing=1.35)
    # right diagram
    add_text(s, Inches(7.3), Inches(1.45), Inches(5.4), Inches(0.45),
             "Memory layout (heap)", size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_rect(s, Inches(7.3), Inches(1.95), Inches(5.4), Inches(4.6),
             LIGHT_GREY)
    # node
    nx = Inches(8.0); ny = Inches(3.0); nw = Inches(4.0); nh = Inches(1.4)
    add_node_box(s, nx, ny, nw, nh, "data = 42", value_size=18)
    # next arrow points to nothing for now
    add_text(s, nx + nw + Inches(0.15), ny + nh / 2 - Inches(0.2),
             Inches(1.0), Inches(0.4),
             "→ next",
             size=14, bold=True, color=TEAL,
             align=PP_ALIGN.LEFT)
    add_text(s, Inches(7.3), Inches(5.7), Inches(5.4), Inches(0.7),
             "A node is to a list what a stone is to a path.",
             size=13, italic=True, color=DARK_GREY,
             align=PP_ALIGN.CENTER)


def slide_types(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "3.  Three Variants of Linked Lists",
               "Same atom, different topologies", n, total)
    # three columns
    cols = [
        ("Singly Linked", NAVY,
         "Each node has one pointer: next.",
         ["Lightest representation.",
          "Forward traversal only.",
          "Default choice when you don't need backwards walks."]),
        ("Doubly Linked", TEAL,
         "Each node has two pointers: prev and next.",
         ["Constant-time deletion when you already hold the node.",
          "Used by LRU caches and OS schedulers.",
          "Costs an extra pointer per node."]),
        ("Circular", GOLD,
         "The last node's next loops back to the head.",
         ["Natural for round-robin scheduling.",
          "Music players, buffer queues, turn-taking games.",
          "Beware of infinite traversal — always test termination."]),
    ]
    cw = Inches(4.0); gap = Inches(0.2); x0 = Inches(0.6); y0 = Inches(1.55)
    for i, (title, color, desc, bullets) in enumerate(cols):
        cx = x0 + (cw + gap) * i
        add_rect(s, cx, y0, cw, Inches(0.7), color)
        add_text(s, cx, y0, cw, Inches(0.7), title,
                 size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, cx, y0 + Inches(0.7), cw, Inches(4.6), LIGHT_GREY)
        add_text(s, cx + Inches(0.2), y0 + Inches(0.85),
                 cw - Inches(0.4), Inches(0.9),
                 desc, size=14, italic=True, color=DARK_GREY)
        add_bullets(s, cx + Inches(0.2), y0 + Inches(2.0),
                    cw - Inches(0.4), Inches(3.2),
                    bullets, size=13, line_spacing=1.25,
                    bullet_color=color)
    # mini diagrams under each column
    diag_y = Inches(6.4)
    # singly: A -> B -> C -> ∅
    sx = Inches(0.7); nw = Inches(0.85); nh = Inches(0.45); gp = Inches(0.18)
    for i, v in enumerate(["A", "B", "C"]):
        add_node_box(s, sx + (nw + gp) * i, diag_y, nw, nh, v, value_size=12)
    for i in range(2):
        x1 = sx + (nw + gp) * i + nw
        x2 = sx + (nw + gp) * (i + 1)
        add_arrow(s, x1, diag_y + nh / 2, x2, diag_y + nh / 2)
    add_text(s, sx + (nw + gp) * 3, diag_y, Inches(0.5), nh, "∅",
             size=16, bold=True, color=RED,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # doubly: ∅ <-> A <-> B <-> C <-> ∅
    dx = Inches(4.95)
    add_text(s, dx - Inches(0.4), diag_y, Inches(0.4), nh, "∅",
             size=14, bold=True, color=RED,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for i, v in enumerate(["A", "B", "C"]):
        add_node_box(s, dx + (nw + gp) * i, diag_y, nw, nh, v,
                     value_size=12, has_prev=True)
    for i in range(2):
        x1 = dx + (nw + gp) * i + nw
        x2 = dx + (nw + gp) * (i + 1)
        add_arrow(s, x1, diag_y + nh / 2 - Emu(40000),
                  x2, diag_y + nh / 2 - Emu(40000))
        add_arrow(s, x2, diag_y + nh / 2 + Emu(40000),
                  x1, diag_y + nh / 2 + Emu(40000), color=GOLD)
    # circular: A -> B -> C ↺
    cx0 = Inches(9.2)
    for i, v in enumerate(["A", "B", "C"]):
        add_node_box(s, cx0 + (nw + gp) * i, diag_y, nw, nh, v, value_size=12)
    for i in range(2):
        x1 = cx0 + (nw + gp) * i + nw
        x2 = cx0 + (nw + gp) * (i + 1)
        add_arrow(s, x1, diag_y + nh / 2, x2, diag_y + nh / 2)
    # loop arrow back
    last_right = cx0 + (nw + gp) * 2 + nw
    add_arrow(s, last_right, diag_y + nh,
              cx0, diag_y + nh + Inches(0.25), color=GOLD)
    add_arrow(s, cx0, diag_y + nh + Inches(0.25),
              cx0, diag_y + nh / 2, color=GOLD)


def slide_complexity(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "4.  Complexity at a Glance",
               "When does this structure pay off?", n, total)
    # table
    rows = [
        ("Operation",                "Array",    "Singly LL", "Doubly LL"),
        ("Random access list[i]",    "O(1)",     "O(N)",      "O(N)"),
        ("Insert at head",           "O(N)",     "O(1)",      "O(1)"),
        ("Insert at tail (no ptr)",  "O(1)*",    "O(N)",      "O(N)"),
        ("Insert at tail (w/ ptr)",  "O(1)",     "O(1)",      "O(1)"),
        ("Delete given node ptr",    "O(N)",     "O(N)",      "O(1)"),
        ("Search by value",          "O(N)",     "O(N)",      "O(N)"),
        ("Memory overhead",          "Low",      "1 ptr/node","2 ptrs/node"),
    ]
    x0 = Inches(0.9); y0 = Inches(1.55)
    col_w = [Inches(4.0), Inches(2.5), Inches(2.7), Inches(2.7)]
    rh = Inches(0.55)
    for r, row in enumerate(rows):
        cx = x0
        for c, txt in enumerate(row):
            fill = NAVY if r == 0 else (LIGHT_GREY if r % 2 == 1 else WHITE)
            color = WHITE if r == 0 else DARK_GREY
            bold = (r == 0) or (c == 0)
            add_rect(s, cx, y0 + rh * r, col_w[c], rh, fill,
                     line=RGBColor(0xC0, 0xC4, 0xCC))
            add_text(s, cx + Inches(0.15), y0 + rh * r,
                     col_w[c] - Inches(0.2), rh,
                     txt, size=14, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[c]
    add_text(s, Inches(0.9), Inches(6.2), Inches(11), Inches(0.4),
             "*Amortized for dynamic arrays — individual resizes are O(N).",
             size=12, italic=True, color=DARK_GREY)
    add_text(s, Inches(0.9), Inches(6.6), Inches(11), Inches(0.4),
             "Rule of thumb: choose linked structures when modifications dominate reads.",
             size=14, bold=True, color=TEAL)


def slide_operations_overview(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "5.  Core Operations — Mental Models",
               "Every operation reduces to: which pointers do I rewrite, and in what order?",
               n, total)
    # 4 quadrants
    quads = [
        ("Append (insert at tail)", NAVY,
         ["Walk to last node where next is None.",
          "Set last.next = new_node.",
          "O(N) without a tail pointer; O(1) with one."]),
        ("Prepend (insert at head)", TEAL,
         ["new_node.next = head.",
          "head = new_node.",
          "Always O(1) — no traversal needed."]),
        ("Insert at position k", GOLD,
         ["Walk to node k-1 (the predecessor).",
          "new_node.next = predecessor.next.",
          "predecessor.next = new_node."]),
        ("Delete node with value v", RED,
         ["Track previous pointer while traversing.",
          "previous.next = current.next.",
          "Mind the head: deleting head shifts the list."]),
    ]
    qw = Inches(5.9); qh = Inches(2.45); gx = Inches(0.3); gy = Inches(0.25)
    x0 = Inches(0.6); y0 = Inches(1.6)
    for i, (title, color, bullets) in enumerate(quads):
        col = i % 2; row = i // 2
        cx = x0 + (qw + gx) * col
        cy = y0 + (qh + gy) * row
        add_rect(s, cx, cy, qw, Inches(0.55), color)
        add_text(s, cx, cy, qw, Inches(0.55), title,
                 size=17, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, cx, cy + Inches(0.55), qw, qh - Inches(0.55), LIGHT_GREY)
        add_bullets(s, cx + Inches(0.25), cy + Inches(0.7),
                    qw - Inches(0.5), qh - Inches(0.7),
                    bullets, size=14, line_spacing=1.3,
                    bullet_color=color)


def slide_insert_walkthrough(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "6.  Insertion: A Visual Walkthrough",
               "Inserting node X between B and C", n, total)
    # before
    add_text(s, Inches(0.9), Inches(1.5), Inches(11), Inches(0.4),
             "Step 1 — Initial list",
             size=16, bold=True, color=NAVY)
    nw = Inches(1.1); nh = Inches(0.7); gp = Inches(0.45)
    sx = Inches(1.0); sy = Inches(2.0)
    for i, v in enumerate(["A", "B", "C", "D"]):
        add_node_box(s, sx + (nw + gp) * i, sy, nw, nh, v)
    for i in range(3):
        x1 = sx + (nw + gp) * i + nw
        x2 = sx + (nw + gp) * (i + 1)
        add_arrow(s, x1, sy + nh / 2, x2, sy + nh / 2)
    # step 2
    add_text(s, Inches(0.9), Inches(3.1), Inches(11), Inches(0.4),
             "Step 2 — Allocate X. Set X.next = B.next  (i.e. point X to C).",
             size=16, bold=True, color=NAVY)
    sy2 = Inches(3.65)
    for i, v in enumerate(["A", "B", "C", "D"]):
        add_node_box(s, sx + (nw + gp) * i, sy2, nw, nh, v)
    for i in range(3):
        x1 = sx + (nw + gp) * i + nw
        x2 = sx + (nw + gp) * (i + 1)
        add_arrow(s, x1, sy2 + nh / 2, x2, sy2 + nh / 2)
    # X above
    xnx = sx + (nw + gp) * 1 + Inches(0.5); xny = sy2 + Inches(1.05)
    add_node_box(s, xnx, xny, nw, nh, "X", fill=GOLD)
    # arrow X.next -> C
    cx = sx + (nw + gp) * 2 + nw / 2
    add_arrow(s, xnx + nw / 2, xny, cx, sy2 + nh, color=GOLD)
    # step 3
    add_text(s, Inches(0.9), Inches(5.55), Inches(11), Inches(0.4),
             "Step 3 — Set B.next = X. Insertion done in O(1) once predecessor is known.",
             size=16, bold=True, color=NAVY)
    sy3 = Inches(6.1)
    seq = ["A", "B", "X", "C", "D"]
    nw2 = Inches(0.95); gp2 = Inches(0.3)
    sx3 = Inches(1.0)
    for i, v in enumerate(seq):
        fill = GOLD if v == "X" else WHITE
        add_node_box(s, sx3 + (nw2 + gp2) * i, sy3, nw2, Inches(0.6), v,
                     fill=fill, value_size=16)
    for i in range(len(seq) - 1):
        x1 = sx3 + (nw2 + gp2) * i + nw2
        x2 = sx3 + (nw2 + gp2) * (i + 1)
        add_arrow(s, x1, sy3 + Inches(0.3), x2, sy3 + Inches(0.3))


def slide_node_code(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "7.  Implementation — The Node",
               "Minimal, immutable in spirit, defensible in production",
               n, total)
    code = (
        "class Node:\n"
        "    \"\"\"A single element of a singly linked list.\"\"\"\n"
        "\n"
        "    __slots__ = (\"data\", \"next\")\n"
        "\n"
        "    def __init__(self, data):\n"
        "        self.data = data\n"
        "        self.next = None\n"
        "\n"
        "    def __repr__(self):\n"
        "        return f\"Node({self.data!r})\""
    )
    add_code_block(s, Inches(0.9), Inches(1.5), Inches(7.2), Inches(4.5),
                   code, size=16)
    # right column: notes
    add_text(s, Inches(8.4), Inches(1.5), Inches(4.4), Inches(0.5),
             "Design notes", size=18, bold=True, color=NAVY)
    notes = [
        "__slots__ shrinks per-node memory by skipping __dict__.",
        "next is initialised to None — no dangling pointers.",
        "data is the abstract payload: int, str, object, anything.",
        "__repr__ aids debugging and pointer-tracing in tests.",
    ]
    add_bullets(s, Inches(8.4), Inches(2.0), Inches(4.4), Inches(4.0),
                notes, size=14, line_spacing=1.35)


def slide_list_code(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "8.  Implementation — The Singly Linked List",
               "Append in O(N); print_list in O(N); both stream-friendly",
               n, total)
    code = (
        "class LinkedList:\n"
        "    def __init__(self):\n"
        "        self.head = None\n"
        "\n"
        "    def append(self, data):\n"
        "        new_node = Node(data)\n"
        "        if self.head is None:        # empty list\n"
        "            self.head = new_node\n"
        "            return\n"
        "        current = self.head          # walk to last node\n"
        "        while current.next:\n"
        "            current = current.next\n"
        "        current.next = new_node      # link in place\n"
        "\n"
        "    def print_list(self):\n"
        "        elements, current = [], self.head\n"
        "        while current:\n"
        "            elements.append(str(current.data))\n"
        "            current = current.next\n"
        "        print(\" -> \".join(elements))"
    )
    add_code_block(s, Inches(0.9), Inches(1.45), Inches(8.2), Inches(5.4),
                   code, size=14)
    # right notes
    add_text(s, Inches(9.4), Inches(1.45), Inches(3.5), Inches(0.5),
             "Why this shape?", size=18, bold=True, color=NAVY)
    notes = [
        "Empty-list branch protects against AttributeError on None.",
        "Loop variable current never escapes the method — no aliasing risk.",
        "Tail pointer would reduce append to O(1); deferred to keep the lab focused.",
        "print_list builds the string once — O(N) joins are cheaper than repeated I/O.",
    ]
    add_bullets(s, Inches(9.4), Inches(1.95), Inches(3.6), Inches(4.6),
                notes, size=13, line_spacing=1.3)


def slide_reverse_concept(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "9.  Reversing a Linked List — Intuition",
               "The interview classic. Three pointers, no extra memory.",
               n, total)
    add_text(s, Inches(0.9), Inches(1.45), Inches(11), Inches(0.5),
             "We walk the list once, redirecting each node's next to its predecessor.",
             size=17, color=DARK_GREY)
    # before
    add_text(s, Inches(0.9), Inches(2.05), Inches(5), Inches(0.4),
             "Before",
             size=15, bold=True, color=TEAL)
    nw = Inches(0.95); nh = Inches(0.6); gp = Inches(0.3)
    sx = Inches(0.9); sy = Inches(2.55)
    for i, v in enumerate(["1", "2", "3", "4"]):
        add_node_box(s, sx + (nw + gp) * i, sy, nw, nh, v, value_size=16)
    for i in range(3):
        x1 = sx + (nw + gp) * i + nw
        x2 = sx + (nw + gp) * (i + 1)
        add_arrow(s, x1, sy + nh / 2, x2, sy + nh / 2)
    add_text(s, sx + (nw + gp) * 4, sy, Inches(0.5), nh, "∅",
             size=18, bold=True, color=RED,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # after
    add_text(s, Inches(0.9), Inches(3.6), Inches(5), Inches(0.4),
             "After",
             size=15, bold=True, color=TEAL)
    sy2 = Inches(4.1)
    for i, v in enumerate(["4", "3", "2", "1"]):
        add_node_box(s, sx + (nw + gp) * i, sy2, nw, nh, v, value_size=16)
    for i in range(3):
        x1 = sx + (nw + gp) * i + nw
        x2 = sx + (nw + gp) * (i + 1)
        add_arrow(s, x1, sy2 + nh / 2, x2, sy2 + nh / 2, color=GOLD)
    add_text(s, sx + (nw + gp) * 4, sy2, Inches(0.5), nh, "∅",
             size=18, bold=True, color=RED,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # invariant
    add_rect(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.6),
             LIGHT_GREY)
    add_text(s, Inches(1.1), Inches(5.3), Inches(11), Inches(0.5),
             "Loop invariant",
             size=16, bold=True, color=NAVY)
    add_text(s, Inches(1.1), Inches(5.7), Inches(11), Inches(1.1),
             "At every step, the prefix of the list up to and including prev "
             "has been reversed; current points to the head of the unprocessed "
             "tail. We advance until current is None — at which point prev is "
             "the new head.",
             size=14, color=DARK_GREY)


def slide_reverse_code(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Reversing a Linked List — Code",
               "Three named pointers: prev, current, next_node",
               n, total)
    code = (
        "def reverse_linked_list(head):\n"
        "    prev = None\n"
        "    current = head\n"
        "\n"
        "    while current:\n"
        "        next_node   = current.next   # remember the tail\n"
        "        current.next = prev          # flip the pointer\n"
        "        prev         = current       # advance prev\n"
        "        current      = next_node     # advance current\n"
        "\n"
        "    return prev                      # new head of the list"
    )
    add_code_block(s, Inches(0.9), Inches(1.5), Inches(7.5), Inches(4.0),
                   code, size=16)
    # right side: trace
    add_text(s, Inches(8.7), Inches(1.5), Inches(4.3), Inches(0.5),
             "Trace on  1 → 2 → 3", size=18, bold=True, color=NAVY)
    trace = [
        "Start    prev=∅   curr=1",
        "Iter 1   prev=1   curr=2",
        "Iter 2   prev=2→1 curr=3",
        "Iter 3   prev=3→2→1 curr=∅",
        "Return   3 → 2 → 1",
    ]
    add_bullets(s, Inches(8.7), Inches(2.0), Inches(4.3), Inches(3.5),
                trace, size=14, line_spacing=1.4, bullet_color=GOLD,
                font="Consolas")
    # complexity
    add_rect(s, Inches(0.9), Inches(5.8), Inches(11.5), Inches(1.0),
             LIGHT_GREY)
    add_text(s, Inches(1.1), Inches(5.85), Inches(11), Inches(0.45),
             "Complexity",
             size=16, bold=True, color=NAVY)
    add_text(s, Inches(1.1), Inches(6.25), Inches(11), Inches(0.5),
             "Time  O(N)  ·  visit each node once.   "
             "Space  O(1)  ·  three named pointers regardless of input size.",
             size=14, color=DARK_GREY)


def slide_cycle_concept(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "10.  Cycle Detection — The Question",
               "Does the chain end at None, or does it eat its own tail?",
               n, total)
    # diagram cycle
    add_text(s, Inches(0.9), Inches(1.45), Inches(11), Inches(0.5),
             "A buggy linked list whose final pointer wraps back into itself:",
             size=17, color=DARK_GREY)
    nw = Inches(0.95); nh = Inches(0.7); gp = Inches(0.4)
    sx = Inches(1.5); sy = Inches(2.5)
    for i, v in enumerate(["5", "6", "7"]):
        add_node_box(s, sx + (nw + gp) * i, sy, nw, nh, v, value_size=18)
    # 5->6, 6->7
    for i in range(2):
        x1 = sx + (nw + gp) * i + nw
        x2 = sx + (nw + gp) * (i + 1)
        add_arrow(s, x1, sy + nh / 2, x2, sy + nh / 2)
    # 7 -> 6 (loop)
    x_last_right = sx + (nw + gp) * 2 + nw
    x_b_top = sx + (nw + gp) * 1 + nw / 2
    add_arrow(s, x_last_right, sy + nh,
              x_last_right + Inches(0.3), sy + nh + Inches(0.6),
              color=RED)
    add_arrow(s, x_last_right + Inches(0.3), sy + nh + Inches(0.6),
              x_b_top, sy + nh + Inches(0.6), color=RED)
    add_arrow(s, x_b_top, sy + nh + Inches(0.6),
              x_b_top, sy + nh, color=RED)
    add_text(s, x_last_right + Inches(0.4), sy + nh + Inches(0.7),
             Inches(2.5), Inches(0.4),
             "cycle!", size=14, bold=True, color=RED, italic=True)
    # right side
    add_text(s, Inches(7.5), Inches(1.6), Inches(5), Inches(0.5),
             "Naive solutions", size=18, bold=True, color=NAVY)
    naive = [
        "Hash-set of seen nodes — O(N) time but O(N) extra memory.",
        "Mutate visited flag — destructive; pollutes user data.",
        "Limit traversal to N steps — fragile; you may not know N.",
    ]
    add_bullets(s, Inches(7.5), Inches(2.1), Inches(5.5), Inches(2.5),
                naive, size=14, line_spacing=1.3, bullet_color=RED)
    add_text(s, Inches(7.5), Inches(4.7), Inches(5.5), Inches(0.5),
             "We want O(N) time and O(1) memory — without mutation.",
             size=15, italic=True, color=TEAL)
    # punchline
    add_rect(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.2),
             LIGHT_GREY)
    add_text(s, Inches(1.1), Inches(5.8), Inches(11), Inches(0.5),
             "Floyd's insight",
             size=16, bold=True, color=NAVY)
    add_text(s, Inches(1.1), Inches(6.2), Inches(11), Inches(0.7),
             "If two runners start together on a circular track, the faster "
             "one will eventually lap the slower one — and they will meet at "
             "the same node.",
             size=14, italic=True, color=DARK_GREY)


def slide_floyd_code(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "11.  Floyd's Tortoise & Hare",
               "Two pointers, two speeds, one elegant proof",
               n, total)
    code = (
        "def has_cycle(head):\n"
        "    if head is None or head.next is None:\n"
        "        return False\n"
        "\n"
        "    slow = head           # tortoise: 1 step / iteration\n"
        "    fast = head.next      # hare:     2 steps / iteration\n"
        "\n"
        "    while slow != fast:\n"
        "        if fast is None or fast.next is None:\n"
        "            return False         # hare reached the end\n"
        "        slow = slow.next\n"
        "        fast = fast.next.next\n"
        "\n"
        "    return True                  # collision => cycle"
    )
    add_code_block(s, Inches(0.9), Inches(1.5), Inches(7.5), Inches(4.7),
                   code, size=14)
    # right notes
    add_text(s, Inches(8.7), Inches(1.5), Inches(4.3), Inches(0.5),
             "Why it works", size=18, bold=True, color=NAVY)
    notes = [
        "If no cycle exists, the hare reaches None and we stop.",
        "If a cycle exists, both pointers eventually enter the loop.",
        "Inside the loop, the gap between hare and tortoise shrinks by 1 each step.",
        "Therefore in at most L iterations (L = loop length), they collide.",
    ]
    add_bullets(s, Inches(8.7), Inches(2.0), Inches(4.3), Inches(4.5),
                notes, size=13, line_spacing=1.3)
    add_rect(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.7),
             LIGHT_GREY)
    add_text(s, Inches(1.1), Inches(6.45), Inches(11), Inches(0.5),
             "Time  O(N)   ·   Space  O(1)   ·   Non-mutating.",
             size=15, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)


def slide_pitfalls(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "12.  Common Pitfalls",
               "Where students lose marks — and where production code crashes",
               n, total)
    items = [
        ("Losing the head.  ",
         "Mutating self.head without saving a reference; the rest of the list becomes unreachable."),
        ("Off-by-one in traversal.  ",
         "while current.next vs while current — they differ by exactly one node."),
        ("Forgetting the empty list.  ",
         "Calling current.next on None raises AttributeError; always guard self.head is None."),
        ("Pointer-overwrite order.  ",
         "In reversal, save next_node before overwriting current.next, never after."),
        ("Cycle-induced infinite loops.  ",
         "Plain traversals over a circular list never terminate. Validate with Floyd before walking."),
        ("Memory leaks in non-GC languages.  ",
         "Free deleted nodes in C/C++; Python's refcounting handles it for you."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(12), Inches(5.5),
                items, size=16, line_spacing=1.35, bullet_color=RED)


def slide_applications(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "13.  Where Linked Lists Live in the Real World",
               "Not academic curios — these are running on your laptop right now",
               n, total)
    cards = [
        ("Operating Systems",  NAVY,
         "Process scheduler queues, free-memory lists, file system inode chains."),
        ("Browsers",           TEAL,
         "Back/forward navigation history is a doubly linked list."),
        ("Music & Media",      GOLD,
         "Playlists and round-robin queues — circular linked lists."),
        ("Editors & IDEs",     RED,
         "Undo/redo stacks; piece tables for text buffers (Word, VS Code)."),
        ("Caches",             GREEN,
         "LRU caches combine a hash map with a doubly linked list — O(1) eviction."),
        ("Polynomials & Math", NAVY,
         "Sparse polynomial representations; symbolic computation engines."),
    ]
    cw = Inches(4.0); ch = Inches(2.4); gx = Inches(0.2); gy = Inches(0.2)
    x0 = Inches(0.6); y0 = Inches(1.55)
    for i, (title, color, body) in enumerate(cards):
        col = i % 3; row = i // 3
        cx = x0 + (cw + gx) * col
        cy = y0 + (ch + gy) * row
        add_rect(s, cx, cy, cw, Inches(0.55), color)
        add_text(s, cx, cy, cw, Inches(0.55), title,
                 size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, cx, cy + Inches(0.55), cw, ch - Inches(0.55), LIGHT_GREY)
        add_text(s, cx + Inches(0.2), cy + Inches(0.7),
                 cw - Inches(0.4), ch - Inches(0.7),
                 body, size=14, color=DARK_GREY)


def slide_lab(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "14.  Practical Lab — practice.py",
               "Three exercises. Build, manipulate, detect.",
               n, total)
    items = [
        ("Exercise 1.  ",
         "Implement Node and a Singly LinkedList with append() and print_list()."),
        ("Exercise 2.  ",
         "Implement reverse_linked_list(head) using only three pointers; in-place."),
        ("Exercise 3.  ",
         "Implement has_cycle(head) using Floyd's algorithm in O(1) extra memory."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.55), Inches(12), Inches(2.5),
                items, size=18, line_spacing=1.4)
    add_text(s, Inches(0.9), Inches(4.1), Inches(12), Inches(0.5),
             "Submission expectations",
             size=18, bold=True, color=NAVY)
    rules = [
        "Provide test cases for empty list, single node, and a cycle.",
        "Use meaningful names: prev, current, next_node — never x, y, z.",
        "Include complexity comments above each function (Time / Space).",
        "Do not import any standard library structure — implement the chain yourself.",
    ]
    add_bullets(s, Inches(0.9), Inches(4.55), Inches(12), Inches(2.5),
                rules, size=15, line_spacing=1.3, bullet_color=GOLD)


def slide_summary(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "15.  Summary",
               "Five sentences you should be able to recite at the door",
               n, total)
    items = [
        "A linked list is a chain of self-referential nodes; logical order is encoded in pointers, not memory addresses.",
        "Linked lists trade O(1) random access for O(1) structural modification.",
        "Singly, doubly, and circular variants share the same atom and differ only in their pointer topology.",
        "Reversal redirects each node's next using three named pointers — prev, current, next_node — in O(N) time and O(1) space.",
        "Floyd's Tortoise and Hare detects cycles in O(N) time using only two pointers, no auxiliary memory, and no mutation.",
    ]
    add_bullets(s, Inches(0.9), Inches(1.55), Inches(12), Inches(5.5),
                items, size=18, line_spacing=1.4, bullet_color=GOLD)


def slide_references(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "16.  Further Reading",
               "Trusted sources for deeper study", n, total)
    refs = [
        ("CLRS — Introduction to Algorithms (4th ed.).  ",
         "Cormen, Leiserson, Rivest & Stein. MIT Press. Chapter on Elementary Data Structures."),
        ("The Art of Computer Programming, Vol. 1.  ",
         "Donald E. Knuth. Section 2.2 — Linear Lists."),
        ("Algorithms (4th ed.).  ",
         "Sedgewick & Wayne. Princeton University. Chapter 1.3 — Bags, Queues, Stacks."),
        ("Floyd, R. W. (1967).  ",
         "Nondeterministic Algorithms. JACM 14(4). Origin of cycle-finding."),
        ("LeetCode practice set.  ",
         "Problems 206 (Reverse), 141 (Cycle), 142 (Cycle II), 21 (Merge), 19 (Remove Nth)."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.55), Inches(12), Inches(5.5),
                refs, size=16, line_spacing=1.4)


def slide_thanks():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(5.5), SLIDE_W, Inches(0.12), GOLD)
    add_text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.5),
             "Questions?", size=80, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.7),
             "Bring your code, your traces, and your edge cases.",
             size=22, italic=True, color=GOLD)
    add_text(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.4),
             "Next week  ·  Stacks & Queues — bounded chains with discipline.",
             size=16, color=LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4),
             "Data Structures & Algorithms  ·  Week 03  ·  Linked Lists",
             size=12, color=LIGHT_GREY, italic=True)


# -------- Build ----------
# We need page numbers; build a list and run.
TOTAL = 18  # title + 16 content + thanks

slide_title()
# placeholder list of (callable, needs_page) -- we know order
content = [
    slide_outline,
    slide_objectives,
    slide_motivation,
    slide_node_concept,
    slide_types,
    slide_complexity,
    slide_operations_overview,
    slide_insert_walkthrough,
    slide_node_code,
    slide_list_code,
    slide_reverse_concept,
    slide_reverse_code,
    slide_cycle_concept,
    slide_floyd_code,
    slide_pitfalls,
    slide_applications,
    slide_lab,
    slide_summary,
    slide_references,
]
TOTAL = 1 + len(content) + 1  # title + content + thanks
for i, fn in enumerate(content, start=2):
    fn(i, TOTAL)
slide_thanks()

out = "/Users/silasgah/Downloads/enquire.ai/ds/week_03_linked_lists/linked_lists_dsa_final.pptx"
prs.save(out)
print(f"Wrote {out}  ({len(prs.slides)} slides)")
