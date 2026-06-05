"""
Floyd's Tortoise & Hare — Deep-Dive Lecture Deck
=================================================
Generates a university-grade, stand-alone PowerPoint presentation that
explains Floyd's Cycle-Finding Algorithm from first principles:

  • Motivation & naive approaches
  • Intuition (two runners on a track)
  • Step-by-step animated diagrams
  • Formal correctness proof
  • Full annotated Python implementation
  • Extension: finding the cycle entry-point (Floyd Part II)
  • Complexity analysis
  • Real-world applications
  • Practice problems

Output : floyd_tortoise_hare.pptx  (16:9 widescreen)
Author : DSA Curriculum  –  Week 03 supplement
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Colour palette (matches the main Week 03 deck) ──────────────────────────
NAVY       = RGBColor(0x0B, 0x2A, 0x4A)
TEAL       = RGBColor(0x12, 0x7A, 0x8A)
GOLD       = RGBColor(0xE0, 0xA8, 0x1F)
LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF7)
DARK_GREY  = RGBColor(0x33, 0x3D, 0x4A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
RED        = RGBColor(0xC0, 0x39, 0x2B)
GREEN      = RGBColor(0x2E, 0x7D, 0x32)
PURPLE     = RGBColor(0x6A, 0x1B, 0x9A)
ORANGE     = RGBColor(0xE6, 0x5C, 0x00)
CODE_BG    = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG    = RGBColor(0xE6, 0xE6, 0xE6)
CODE_KW    = RGBColor(0xBB, 0x86, 0xFC)   # keyword purple
CODE_CM    = RGBColor(0x6A, 0x99, 0x55)   # comment green

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── Low-level helpers ────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, line=None, line_width=1):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    return shp


def add_rrect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(2)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top   = Inches(0.02)
    tf.margin_bottom= Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=18, color=DARK_GREY,
                bullet_color=TEAL, line_spacing=1.3, font="Calibri",
                indent=0):
    """
    items: list of strings  OR  (bold_label, rest_text) tuples.
    indent: 0 = top-level bullet (▸), 1 = sub-bullet (–)
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after  = Pt(4)
        glyph = "  –  " if indent else "▸  "
        rb = p.add_run()
        rb.text = glyph
        rb.font.name  = font
        rb.font.size  = Pt(size)
        rb.font.bold  = not indent
        rb.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run()
            r1.text = head
            r1.font.name  = font
            r1.font.size  = Pt(size)
            r1.font.bold  = True
            r1.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = tail
            r2.font.name  = font
            r2.font.size  = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name  = font
            r.font.size  = Pt(size)
            r.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None, page_num=None, total=None):
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    add_rect(slide, Inches(0.35), 0, SLIDE_W - Inches(0.35), Inches(0.08), GOLD)
    add_text(slide, Inches(0.6), Inches(0.18), Inches(11.8), Inches(0.72),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.82), Inches(11.8), Inches(0.42),
                 subtitle, size=15, italic=True, color=TEAL)
    add_rect(slide, Inches(0.6), Inches(1.26), Inches(2.0), Inches(0.04), GOLD)
    add_text(slide, Inches(0.6), SLIDE_H - Inches(0.45), Inches(9), Inches(0.3),
             "Floyd's Tortoise & Hare  ·  Cycle Detection  ·  Week 03 Supplement",
             size=10, color=DARK_GREY, italic=True)
    if page_num is not None and total is not None:
        add_text(slide, SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.45),
                 Inches(1.2), Inches(0.3),
                 f"Slide {page_num} / {total}", size=10, color=DARK_GREY,
                 align=PP_ALIGN.RIGHT)


def add_code_block(slide, x, y, w, h, code, size=13):
    add_rect(slide, x, y, w, h, CODE_BG)
    add_rect(slide, x, y, Inches(0.08), h, GOLD)
    tb = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.1),
                                  w - Inches(0.28), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    lines = code.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.18
        r = p.add_run()
        r.text = ln if ln else " "
        r.font.name  = "Consolas"
        r.font.size  = Pt(size)
        r.font.color.rgb = CODE_FG


def add_arrow(slide, x1, y1, x2, y2, color=TEAL, weight=2.25, head=True):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if head:
        ln = conn.line._get_or_add_ln()
        etree.SubElement(ln, qn("a:tailEnd"),
                         {"type": "triangle", "w": "med", "len": "med"})
    return conn


def add_node(slide, x, y, w, h, label, *,
             fill=WHITE, border=NAVY, label_size=18,
             sub=None, sub_color=TEAL):
    """Draw a simple node box (data only, no next cell)."""
    shp = add_rrect(slide, x, y, w, h, fill, line=border)
    add_text(slide, x, y, w, h, str(label),
             size=label_size, bold=True, color=NAVY if fill == WHITE else WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        add_text(slide, x, y + h, w, Inches(0.3), sub,
                 size=10, bold=True, color=sub_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    return shp


def add_pointer_label(slide, x, y, label, color=TEAL, size=12):
    add_text(slide, x - Inches(0.25), y, Inches(0.7), Inches(0.35),
             label, size=size, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def callout_box(slide, x, y, w, h, title, body, *, accent=TEAL):
    add_rect(slide, x, y, Inches(0.06), h, accent)
    add_rect(slide, x + Inches(0.06), y, w - Inches(0.06), h, LIGHT_GREY)
    add_text(slide, x + Inches(0.2), y + Inches(0.05), w - Inches(0.25),
             Inches(0.4), title, size=14, bold=True, color=accent)
    add_text(slide, x + Inches(0.2), y + Inches(0.42), w - Inches(0.25),
             h - Inches(0.5), body, size=13, color=DARK_GREY)


# ── Slide builders ───────────────────────────────────────────────────────────

def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # accent stripes
    add_rect(s, 0, Inches(5.6), SLIDE_W, Inches(0.14), GOLD)
    add_rect(s, 0, Inches(5.8), SLIDE_W, Inches(0.05), TEAL)
    # week chip
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.9), Inches(0.9),
                               Inches(3.2), Inches(0.58))
    chip.fill.solid(); chip.fill.fore_color.rgb = TEAL
    chip.line.fill.background(); chip.shadow.inherit = False
    add_text(s, Inches(0.9), Inches(0.9), Inches(3.2), Inches(0.58),
             "WEEK 03  ·  DEEP DIVE", size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # main title
    add_text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.6),
             "Floyd's Tortoise\n& Hare", size=60, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.35), Inches(11.5), Inches(0.65),
             "Cycle Detection in O(N) Time and O(1) Space",
             size=24, color=GOLD, italic=True)
    add_text(s, Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.55),
             "A deep-dive into Floyd's Cycle-Finding Algorithm — from intuition to proof",
             size=17, color=LIGHT_GREY)
    # bottom metadata
    add_text(s, Inches(0.9), Inches(6.1), Inches(8), Inches(0.4),
             "Algorithm   ·   Robert W. Floyd  (1967)",
             size=13, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.5), Inches(8), Inches(0.4),
             "Module      ·   Data Structures & Algorithms  —  Week 03  Supplement",
             size=13, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.9), Inches(8), Inches(0.4),
             "Complexity  ·   Time O(N)  ·  Space O(1)  ·  Non-mutating",
             size=13, color=WHITE)
    # tortoise & hare emoji-style icons (text stand-ins)
    add_text(s, Inches(10.5), Inches(1.8), Inches(2.5), Inches(1.2),
             "🐢", size=72, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(10.5), Inches(2.9), Inches(2.5), Inches(1.2),
             "🐇", size=72, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_outline(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Lecture Outline",
               "What this deep-dive covers and why every part matters",
               n, total)
    items = [
        ("1. Motivation.  ",    "Why cycle detection is harder than it looks."),
        ("2. Naive Approaches.  ", "Hash sets, visited flags, step limits — and their costs."),
        ("3. The Intuition.  ", "Two runners on a circular track; relative speed."),
        ("4. How It Works.  ", "Step-by-step trace through a cyclic list."),
        ("5. Correctness Proof.  ", "Mathematical argument — why they must meet."),
        ("6. The Code.  ",     "Annotated Python implementation from scratch."),
        ("7. Finding the Entry Point.  ", "Floyd Part II: locating where the cycle begins."),
        ("8. Complexity.  ",  "Tight O(N) time, O(1) space — with proof sketch."),
        ("9. Variants & Extensions.  ", "Happy Number, middle of list, duplicate detection."),
        ("10. Applications.  ", "Real systems that rely on this exact technique."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(12.1), Inches(5.6),
                items, size=17, line_spacing=1.35)


def slide_objectives(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Learning Outcomes",
               "By the end of this session you will be able to…", n, total)
    objs = [
        "State precisely what Floyd's algorithm detects and why it is preferable to hash-set approaches.",
        "Trace the slow and fast pointers step-by-step through any cyclic or acyclic linked list.",
        "Explain, in plain English, why two pointers at different speeds must meet inside a cycle.",
        "Reproduce the mathematical proof that the meeting is guaranteed within O(N) steps.",
        "Implement has_cycle() and find_cycle_entry() from memory, with correct edge-case handling.",
        "Identify real-world problems (Happy Number, tortoise in physics, rho-algorithm) that reduce to Floyd's.",
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(12.1), Inches(5.6),
                objs, size=17, line_spacing=1.35)


def slide_motivation(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "1.  Motivation — Why Is This Hard?",
               "Cycle detection sounds simple. It isn't.", n, total)

    add_text(s, Inches(0.9), Inches(1.45), Inches(12), Inches(0.5),
             "Consider a linked list whose tail node loops back to an earlier node:",
             size=16, color=DARK_GREY)

    # Draw a cyclic list  5 → 6 → 7 → 8 → 6 (cycle)
    nw = Inches(0.9); nh = Inches(0.65); gp = Inches(0.42)
    sx = Inches(1.2); sy = Inches(2.2)
    vals = ["5", "6", "7", "8"]
    for i, v in enumerate(vals):
        fill = LIGHT_GREY if v == "6" else WHITE
        add_node(s, sx + (nw + gp) * i, sy, nw, nh, v,
                 fill=fill, border=NAVY, label_size=18)
    for i in range(3):
        x1 = sx + (nw + gp) * i + nw
        x2 = sx + (nw + gp) * (i + 1)
        add_arrow(s, x1, sy + nh / 2, x2, sy + nh / 2)
    # loop: 8 → 6
    x_8_right = sx + (nw + gp) * 3 + nw
    x_6_mid   = sx + (nw + gp) * 1 + nw / 2
    add_arrow(s, x_8_right, sy + nh / 2,
              x_8_right + Inches(0.35), sy + nh / 2, color=RED)
    add_arrow(s, x_8_right + Inches(0.35), sy + nh / 2,
              x_8_right + Inches(0.35), sy - Inches(0.5), color=RED, head=False)
    add_arrow(s, x_8_right + Inches(0.35), sy - Inches(0.5),
              x_6_mid, sy - Inches(0.5), color=RED, head=False)
    add_arrow(s, x_6_mid, sy - Inches(0.5),
              x_6_mid, sy, color=RED)
    add_text(s, x_8_right + Inches(0.4), sy - Inches(0.65),
             Inches(1.5), Inches(0.35),
             "← cycle back", size=11, bold=True, color=RED, italic=True)

    # Problem statement
    add_text(s, Inches(0.9), Inches(3.35), Inches(12), Inches(0.45),
             "Problem:  Determine whether such a cycle exists — without knowing the length of the list.",
             size=16, bold=True, color=NAVY)

    # Two issue columns
    add_rect(s, Inches(0.9), Inches(3.95), Inches(5.7), Inches(2.65), LIGHT_GREY)
    add_text(s, Inches(1.1), Inches(4.0), Inches(5.3), Inches(0.5),
             "Why you can't just walk to the end", size=14, bold=True, color=RED)
    add_text(s, Inches(1.1), Inches(4.45), Inches(5.3), Inches(2.1),
             "A plain while current: loop will spin forever if there is a cycle.\n"
             "You have no termination signal — None is never reached.",
             size=13, color=DARK_GREY)

    add_rect(s, Inches(6.9), Inches(3.95), Inches(5.9), Inches(2.65), LIGHT_GREY)
    add_text(s, Inches(7.1), Inches(4.0), Inches(5.5), Inches(0.5),
             "Why you can't just count steps", size=14, bold=True, color=RED)
    add_text(s, Inches(7.1), Inches(4.45), Inches(5.5), Inches(2.1),
             "You don't know N — the list length is unknown.\n"
             "Any fixed limit will falsely report \"no cycle\" on large lists.",
             size=13, color=DARK_GREY)

    # key question
    callout_box(s, Inches(0.9), Inches(6.75), Inches(12.0), Inches(0.55),
                "The question:",
                "Can we detect a cycle in O(N) time with only O(1) extra memory — without modifying the list?",
                accent=GOLD)


def slide_naive(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "2.  Naive Approaches & Their Costs",
               "Every simpler solution has at least one fatal flaw", n, total)

    rows = [
        ("Approach", "How It Works", "Time", "Space", "Fatal Flaw"),
        ("Hash set of\nseen nodes",
         "Record every visited node's id(). If we see one twice, cycle found.",
         "O(N)", "O(N)", "Requires O(N) auxiliary memory — fails at scale."),
        ("Visited flag\n(mutate node)",
         "Add a boolean field visited to each node and mark as traversed.",
         "O(N)", "O(1)*", "Destructive — corrupts user data; not acceptable in libraries."),
        ("Step counter",
         "Abort traversal after N steps for some chosen N.",
         "O(N)", "O(1)", "Requires knowing N. Fragile; false negatives on large lists."),
        ("Floyd's\nTortoise & Hare",
         "Two pointers at speeds 1× and 2×. They collide iff a cycle exists.",
         "O(N)", "O(1)", "None — this is the correct solution. ✓"),
    ]

    x0 = Inches(0.6); y0 = Inches(1.5)
    col_w = [Inches(1.8), Inches(4.0), Inches(0.9), Inches(0.9), Inches(4.9)]
    rh_header = Inches(0.55); rh = Inches(0.88)

    for r, row in enumerate(rows):
        cy = y0 + (rh_header if r == 0 else rh_header + rh * (r - 1))
        cx = x0
        for c, txt in enumerate(row):
            is_header = (r == 0)
            is_solution = (r == len(rows) - 1)
            fill = NAVY if is_header else (
                   RGBColor(0xE8, 0xF8, 0xE8) if is_solution else
                   (LIGHT_GREY if r % 2 == 1 else WHITE))
            txt_color = WHITE if is_header else (GREEN if is_solution else DARK_GREY)
            bold = is_header or (c == 0)
            h = rh_header if is_header else rh
            add_rect(s, cx, cy, col_w[c], h, fill,
                     line=RGBColor(0xC0, 0xC4, 0xCC))
            add_text(s, cx + Inches(0.1), cy, col_w[c] - Inches(0.15), h,
                     txt, size=12, bold=bold, color=txt_color,
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[c]

    add_text(s, Inches(0.6), Inches(7.1), Inches(12), Inches(0.35),
             "* The visited-flag space is O(1) per se, but the flag alters the data structure — a side-effect we cannot accept.",
             size=11, italic=True, color=DARK_GREY)


def slide_intuition(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "3.  The Core Intuition — Two Runners on a Track",
               "Think of a circular running track, not a linked list", n, total)

    # Left column: text
    add_text(s, Inches(0.9), Inches(1.45), Inches(6.2), Inches(0.55),
             "The Runners Analogy", size=20, bold=True, color=NAVY)
    analogy = [
        ("Tortoise (slow).  ", "Runs 1 step per second around the track."),
        ("Hare (fast).  ",     "Runs 2 steps per second — always ahead."),
        ("Flat track.  ",      "They diverge forever; hare exits, tortoise follows."),
        ("Circular track.  ",  "Hare laps tortoise. They will always meet."),
    ]
    add_bullets(s, Inches(0.9), Inches(2.05), Inches(6.2), Inches(2.8),
                analogy, size=16, line_spacing=1.4)

    add_text(s, Inches(0.9), Inches(4.95), Inches(6.2), Inches(0.45),
             "Why they meet (intuition):", size=16, bold=True, color=TEAL)
    add_text(s, Inches(0.9), Inches(5.4), Inches(6.2), Inches(1.7),
             "Once both are inside the loop, the hare closes the gap by exactly "
             "1 node per iteration (it moves 2, the tortoise moves 1 — net gain 1). "
             "Starting from any gap d ≤ L, after d steps the gap is 0.",
             size=14, color=DARK_GREY)

    # Right column: circular track diagram
    # Draw a circle-ish set of nodes arranged in an arc
    add_rect(s, Inches(7.2), Inches(1.45), Inches(5.7), Inches(5.8), LIGHT_GREY)
    add_text(s, Inches(7.2), Inches(1.5), Inches(5.7), Inches(0.4),
             "Circular track — abstract view", size=13, bold=True,
             color=NAVY, align=PP_ALIGN.CENTER)

    # draw 6 nodes in a rough circle using computed positions
    import math
    cx_c = Inches(10.05); cy_c = Inches(4.45)
    radius = Inches(1.8)
    node_r = Inches(0.38)
    node_labels = ["A", "B", "C", "D", "E", "F"]
    positions = []
    for k in range(6):
        angle = math.pi / 2 - k * (2 * math.pi / 6)   # start top, go clockwise
        nx = cx_c + radius * math.cos(angle) - node_r
        ny = cy_c - radius * math.sin(angle) - node_r
        positions.append((nx, ny))
        fill = TEAL if k == 0 else WHITE
        add_node(s, nx, ny, node_r * 2, node_r * 2, node_labels[k],
                 fill=fill, border=NAVY, label_size=16)

    # arrows between consecutive nodes (clockwise)
    for k in range(6):
        x1n, y1n = positions[k]
        x2n, y2n = positions[(k + 1) % 6]
        # mid-points of node edges
        x1 = x1n + node_r; y1 = y1n + node_r
        x2 = x2n + node_r; y2 = y2n + node_r
        add_arrow(s, x1, y1, x2, y2, color=NAVY, weight=1.75)

    # tortoise marker (green circle on node B)
    tx, ty = positions[1]
    add_rrect(s, tx - Inches(0.12), ty - Inches(0.12),
              node_r * 2 + Inches(0.24), node_r * 2 + Inches(0.24),
              RGBColor(0x2E, 0x7D, 0x32), line=GREEN)
    add_text(s, tx - Inches(0.12), ty + node_r * 2 + Inches(0.03),
             node_r * 2 + Inches(0.24), Inches(0.28),
             "🐢 slow", size=10, bold=True, color=GREEN,
             align=PP_ALIGN.CENTER)

    # hare marker (orange circle on node D)
    hx, hy = positions[3]
    add_rrect(s, hx - Inches(0.12), hy - Inches(0.12),
              node_r * 2 + Inches(0.24), node_r * 2 + Inches(0.24),
              RGBColor(0xE6, 0x5C, 0x00), line=ORANGE)
    add_text(s, hx - Inches(0.12), hy + node_r * 2 + Inches(0.03),
             node_r * 2 + Inches(0.24), Inches(0.28),
             "🐇 fast", size=10, bold=True, color=ORANGE,
             align=PP_ALIGN.CENTER)


def _draw_trace_row(slide, y, nodes, slow_idx, fast_idx,
                    nw, nh, gp, sx, step_label, collision=False):
    """Helper: draw one row of the step-by-step trace."""
    add_text(slide, Inches(0.55), y, Inches(1.2), nh,
             step_label, size=12, bold=True, color=NAVY,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    for i, v in enumerate(nodes):
        bdr = RED if (collision and i == slow_idx) else NAVY
        fill = WHITE
        if i == slow_idx and i == fast_idx and collision:
            fill = RED
        elif i == slow_idx:
            fill = RGBColor(0xD0, 0xF0, 0xD0)   # light green for slow
        elif i == fast_idx:
            fill = RGBColor(0xFF, 0xE0, 0xB2)   # light orange for fast
        add_node(slide, sx + (nw + gp) * i, y, nw, nh, v,
                 fill=fill, border=bdr, label_size=13)
    for i in range(len(nodes) - 1):
        x1 = sx + (nw + gp) * i + nw
        x2 = sx + (nw + gp) * (i + 1)
        add_arrow(slide, x1, y + nh / 2, x2, y + nh / 2, weight=1.5)
    # cycle arrow: last node → node at index 2 (for the 1,2,3,4,5,6,7 example)
    # 7 → 3  (0-indexed: node 6 → node 2)
    last_x = sx + (nw + gp) * (len(nodes) - 1) + nw
    target_x = sx + (nw + gp) * 2 + nw / 2
    add_arrow(slide, last_x, y + nh / 2,
              last_x + Inches(0.2), y + nh / 2, color=RED, weight=1.5, head=False)
    add_arrow(slide, last_x + Inches(0.2), y + nh / 2,
              last_x + Inches(0.2), y - Inches(0.18), color=RED, weight=1.5, head=False)
    add_arrow(slide, last_x + Inches(0.2), y - Inches(0.18),
              target_x, y - Inches(0.18), color=RED, weight=1.5, head=False)
    add_arrow(slide, target_x, y - Inches(0.18),
              target_x, y, color=RED, weight=1.5)

    # pointer labels
    if slow_idx is not None:
        sx_label = sx + (nw + gp) * slow_idx + nw / 2
        add_text(slide, sx_label - Inches(0.3), y + nh + Inches(0.04),
                 Inches(0.6), Inches(0.25),
                 "S", size=11, bold=True,
                 color=GREEN if not collision else RED,
                 align=PP_ALIGN.CENTER)
    if fast_idx is not None:
        fx_label = sx + (nw + gp) * fast_idx + nw / 2
        add_text(slide, fx_label - Inches(0.3), y + nh + Inches(0.04),
                 Inches(0.6), Inches(0.25),
                 "F" if not (fast_idx == slow_idx and collision) else "S=F",
                 size=11, bold=True,
                 color=ORANGE if not collision else RED,
                 align=PP_ALIGN.CENTER)


def slide_trace(n, total):
    """Step-by-step trace through  1→2→3→4→5→6→7→3 (cycle at index 2)."""
    s = prs.slides.add_slide(BLANK)
    add_header(s, "4.  Step-by-Step Trace",
               "List:  1 → 2 → 3 → 4 → 5 → 6 → 7 → (back to 3)   |   S = slow  F = fast",
               n, total)

    nodes = ["1", "2", "3", "4", "5", "6", "7"]
    nw = Inches(0.82); nh = Inches(0.52); gp = Inches(0.22)
    sx = Inches(1.55)

    # Steps: (step_label, slow_idx, fast_idx)
    # Cycle: 7.next → 3 (index 2)
    # slow starts at 0, fast starts at 0 then moves to index 1 first
    steps = [
        ("Start",   0, 0),
        ("Step 1",  1, 2),   # slow→1, fast→2 (0-indexed)
        ("Step 2",  2, 4),
        ("Step 3",  3, 6),
        ("Step 4",  4, 2),   # fast wraps: 7→3 so index 2
        ("Step 5",  5, 4),
        ("Step 6 ✓", 6, 6),  # collision at index 6 (node "7")
    ]

    y_start = Inches(1.52)
    row_h = Inches(0.82)

    for idx, (label, si, fi) in enumerate(steps):
        collision = (si == fi and idx > 0)
        _draw_trace_row(s, y_start + row_h * idx, nodes, si, fi,
                        nw, nh, gp, sx, label, collision=collision)

    # legend
    add_rrect(s, SLIDE_W - Inches(1.85), Inches(1.52), Inches(1.45), Inches(1.6),
              LIGHT_GREY, line=DARK_GREY)
    add_text(s, SLIDE_W - Inches(1.8), Inches(1.58), Inches(1.35), Inches(0.35),
             "Legend", size=12, bold=True, color=NAVY)
    add_node(s, SLIDE_W - Inches(1.78), Inches(1.98),
             Inches(0.42), Inches(0.32), "S",
             fill=RGBColor(0xD0, 0xF0, 0xD0), border=NAVY, label_size=11)
    add_text(s, SLIDE_W - Inches(1.3), Inches(1.98), Inches(1.1), Inches(0.32),
             "= slow", size=11, color=DARK_GREY, anchor=MSO_ANCHOR.MIDDLE)
    add_node(s, SLIDE_W - Inches(1.78), Inches(2.35),
             Inches(0.42), Inches(0.32), "F",
             fill=RGBColor(0xFF, 0xE0, 0xB2), border=NAVY, label_size=11)
    add_text(s, SLIDE_W - Inches(1.3), Inches(2.35), Inches(1.1), Inches(0.32),
             "= fast", size=11, color=DARK_GREY, anchor=MSO_ANCHOR.MIDDLE)
    add_node(s, SLIDE_W - Inches(1.78), Inches(2.72),
             Inches(0.42), Inches(0.32), "!",
             fill=RED, border=RED, label_size=11)
    add_text(s, SLIDE_W - Inches(1.3), Inches(2.72), Inches(1.1), Inches(0.32),
             "= meet", size=11, color=DARK_GREY, anchor=MSO_ANCHOR.MIDDLE)


def slide_proof(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "5.  Correctness Proof — Why They Must Meet",
               "A tight mathematical argument for the collision guarantee", n, total)

    # Left: setup
    add_text(s, Inches(0.9), Inches(1.45), Inches(6.0), Inches(0.45),
             "Setup", size=18, bold=True, color=NAVY)

    setup_items = [
        ("F.  ", "Length of the acyclic prefix (nodes before the cycle)."),
        ("L.  ", "Length of the cycle."),
        ("At time T.  ", "Both pointers are inside the cycle."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.95), Inches(6.0), Inches(1.8),
                setup_items, size=15, line_spacing=1.35)

    # Proof steps
    add_text(s, Inches(0.9), Inches(3.85), Inches(6.0), Inches(0.45),
             "Proof", size=18, bold=True, color=NAVY)
    proof_items = [
        "Slow pointer enters the cycle at step F. At that moment it is at position 0 mod L inside the cycle.",
        "Fast pointer entered the cycle earlier and is at position (F mod L) inside the cycle.",
        "Relative speed = 2 – 1 = 1. Hare gains 1 node per iteration on the tortoise.",
        "The gap closes to 0 in at most L iterations. ∴ they collide.",
        "Total steps ≤ F + L = O(N). ∎",
    ]
    add_bullets(s, Inches(0.9), Inches(4.35), Inches(6.0), Inches(2.85),
                proof_items, size=14, line_spacing=1.35, bullet_color=GOLD)

    # Right: visual gap-closing diagram
    add_rect(s, Inches(7.2), Inches(1.45), Inches(5.7), Inches(5.8), LIGHT_GREY)
    add_text(s, Inches(7.4), Inches(1.52), Inches(5.3), Inches(0.4),
             "Gap-closing inside the loop  (L = 6)", size=13, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)

    gap_rows = [
        ("Iter 0", 0, 3, "Gap = 3"),
        ("Iter 1", 1, 4, "Gap = 2"),
        ("Iter 2", 2, 5, "Gap = 1"),
        ("Iter 3", 3, 6, "Gap = 0  ✓"),
    ]
    loop_labels = ["0", "1", "2", "3", "4", "5"]
    nw2 = Inches(0.6); nh2 = Inches(0.42); gp2 = Inches(0.1)
    sx2 = Inches(7.5)
    row_y0 = Inches(2.05)
    row_gap = Inches(0.98)

    for ri, (label, si2, fi2, gap_txt) in enumerate(gap_rows):
        ry = row_y0 + row_gap * ri
        add_text(s, Inches(7.22), ry, Inches(0.95), nh2,
                 label, size=11, bold=True, color=DARK_GREY,
                 anchor=MSO_ANCHOR.MIDDLE)
        for k, lbl in enumerate(loop_labels):
            fill = RGBColor(0xD0, 0xF0, 0xD0) if k == si2 % 6 else (
                   RGBColor(0xFF, 0xE0, 0xB2) if k == fi2 % 6 else WHITE)
            brd = RED if (si2 % 6 == fi2 % 6) else NAVY
            add_node(s, sx2 + (nw2 + gp2) * k, ry, nw2, nh2, lbl,
                     fill=fill, border=brd, label_size=11)
        # circular arrows between loop nodes
        for k in range(5):
            add_arrow(s, sx2 + (nw2 + gp2) * k + nw2, ry + nh2 / 2,
                      sx2 + (nw2 + gp2) * (k + 1), ry + nh2 / 2,
                      color=NAVY, weight=1.0)
        # gap annotation
        gap_col = RED if "0  ✓" in gap_txt else TEAL
        add_text(s, sx2 + (nw2 + gp2) * 6 + Inches(0.05), ry,
                 Inches(1.3), nh2, gap_txt,
                 size=11, bold=True, color=gap_col,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Complexity box
    callout_box(s, Inches(7.2), Inches(6.1), Inches(5.7), Inches(0.7),
                "Key insight:",
                "The relative speed of 1 guarantees the gap reaches exactly 0 (not skips over). "
                "If fast moved 3 steps and slow moved 1, the gap could be skipped!",
                accent=RED)


def slide_code(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "6.  The Code — has_cycle()",
               "Complete, annotated Python implementation", n, total)

    code = (
        "def has_cycle(head: Node) -> bool:\n"
        "    \"\"\"\n"
        "    Detect a cycle in a singly linked list.\n"
        "    Time  O(N)  ·  Space  O(1)  ·  Non-mutating.\n"
        "    \"\"\"\n"
        "    # Edge case: empty or single-node list cannot form a cycle.\n"
        "    if head is None or head.next is None:\n"
        "        return False\n"
        "\n"
        "    slow = head           # tortoise — advances 1 step per iteration\n"
        "    fast = head.next      # hare     — advances 2 steps per iteration\n"
        "    # Note: fast starts one ahead to avoid an immediate false positive.\n"
        "\n"
        "    while slow != fast:\n"
        "        # If the hare reaches the end, there is no cycle.\n"
        "        if fast is None or fast.next is None:\n"
        "            return False\n"
        "\n"
        "        slow = slow.next           # tortoise: 1 step\n"
        "        fast = fast.next.next      # hare:     2 steps\n"
        "\n"
        "    # Loop exited because slow == fast → collision → cycle exists.\n"
        "    return True"
    )
    add_code_block(s, Inches(0.9), Inches(1.48), Inches(8.0), Inches(5.5),
                   code, size=13)

    # Right: line-by-line commentary
    add_text(s, Inches(9.2), Inches(1.48), Inches(3.85), Inches(0.45),
             "Line-by-line notes", size=17, bold=True, color=NAVY)
    notes = [
        ("Edge guard.  ",
         "Prevents NoneType errors on empty or one-node lists — always first."),
        ("fast = head.next.  ",
         "Offset by 1 so the while condition is valid immediately."),
        ("while slow != fast.  ",
         "Identity check (is) compares node addresses, not values."),
        ("fast is None.  ",
         "Hare stepped off the end — no cycle."),
        ("fast.next is None.  ",
         "Hare would step off next iteration — stop early."),
        ("slow.next / fast.next.next.  ",
         "The 1×/2× speed differential — the heart of the algorithm."),
        ("return True.  ",
         "Collision guaranteed; cycle confirmed in O(N)."),
    ]
    add_bullets(s, Inches(9.2), Inches(2.0), Inches(3.85), Inches(5.0),
                notes, size=12, line_spacing=1.35, bullet_color=GOLD)


def slide_entry_point(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "7.  Floyd Part II — Finding the Cycle Entry",
               "Where exactly does the loop begin? O(N) time, O(1) space.", n, total)

    add_text(s, Inches(0.9), Inches(1.45), Inches(12), Inches(0.55),
             "After collision is confirmed, a second phase pinpoints the entry node:",
             size=16, color=DARK_GREY)

    code = (
        "def find_cycle_entry(head: Node) -> Node | None:\n"
        "    \"\"\"\n"
        "    Returns the node where the cycle begins, or None if no cycle.\n"
        "    Relies on Floyd's mathematical property: after meeting, reset one\n"
        "    pointer to head and advance both at speed 1 — they meet at entry.\n"
        "    \"\"\"\n"
        "    if head is None or head.next is None:\n"
        "        return None\n"
        "\n"
        "    slow = fast = head\n"
        "\n"
        "    # Phase 1: detect collision\n"
        "    while True:\n"
        "        if fast is None or fast.next is None:\n"
        "            return None          # no cycle\n"
        "        slow = slow.next\n"
        "        fast = fast.next.next\n"
        "        if slow == fast:\n"
        "            break\n"
        "\n"
        "    # Phase 2: find entry\n"
        "    slow = head              # reset slow to head\n"
        "    while slow != fast:      # both move 1 step — they meet at entry\n"
        "        slow = slow.next\n"
        "        fast = fast.next\n"
        "    return slow              # the cycle entry node"
    )
    add_code_block(s, Inches(0.9), Inches(2.1), Inches(7.8), Inches(5.1),
                   code, size=12)

    # Right: why it works
    add_text(s, Inches(9.0), Inches(2.1), Inches(4.0), Inches(0.45),
             "Why Phase 2 works", size=17, bold=True, color=NAVY)
    why = [
        "Let F = prefix length, L = cycle length.",
        "At collision: slow has moved F + a steps; fast has moved F + a + kL steps (some integer k).",
        "Reset slow to head. Both advance 1 step / iteration.",
        "After F more steps, slow reaches entry. Fast (from collision point) also reaches entry — proven by modular arithmetic.",
        "They meet exactly at the cycle entry. ∎",
    ]
    add_bullets(s, Inches(9.0), Inches(2.65), Inches(4.1), Inches(3.5),
                why, size=13, line_spacing=1.35, bullet_color=TEAL)

    callout_box(s, Inches(9.0), Inches(6.2), Inches(4.1), Inches(0.9),
                "LeetCode 142",
                "\"Linked List Cycle II\" — directly tests this exact algorithm. "
                "Expected: O(N) / O(1). Standard solution is Floyd Part II.",
                accent=GOLD)


def slide_complexity(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "8.  Complexity Analysis",
               "Tight bounds — not just big-O hand-waving", n, total)

    # Two columns: Time and Space
    for col_i, (heading, color, content) in enumerate([
        ("Time Complexity  O(N)", TEAL, [
            "Phase 1 — Prefix traversal: O(F) steps until slow enters the cycle.",
            "Phase 1 — Cycle traversal: at most O(L) additional steps to meet.",
            "Since F + L ≤ N (total list length), Phase 1 is O(N).",
            "Phase 2 — Re-traverse prefix: O(F) ≤ O(N) steps.",
            "Overall: O(N) — linear in list length.",
        ]),
        ("Space Complexity  O(1)", GOLD, [
            "Only two pointer variables (slow, fast) regardless of list size.",
            "No hash set, no stack, no visited array.",
            "Memory usage is constant — scales to billion-node lists.",
            "Compare: hash-set approach uses O(N) bytes — catastrophic at scale.",
        ]),
    ]):
        cx = Inches(0.6) + Inches(6.3) * col_i
        add_rect(s, cx, Inches(1.5), Inches(6.0), Inches(0.6), color)
        add_text(s, cx, Inches(1.5), Inches(6.0), Inches(0.6),
                 heading, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, cx, Inches(2.1), Inches(6.0), Inches(4.0), LIGHT_GREY)
        add_bullets(s, cx + Inches(0.2), Inches(2.25),
                    Inches(5.6), Inches(3.7),
                    content, size=14, line_spacing=1.4,
                    bullet_color=color)

    # Comparison table
    add_text(s, Inches(0.6), Inches(6.25), Inches(12), Inches(0.42),
             "Algorithm Comparison", size=16, bold=True, color=NAVY)
    rows2 = [
        ("Method", "Time", "Space", "Mutates?", "Needs N?"),
        ("Hash Set",        "O(N)", "O(N)", "No",  "No"),
        ("Visited Flag",    "O(N)", "O(1)", "Yes", "No"),
        ("Step Limit",      "O(N)", "O(1)", "No",  "Yes"),
        ("Floyd's T&H",     "O(N)", "O(1)", "No",  "No"),
    ]
    x0 = Inches(0.6); y0 = Inches(6.7)
    cw2 = [Inches(3.2), Inches(1.6), Inches(1.6), Inches(1.8), Inches(1.8)]
    rh2 = Inches(0.42)
    for r, row in enumerate(rows2):
        cx2 = x0
        for c, txt in enumerate(row):
            fill2 = NAVY if r == 0 else (
                    RGBColor(0xE8, 0xF8, 0xE8) if r == 4 else
                    (LIGHT_GREY if r % 2 == 0 else WHITE))
            col2 = WHITE if r == 0 else (GREEN if r == 4 else DARK_GREY)
            add_rect(s, cx2, y0 + rh2 * r, cw2[c], rh2,
                     fill2, line=RGBColor(0xC0, 0xC4, 0xCC))
            add_text(s, cx2 + Inches(0.08), y0 + rh2 * r,
                     cw2[c] - Inches(0.1), rh2,
                     txt, size=12, bold=(r == 0 or c == 0),
                     color=col2, anchor=MSO_ANCHOR.MIDDLE)
            cx2 += cw2[c]


def slide_variants(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "9.  Variants & Extensions of the Technique",
               "The same two-pointer pattern solves a surprising range of problems", n, total)

    cards = [
        ("Find Middle of List", NAVY,
         "Move slow 1×, fast 2×. When fast reaches the end, slow is at the midpoint.\n"
         "Used in merge-sort on linked lists. O(N) / O(1)."),
        ("Happy Number\n(LeetCode 202)", TEAL,
         "A number is 'happy' if repeatedly summing squares of digits reaches 1.\n"
         "Non-happy numbers cycle. Detect with Floyd on the digit-sum sequence."),
        ("Find Duplicate\n(LeetCode 287)", GOLD,
         "Array of N+1 integers in [1, N]. Treat values as next-pointers.\n"
         "The duplicate creates a cycle. Floyd finds it in O(N) / O(1)."),
        ("Pollar's Rho\nFactorisation", PURPLE,
         "Pollard (1975) used Floyd's cycle detection to factorise large integers.\n"
         "Powers modern cryptographic attacks on weak RSA keys."),
        ("Nth Node From End", GREEN,
         "Move fast N steps ahead, then advance both at 1×.\n"
         "When fast reaches end, slow is at the Nth-from-last node."),
        ("Palindrome Check", RED,
         "Find middle (Floyd), reverse second half, compare.\n"
         "O(N) time, O(1) space — no auxiliary stack needed."),
    ]

    cw = Inches(4.0); ch = Inches(2.4); gx = Inches(0.25); gy = Inches(0.2)
    x0 = Inches(0.6); y0 = Inches(1.55)
    for i, (title, color, body) in enumerate(cards):
        col = i % 3; row = i // 3
        cx = x0 + (cw + gx) * col
        cy = y0 + (ch + gy) * row
        add_rect(s, cx, cy, cw, Inches(0.55), color)
        add_text(s, cx, cy, cw, Inches(0.55), title,
                 size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, cx, cy + Inches(0.55), cw, ch - Inches(0.55), LIGHT_GREY)
        add_text(s, cx + Inches(0.18), cy + Inches(0.72),
                 cw - Inches(0.35), ch - Inches(0.8),
                 body, size=12, color=DARK_GREY)


def slide_applications(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "10.  Real-World Applications",
               "From OS kernels to cryptography — Floyd's algorithm runs everywhere", n, total)

    apps = [
        ("Operating Systems", NAVY,
         "The Linux kernel uses cycle-detection variants when walking page-table chains "
         "and process-wait graphs to identify deadlocks."),
        ("Cryptography", TEAL,
         "Pollard's rho algorithm — built on Floyd's — breaks RSA keys by finding "
         "non-trivial factors of composite numbers efficiently."),
        ("Compilers & Linkers", GOLD,
         "Dependency graphs may contain circular imports. Compilers detect these cycles "
         "before attempting topological ordering."),
        ("Garbage Collectors", RED,
         "Reference-counting GCs need cycle detection to collect objects in "
         "reference cycles. CPython's gc module does exactly this."),
        ("Network Protocols", GREEN,
         "Routing loops (packets cycling indefinitely) are detected by tracking "
         "the path signature — conceptually equivalent to Floyd's."),
        ("Databases", PURPLE,
         "Transaction deadlock detection — find cycles in the wait-for graph, "
         "then abort the youngest transaction."),
    ]

    cw = Inches(4.1); ch = Inches(2.35); gx = Inches(0.17); gy = Inches(0.2)
    x0 = Inches(0.55); y0 = Inches(1.55)
    for i, (title, color, body) in enumerate(apps):
        col = i % 3; row = i // 3
        cx = x0 + (cw + gx) * col
        cy = y0 + (ch + gy) * row
        add_rect(s, cx, cy, cw, Inches(0.52), color)
        add_text(s, cx, cy, cw, Inches(0.52), title,
                 size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, cx, cy + Inches(0.52), cw, ch - Inches(0.52), LIGHT_GREY)
        add_text(s, cx + Inches(0.18), cy + Inches(0.68),
                 cw - Inches(0.35), ch - Inches(0.75),
                 body, size=12, color=DARK_GREY)


def slide_practice(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "11.  Practice Problems",
               "Consolidate your understanding — solve these before the next class", n, total)

    problems = [
        ("Easy — LeetCode 141", TEAL,
         "Linked List Cycle\n"
         "Implement has_cycle(). Expected: O(N) time, O(1) space.\n"
         "Edge cases: empty list, single node, two-node cycle."),
        ("Medium — LeetCode 142", GOLD,
         "Linked List Cycle II\n"
         "Return the node where the cycle begins, or None.\n"
         "Implement Floyd Part II (find_cycle_entry)."),
        ("Medium — LeetCode 287", NAVY,
         "Find the Duplicate Number\n"
         "Array of N+1 ints in [1, N]. Find the duplicate.\n"
         "Constraint: O(1) space — no sorting, no hash set."),
        ("Easy — LeetCode 202", GREEN,
         "Happy Number\n"
         "Detect whether the digit-sum sequence cycles.\n"
         "Apply Floyd to a non-pointer sequence."),
        ("Medium — LeetCode 876", PURPLE,
         "Middle of the Linked List\n"
         "Return the middle node (second middle if even length).\n"
         "One-pass, O(1) space using fast/slow pointers."),
        ("Hard — Conceptual", RED,
         "Prove that Floyd's algorithm requires at most F + L steps in Phase 1.\n"
         "Write your proof formally using modular arithmetic."),
    ]

    cw = Inches(4.0); ch = Inches(2.35); gx = Inches(0.25); gy = Inches(0.2)
    x0 = Inches(0.6); y0 = Inches(1.55)
    for i, (title, color, body) in enumerate(problems):
        col = i % 3; row = i // 3
        cx = x0 + (cw + gx) * col
        cy = y0 + (ch + gy) * row
        add_rect(s, cx, cy, cw, Inches(0.52), color)
        add_text(s, cx, cy, cw, Inches(0.52), title,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, cx, cy + Inches(0.52), cw, ch - Inches(0.52), LIGHT_GREY)
        add_text(s, cx + Inches(0.18), cy + Inches(0.68),
                 cw - Inches(0.35), ch - Inches(0.75),
                 body, size=12, color=DARK_GREY)


def slide_summary(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "12.  Summary — Five Things to Remember",
               "Recite these before leaving the room", n, total)
    items = [
        "Floyd's algorithm uses two pointers at speeds 1× and 2×. The hare is always ahead of (or with) the tortoise.",
        "If a cycle exists, both pointers enter the loop and the hare closes the gap by exactly 1 per step — collision is guaranteed.",
        "The algorithm runs in O(N) time and O(1) space. It does not mutate the list and does not require knowing N.",
        "Phase II (resetting slow to head, both at speed 1) finds the precise cycle entry point in O(F) additional steps.",
        "The same two-pointer pattern solves: middle-of-list, Nth-from-end, Happy Number, duplicate detection, and more.",
    ]
    add_bullets(s, Inches(0.9), Inches(1.55), Inches(12.1), Inches(5.55),
                items, size=18, line_spacing=1.45, bullet_color=GOLD)


def slide_references(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "13.  References & Further Reading",
               "Primary sources and practice resources", n, total)
    refs = [
        ("Floyd, R. W. (1967).  ",
         "Nondeterministic Algorithms. JACM 14(4):636–644. Original publication of the cycle-finding method."),
        ("Knuth, D. E.  ",
         "The Art of Computer Programming, Vol. 2, §3.1. Rho sequence and cycle detection."),
        ("CLRS 4th Ed.  ",
         "Cormen et al. Introduction to Algorithms. Chapter on Elementary Data Structures."),
        ("Pollard, J. M. (1975).  ",
         "A Monte Carlo Method for Factorization. BIT 15(3):331–334. Floyd applied to cryptography."),
        ("LeetCode Problems.  ",
         "141 (Cycle), 142 (Cycle II), 287 (Duplicate), 202 (Happy), 876 (Middle)."),
        ("Competitive Programmer's Handbook.  ",
         "Antti Laaksonen. Chapter on graph algorithms — cycle detection in functional graphs."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.55), Inches(12.1), Inches(5.6),
                refs, size=16, line_spacing=1.4)


def slide_thanks():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(5.6), SLIDE_W, Inches(0.14), GOLD)
    add_rect(s, 0, Inches(5.8), SLIDE_W, Inches(0.05), TEAL)
    add_text(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.5),
             "Questions?", size=80, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.2), Inches(11.5), Inches(0.7),
             "Bring your traces, your proofs, and your failing test cases.",
             size=22, italic=True, color=GOLD)
    add_text(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.5),
             "Remember:  🐢  one step at a time is still enough to find the loop.",
             size=16, color=LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.4),
             "Next topic  ·  Stacks & Queues — bounded chains with discipline.",
             size=15, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.4),
             "Data Structures & Algorithms  ·  Week 03 Supplement  ·  Floyd's Tortoise & Hare",
             size=11, color=LIGHT_GREY, italic=True)


# ── Build ────────────────────────────────────────────────────────────────────

content_slides = [
    slide_outline,
    slide_objectives,
    slide_motivation,
    slide_naive,
    slide_intuition,
    slide_trace,
    slide_proof,
    slide_code,
    slide_entry_point,
    slide_complexity,
    slide_variants,
    slide_applications,
    slide_practice,
    slide_summary,
    slide_references,
]

TOTAL = 1 + len(content_slides) + 1   # title + content + thanks

slide_title()
for i, fn in enumerate(content_slides, start=2):
    fn(i, TOTAL)
slide_thanks()

OUT = os.path.join(os.path.dirname(__file__), "floyd_tortoise_hare.pptx")
prs.save(OUT)
print(f"✓  Wrote {OUT}  ({len(prs.slides)} slides)")
