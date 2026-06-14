"""
Generate a university-grade PowerPoint deck for Week 03 Supplement: Stacks & Queues.

Output: stacks_queues_dsa.pptx (16:9 widescreen)
Author : DSA Curriculum
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# -------- Theme --------
NAVY        = RGBColor(0x0B, 0x2A, 0x4A)   # deep navy (primary)
TEAL        = RGBColor(0x12, 0x7A, 0x8A)   # accent
GOLD        = RGBColor(0xE0, 0xA8, 0x1F)   # accent
LIGHT_GREY  = RGBColor(0xF2, 0xF4, 0xF7)   # background panels
DARK_GREY   = RGBColor(0x33, 0x3D, 0x4A)   # body text
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xC0, 0x39, 0x2B)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
PURPLE      = RGBColor(0x6A, 0x1B, 0x9A)
ORANGE      = RGBColor(0xE6, 0x5C, 0x00)
CODE_BG     = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG     = RGBColor(0xE6, 0xE6, 0xE6)

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None, line_width=1):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    return shp


def add_rrect(slide, x, y, w, h, fill, line=None, line_width=2):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=18, color=DARK_GREY,
                bullet_color=TEAL, line_spacing=1.25, font="Calibri", indent=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        glyph = "  –  " if indent else "▸  "
        rb = p.add_run()
        rb.text = glyph
        rb.font.name = font
        rb.font.size = Pt(size)
        rb.font.bold = not indent
        rb.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run()
            r1.text = head
            r1.font.name = font
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = tail
            r2.font.name = font
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None, page_num=None, total=None):
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    add_rect(slide, Inches(0.35), Inches(0.0), SLIDE_W - Inches(0.35), Inches(0.08), GOLD)
    add_text(slide, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.78), Inches(11.5), Inches(0.45),
                 subtitle, size=15, italic=True, color=TEAL)
    add_rect(slide, Inches(0.6), Inches(1.22), Inches(2.0), Inches(0.04), GOLD)
    add_text(slide, Inches(0.6), SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
             "Data Structures & Algorithms · Week 03 Supplement · Stacks & Queues",
             size=10, color=DARK_GREY, italic=True)
    if page_num is not None and total is not None:
        add_text(slide, SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.45),
                 Inches(1.2), Inches(0.3),
                 f"Slide {page_num} / {total}", size=10, color=DARK_GREY,
                 align=PP_ALIGN.RIGHT)


def add_code_block(slide, x, y, w, h, code, size=14):
    add_rect(slide, x, y, w, h, CODE_BG)
    add_rect(slide, x, y, Inches(0.08), h, GOLD)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.25), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = code.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = ln if ln else " "
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = CODE_FG


def add_arrow(slide, x1, y1, x2, y2, color=TEAL, weight=2.25):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(
        ln, qn("a:tailEnd"),
        {"type": "triangle", "w": "med", "len": "med"}
    )
    return conn


def add_node_box(slide, x, y, w, h, value, *, fill=WHITE, border=NAVY,
                  value_size=20, has_next=True, has_prev=False):
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    outer.fill.solid()
    outer.fill.fore_color.rgb = fill
    outer.line.color.rgb = border
    outer.line.width = Pt(2)
    outer.shadow.inherit = False
    
    parts = (1 if has_prev else 0) + 1 + (1 if has_next else 0)
    cell_w = w / parts
    idx = 0
    if has_prev:
        add_rect(slide, x + Emu(0), y, cell_w, h, LIGHT_GREY, line=border)
        add_text(slide, x, y, cell_w, h, "prev", size=12, color=DARK_GREY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        idx += 1
    add_text(slide, x + cell_w * idx, y, cell_w, h, str(value),
             size=value_size, bold=True, color=NAVY if fill == WHITE else WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if has_next:
        idx2 = idx + 1
        add_rect(slide, x + cell_w * idx2, y, cell_w, h, LIGHT_GREY, line=border)
        add_text(slide, x + cell_w * idx2, y, cell_w, h, "next",
                 size=12, color=DARK_GREY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for k in range(1, parts):
        sep_x = x + cell_w * k
        line = slide.shapes.add_connector(1, sep_x, y, sep_x, y + h)
        line.line.color.rgb = border
        line.line.width = Pt(1.25)


# ---------- Slide Builders ----------

def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(5.5), SLIDE_W, Inches(0.12), GOLD)
    add_rect(s, 0, Inches(5.7), SLIDE_W, Inches(0.04), TEAL)
    
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.9), Inches(1.0),
                              Inches(3.8), Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb = TEAL
    chip.line.fill.background(); chip.shadow.inherit = False
    add_text(s, Inches(0.9), Inches(1.0), Inches(3.8), Inches(0.55),
             "WEEK 03  ·  SUPPLEMENTARY DECK", size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
             
    add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(1.4),
             "Stacks & Queues", size=64, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.7),
             "Understanding Linear Abstract Data Types and Implementations",
             size=26, color=GOLD, italic=True)
    add_text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.6),
             "A Deep-Dive into LIFO, FIFO, and Memory Layouts",
             size=18, color=LIGHT_GREY)
             
    add_text(s, Inches(0.9), Inches(6.0), Inches(8), Inches(0.4),
             "Course   ·   CSC: Data Structures & Algorithms",
             size=14, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.4), Inches(8), Inches(0.4),
             "Module   ·   Week 03 Supplement",
             size=14, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.8), Inches(8), Inches(0.4),
             "Format   ·   Lecture · Visual Demonstration",
             size=14, color=WHITE)
             
    # stack icon visual
    add_text(s, Inches(10.5), Inches(1.8), Inches(2.5), Inches(1.2),
             "🥞", size=72, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(10.5), Inches(3.2), Inches(2.5), Inches(1.2),
             "🚶🚶🚶", size=48, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_outline(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Lecture Outline",
               "What we will cover today about Stacks and Queues", n, total)
    items = [
        ("1. Abstract Data Types (ADTs).  ", "Separation of behavior (interface) from memory (implementation)."),
        ("2. The Stack (LIFO).  ", "Plate-stack intuition, operations, and visual representations."),
        ("3. Stack Implementation.  ", "Building stacks using Linked Lists and Dynamic Arrays (Big-O comparison)."),
        ("4. The Queue (FIFO).  ", "Line checkout intuition, operations, and visual layouts."),
        ("5. Queue Implementation.  ", "Double-pointer Linked Lists and Circular Array Buffers."),
        ("6. Double-Ended Queues (Deques).  ", "Adding and removing from both ends."),
        ("7. Real-World Applications.  ", "Call stacks, BFS, buffer queues, print jobs, and backtracking."),
        ("8. Practice Problems.  ", "Valid Parentheses, Stack-using-Queue, and min stacks."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(12), Inches(5.5),
                items, size=18, line_spacing=1.35)


def slide_adt_concept(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "1.  Abstract Data Types (ADTs)",
               "Decoupling the 'What' from the 'How'", n, total)
               
    add_text(s, Inches(0.9), Inches(1.45), Inches(6.0), Inches(0.5),
             "What is an ADT?", size=22, bold=True, color=NAVY)
    
    items = [
        ("Interface vs. Implementation.  ", "An ADT specifies what operations are supported, not how they are coded."),
        ("Encapsulation.  ", "Users interact with a clean API, hiding structural details (pointers, arrays)."),
        ("Flexibility.  ", "We can swap a Stack's backend from an array to a linked list without breaking user code."),
    ]
    add_bullets(s, Inches(0.9), Inches(2.0), Inches(6.0), Inches(4.5),
                items, size=16, line_spacing=1.35)
                
    # Right panel: visual representation of ADT encapsulation
    add_rect(s, Inches(7.3), Inches(1.45), Inches(5.4), Inches(5.3), LIGHT_GREY)
    add_text(s, Inches(7.3), Inches(1.7), Inches(5.4), Inches(0.4),
             "The ADT Shield", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
             
    # API Box
    add_rrect(s, Inches(8.5), Inches(2.3), Inches(3.0), Inches(1.0), TEAL)
    add_text(s, Inches(8.5), Inches(2.3), Inches(3.0), Inches(1.0),
             "USER API\n(push, pop, enqueue)", size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
             
    add_arrow(s, Inches(10.0), Inches(3.3), Inches(10.0), Inches(4.1))
    
    # Implementation Box
    add_rrect(s, Inches(8.0), Inches(4.1), Inches(4.0), Inches(1.8), NAVY)
    add_text(s, Inches(8.0), Inches(4.1), Inches(4.0), Inches(1.8),
             "Concrete Implementations\n\nOption A: Array / List\nOption B: Node Linked Chain",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_stack_concept(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "2.  The Stack (LIFO)",
               "Last-In, First-Out (LIFO) Structure", n, total)
               
    add_text(s, Inches(0.9), Inches(1.45), Inches(6.0), Inches(0.5),
             "The Stack Metaphor", size=22, bold=True, color=NAVY)
    
    items = [
        ("Plate Stack.  ", "You add plates to the top. You remove plates from the top. Trying to pull from the bottom crashes the stack."),
        ("LIFO Policy.  ", "The element most recently pushed onto the stack is the first one popped off."),
        ("Restricted Access.  ", "No random access allowed. You cannot inspect or delete elements in the middle without popping the top first."),
    ]
    add_bullets(s, Inches(0.9), Inches(2.0), Inches(6.0), Inches(4.5),
                items, size=16, line_spacing=1.4)
                
    # Right panel: visual representation of a stack
    add_rect(s, Inches(7.3), Inches(1.45), Inches(5.4), Inches(5.3), LIGHT_GREY)
    add_text(s, Inches(7.3), Inches(1.6), Inches(5.4), Inches(0.4),
             "Stack Visualization", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
             
    # Draw stack bucket
    add_rect(s, Inches(8.8), Inches(2.2), Inches(0.08), Inches(4.0), DARK_GREY)
    add_rect(s, Inches(11.2), Inches(2.2), Inches(0.08), Inches(4.0), DARK_GREY)
    add_rect(s, Inches(8.8), Inches(6.2), Inches(2.48), Inches(0.08), DARK_GREY)
    
    # Push elements in the bucket
    elements = ["10 (Bottom)", "20", "30 (Top)"]
    ey = Inches(5.2)
    for i, el in enumerate(elements):
        col = GOLD if i == 2 else WHITE
        bdr = NAVY
        add_rrect(s, Inches(9.0), ey, Inches(2.0), Inches(0.8), col, line=bdr)
        add_text(s, Inches(9.0), ey, Inches(2.0), Inches(0.8), el,
                 size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        ey -= Inches(1.0)
        
    # Push/pop arrows
    add_arrow(s, Inches(10.0), Inches(1.3), Inches(10.0), Inches(2.0))
    add_text(s, Inches(10.1), Inches(1.4), Inches(1.5), Inches(0.4),
             "Push / Pop here", size=12, bold=True, color=TEAL)


def slide_stack_ops(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Stack Operations & Complexity",
               "The Core Interface of a Stack", n, total)
               
    # Table of operations
    rows = [
        ("Operation", "Description", "Complexity"),
        ("push(item)", "Insert item onto the top of the stack.", "O(1)"),
        ("pop()", "Remove and return the item at the top of the stack.", "O(1)"),
        ("peek() / top()", "Return the item at the top without removing it.", "O(1)"),
        ("is_empty()", "Return True if the stack has no elements, else False.", "O(1)"),
        ("size()", "Return the count of elements currently on the stack.", "O(1)"),
    ]
    
    x0 = Inches(0.9); y0 = Inches(1.6)
    col_w = [Inches(3.2), Inches(5.8), Inches(2.5)]
    rh = Inches(0.7)
    
    for r, row in enumerate(rows):
        cx = x0
        for c, txt in enumerate(row):
            fill = NAVY if r == 0 else (LIGHT_GREY if r % 2 == 1 else WHITE)
            color = WHITE if r == 0 else (GREEN if c == 2 and r > 0 else DARK_GREY)
            bold = (r == 0) or (c == 0) or (c == 2)
            add_rect(s, cx, y0 + rh * r, col_w[c], rh, fill,
                     line=RGBColor(0xC0, 0xC4, 0xCC))
            add_text(s, cx + Inches(0.15), y0 + rh * r,
                     col_w[c] - Inches(0.2), rh,
                     txt, size=15, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[c]
            
    # Key concept callout box
    callout_box_y = Inches(1.6 + 0.7 * len(rows) + 0.3)
    add_rect(s, Inches(0.9), callout_box_y, Inches(11.5), Inches(1.1), LIGHT_GREY)
    add_rect(s, Inches(0.9), callout_box_y, Inches(0.08), Inches(1.1), GOLD)
    add_text(s, Inches(1.1), callout_box_y + Inches(0.1), Inches(11.0), Inches(0.35),
             "Important Rule:", size=14, bold=True, color=GOLD)
    add_text(s, Inches(1.1), callout_box_y + Inches(0.45), Inches(11.0), Inches(0.6),
             "A correct implementation of a Stack MUST execute push, pop, and peek in O(1) time. "
             "If any of these operations scale linearly with list size O(N), the implementation is incorrect.",
             size=13, color=DARK_GREY)


def slide_stack_linked_list(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "3.  Stack Implementation: Linked List",
               "Singly Linked List with insertions/deletions at Head", n, total)
               
    code = (
        "class LinkedStack:\n"
        "    class Node:\n"
        "        def __init__(self, data, next_node=None):\n"
        "            self.data = data\n"
        "            self.next = next_node\n"
        "\n"
        "    def __init__(self):\n"
        "        self.head = None\n"
        "\n"
        "    def push(self, data):\n"
        "        # Prepend: insert at the head of the list\n"
        "        self.head = self.Node(data, self.head)\n"
        "\n"
        "    def pop(self):\n"
        "        if self.head is None:\n"
        "            raise IndexError(\"Pop from empty stack\")\n"
        "        val = self.head.data\n"
        "        self.head = self.head.next   # shift head\n"
        "        return val"
    )
    add_code_block(s, Inches(0.9), Inches(1.48), Inches(7.5), Inches(5.5), code, size=13)
    
    # Right notes column
    add_text(s, Inches(8.7), Inches(1.48), Inches(4.0), Inches(0.45),
             "Design Insights", size=18, bold=True, color=NAVY)
    notes = [
        ("Prepend pushes.  ", "By inserting at the head, we avoid traversing the entire list. Complexity is O(1)."),
        ("Head deletions.  ", "Popping removes the head node, which is also O(1)."),
        ("No Capacity Limits.  ", "Memory is allocated dynamically per Node, avoiding resize halts."),
        ("Pointer Overhead.  ", "Each node requires extra memory to store the reference to .next."),
    ]
    add_bullets(s, Inches(8.7), Inches(2.0), Inches(4.0), Inches(4.8),
                notes, size=13, line_spacing=1.35, bullet_color=GOLD)


def slide_stack_array(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Stack Implementation: Dynamic Array",
               "Array/List Backend with operations at the End", n, total)
               
    code = (
        "class ArrayStack:\n"
        "    def __init__(self):\n"
        "        self._data = []   # dynamic list representation\n"
        "\n"
        "    def push(self, data):\n"
        "        # append is amortized O(1)\n"
        "        self._data.append(data)\n"
        "\n"
        "    def pop(self):\n"
        "        if self.is_empty():\n"
        "            raise IndexError(\"Pop from empty stack\")\n"
        "        # pop() from end is O(1)\n"
        "        return self._data.pop()\n"
        "\n"
        "    def is_empty(self):\n"
        "        return len(self._data) == 0"
    )
    add_code_block(s, Inches(0.9), Inches(1.48), Inches(7.5), Inches(5.5), code, size=13)
    
    # Right notes column
    add_text(s, Inches(8.7), Inches(1.48), Inches(4.0), Inches(0.45),
             "Array vs. Linked List", size=18, bold=True, color=NAVY)
    notes = [
        ("End of Array.  ", "Inserting/deleting at the end avoids element shifting. If we inserted at index 0, it would be O(N)."),
        ("Cache Locality.  ", "Arrays are stored contiguously in memory, making reads faster on hardware due to caching."),
        ("Amortized Cost.  ", "Push is amortized O(1), but occasionally triggers an O(N) resize event."),
        ("Pre-allocation waste.  ", "Unused slots in the array capacity consume memory."),
    ]
    add_bullets(s, Inches(8.7), Inches(2.0), Inches(4.0), Inches(4.8),
                notes, size=13, line_spacing=1.35, bullet_color=TEAL)


def slide_queue_concept(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "4.  The Queue (FIFO)",
               "First-In, First-Out (FIFO) Structure", n, total)
               
    add_text(s, Inches(0.9), Inches(1.45), Inches(6.0), Inches(0.5),
             "The Checkout Line Metaphor", size=22, bold=True, color=NAVY)
    
    items = [
        ("First-come, First-served.  ", "New elements join at the rear. Elements leave from the front."),
        ("FIFO Policy.  ", "The element that has been in the queue the longest is the next one to be removed."),
        ("Dual Endpoints.  ", "Requires two pointer references: FRONT (for popping) and REAR/TAIL (for pushing)."),
    ]
    add_bullets(s, Inches(0.9), Inches(2.0), Inches(6.0), Inches(4.5),
                items, size=16, line_spacing=1.4)
                
    # Right panel: visual representation of a queue
    add_rect(s, Inches(7.3), Inches(1.45), Inches(5.4), Inches(5.3), LIGHT_GREY)
    add_text(s, Inches(7.3), Inches(1.6), Inches(5.4), Inches(0.4),
             "Queue Visualization", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
             
    # Draw queue lane
    add_rect(s, Inches(7.5), Inches(3.2), Inches(5.0), Inches(0.08), DARK_GREY)
    add_rect(s, Inches(7.5), Inches(4.8), Inches(5.0), Inches(0.08), DARK_GREY)
    
    # Queue items inside the lane
    q_items = ["Front (10)", "20", "30 (Rear)"]
    qx = Inches(7.8)
    for i, el in enumerate(q_items):
        col = GOLD if i == 0 else (TEAL if i == 2 else WHITE)
        add_rrect(s, qx, Inches(3.5), Inches(1.2), Inches(1.0), col, line=NAVY)
        add_text(s, qx, Inches(3.5), Inches(1.2), Inches(1.0), el,
                 size=12, bold=True, color=NAVY if col == WHITE else WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        qx += Inches(1.4)
        
    # Enqueue/Dequeue labels
    add_arrow(s, Inches(7.3), Inches(4.0), Inches(7.7), Inches(4.0))
    add_text(s, Inches(7.1), Inches(4.3), Inches(1.2), Inches(0.35),
             "← Dequeue", size=11, bold=True, color=GOLD)
             
    add_arrow(s, Inches(11.9), Inches(4.0), Inches(12.3), Inches(4.0))
    add_text(s, Inches(11.7), Inches(4.3), Inches(1.2), Inches(0.35),
             "Enqueue ←", size=11, bold=True, color=TEAL)


def slide_queue_ops(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Queue Operations & Complexity",
               "The Core Interface of a Queue", n, total)
               
    # Table of operations
    rows = [
        ("Operation", "Description", "Complexity"),
        ("enqueue(item)", "Add an item to the rear of the queue.", "O(1)"),
        ("dequeue()", "Remove and return the item at the front of the queue.", "O(1)"),
        ("peek() / front()", "Return the front item without removing it.", "O(1)"),
        ("is_empty()", "Return True if the queue has no elements, else False.", "O(1)"),
        ("size()", "Return the count of elements currently in the queue.", "O(1)"),
    ]
    
    x0 = Inches(0.9); y0 = Inches(1.6)
    col_w = [Inches(3.2), Inches(5.8), Inches(2.5)]
    rh = Inches(0.7)
    
    for r, row in enumerate(rows):
        cx = x0
        for c, txt in enumerate(row):
            fill = NAVY if r == 0 else (LIGHT_GREY if r % 2 == 1 else WHITE)
            color = WHITE if r == 0 else (GREEN if c == 2 and r > 0 else DARK_GREY)
            bold = (r == 0) or (c == 0) or (c == 2)
            add_rect(s, cx, y0 + rh * r, col_w[c], rh, fill,
                     line=RGBColor(0xC0, 0xC4, 0xCC))
            add_text(s, cx + Inches(0.15), y0 + rh * r,
                     col_w[c] - Inches(0.2), rh,
                     txt, size=15, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[c]
            
    # Key concept callout box
    callout_box_y = Inches(1.6 + 0.7 * len(rows) + 0.3)
    add_rect(s, Inches(0.9), callout_box_y, Inches(11.5), Inches(1.1), LIGHT_GREY)
    add_rect(s, Inches(0.9), callout_box_y, Inches(0.08), Inches(1.1), GOLD)
    add_text(s, Inches(1.1), callout_box_y + Inches(0.1), Inches(11.0), Inches(0.35),
             "Design Contrast:", size=14, bold=True, color=GOLD)
    add_text(s, Inches(1.1), callout_box_y + Inches(0.45), Inches(11.0), Inches(0.6),
             "Unlike stacks, queues operate on BOTH ends. This requires tracking "
             "both endpoints. If we only track one end, either enqueue or dequeue "
             "will fall to O(N) efficiency because we have to traverse or shift elements.",
             size=13, color=DARK_GREY)


def slide_queue_linked_list(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "5.  Queue Implementation: Linked List",
               "Double-pointer Singly Linked List with Front & Rear References", n, total)
               
    code = (
        "class LinkedQueue:\n"
        "    class Node:\n"
        "        def __init__(self, data, next_node=None):\n"
        "            self.data, self.next = data, next_node\n"
        "\n"
        "    def __init__(self):\n"
        "        self.front = self.rear = None\n"
        "\n"
        "    def enqueue(self, data):\n"
        "        new_node = self.Node(data)\n"
        "        if self.rear is None:\n"
        "            self.front = self.rear = new_node\n"
        "            return\n"
        "        self.rear.next = new_node  # append at tail\n"
        "        self.rear = new_node       # shift tail reference\n"
        "\n"
        "    def dequeue(self):\n"
        "        if self.front is None:\n"
        "            raise IndexError(\"Dequeue from empty queue\")\n"
        "        val = self.front.data\n"
        "        self.front = self.front.next\n"
        "        if self.front is None:     # queue now empty\n"
        "            self.rear = None\n"
        "        return val"
    )
    add_code_block(s, Inches(0.9), Inches(1.48), Inches(7.5), Inches(5.5), code, size=12)
    
    # Right notes column
    add_text(s, Inches(8.7), Inches(1.48), Inches(4.0), Inches(0.45),
             "Linked Queue Structure", size=18, bold=True, color=NAVY)
    notes = [
        ("Dual Pointers.  ", "Maintaining front and rear references allows constant-time access to both endpoints."),
        ("Enqueue at Rear.  ", "rear.next = new_node is a constant-time modification (O(1))."),
        ("Dequeue at Front.  ", "head = head.next is also O(1). No list traversal is needed."),
        ("Empty Edge Case.  ", "Careful logic is required to set rear to None if dequeue empties the queue."),
    ]
    add_bullets(s, Inches(8.7), Inches(2.0), Inches(4.0), Inches(4.8),
                notes, size=13, line_spacing=1.35, bullet_color=GOLD)


def slide_queue_circular_array(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Queue Implementation: Circular Buffer",
               "Fixed Array with modulo pointer wrapping to avoid shifting", n, total)
               
    code = (
        "class CircularQueue:\n"
        "    def __init__(self, capacity=10):\n"
        "        self._data = [None] * capacity\n"
        "        self._size = 0\n"
        "        self._front = 0\n"
        "\n"
        "    def enqueue(self, data):\n"
        "        if self._size == len(self._data):\n"
        "            raise IndexError(\"Queue is full\")\n"
        "        rear_idx = (self._front + self._size) % len(self._data)\n"
        "        self._data[rear_idx] = data\n"
        "        self._size += 1\n"
        "\n"
        "    def dequeue(self):\n"
        "        if self._size == 0:\n"
        "            raise IndexError(\"Queue is empty\")\n"
        "        val = self._data[self._front]\n"
        "        self._data[self._front] = None\n"
        "        self._front = (self._front + 1) % len(self._data)\n"
        "        self._size -= 1\n"
        "        return val"
    )
    add_code_block(s, Inches(0.9), Inches(1.48), Inches(7.5), Inches(5.5), code, size=12)
    
    # Right notes column
    add_text(s, Inches(8.7), Inches(1.48), Inches(4.0), Inches(0.45),
             "Why Circular Buffer?", size=18, bold=True, color=NAVY)
    notes = [
        ("The Shifting Problem.  ", "If we dequeue from index 0 of a normal array, we must shift the remaining elements, costing O(N)."),
        ("Circular Pointer.  ", "Instead of shifting data, we shift the FRONT pointer clockwise using modular math: (front + 1) % capacity."),
        ("Boundary pointer checks.  ", "Empty queue state: front == rear. Full queue state: (rear + 1) % capacity == front (leaving one empty slot)."),
        ("No Shifting (O(1)).  ", "Both operations are now guaranteed O(1) time without element copies."),
    ]
    add_bullets(s, Inches(8.7), Inches(2.0), Inches(4.0), Inches(4.8),
                notes, size=13, line_spacing=1.35, bullet_color=TEAL)


def slide_queue_circular_trace(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Circular Buffer: Step-by-Step Pointer Trace",
               "Visualizing Front and Rear pointer wrapping with a capacity of 4", n, total)

    phases = [
        (
            "1. Initial State", 
            ["A", "B", "C", "D"], 
            0, 
            3, 
            [
                ("Full Queue. ", "All 4 slots are occupied."),
                ("Front points to A ", "at index 0."),
                ("Rear points to D ", "at index 3."),
            ]
        ),
        (
            "2. Dequeue A & B", 
            ["", "", "C", "D"], 
            2, 
            3, 
            [
                ("No shifting. ", "We do not shift elements."),
                ("Front moves ", "from 0 to 1, then to 2."),
                ("Empty space ", "created at slots 0 & 1."),
            ]
        ),
        (
            "3. Enqueue E (Wrap)", 
            ["E", "", "C", "D"], 
            2, 
            0, 
            [
                ("Space Reused. ", "We place E at index 0."),
                ("Rear wraps: ", "(3 + 1) % 4 = 0."),
                ("Queue is now ", "divided across boundaries."),
            ]
        ),
        (
            "4. Dequeue C & D", 
            ["E", "", "", ""], 
            0, 
            0, 
            [
                ("Front wraps. ", "C at 2, then D at 3 are removed."),
                ("Modulo wrapping: ", "(3 + 1) % 4 = 0."),
                ("Front = Rear = 0 ", "pointing to E."),
            ]
        ),
    ]

    for col_idx, (title, cells, front_idx, rear_idx, desc) in enumerate(phases):
        X_col = Inches(0.9 + col_idx * 2.95)
        Y_col = Inches(1.5)
        W_col = Inches(2.7)
        H_col = Inches(5.1)
        
        # Panel Background
        add_rrect(s, X_col, Y_col, W_col, H_col, LIGHT_GREY)
        
        # Panel Header
        add_text(s, X_col, Y_col + Inches(0.15), W_col, Inches(0.45), title, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        
        # Draw 4 Cells
        Y_cells = Y_col + Inches(0.7)
        for i in range(4):
            cell_X = X_col + Inches(0.375) + i * Inches(0.5)
            
            # Highlight/Color Logic
            if i == front_idx and i == rear_idx:
                fill = PURPLE
                text_color = WHITE
            elif i == front_idx:
                fill = TEAL
                text_color = WHITE
            elif i == rear_idx:
                fill = GOLD
                text_color = WHITE
            elif cells[i] == "":
                fill = WHITE
                text_color = DARK_GREY
            else:
                fill = WHITE
                text_color = DARK_GREY
                
            # Draw Cell Box
            add_rect(s, cell_X, Y_cells, Inches(0.45), Inches(0.45), fill, line=DARK_GREY, line_width=1)
            
            # Cell Value
            val_str = cells[i] if cells[i] != "" else " "
            add_text(s, cell_X, Y_cells + Inches(0.08), Inches(0.45), Inches(0.35), val_str, size=14, bold=True, color=text_color, align=PP_ALIGN.CENTER)
            
            # Index under the Cell
            add_text(s, cell_X, Y_cells + Inches(0.48), Inches(0.45), Inches(0.25), str(i), size=9, color=DARK_GREY, align=PP_ALIGN.CENTER)
            
        # Draw Status
        Y_status = Y_col + Inches(1.6)
        status_text = f"front: index {front_idx}\nrear: index {rear_idx}"
        add_text(s, X_col, Y_status, W_col, Inches(0.6), status_text, size=12, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)
        
        # Draw Description
        Y_desc = Y_col + Inches(2.25)
        add_bullets(s, X_col + Inches(0.12), Y_desc, W_col - Inches(0.24), Inches(2.7), desc, size=11, line_spacing=1.2, bullet_color=TEAL)


def slide_deque(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "6.  Double-Ended Queue (Deque)",
               "Linear structure supporting insertions and deletions at both ends", n, total)
               
    # Two columns: visual deque vs operations
    add_text(s, Inches(0.9), Inches(1.45), Inches(5.8), Inches(0.5),
             "Generalizing Stacks and Queues", size=20, bold=True, color=NAVY)
             
    items = [
        ("Combination of ADTs.  ", "Supports push/pop at the front, and push/pop at the rear."),
        ("Versatile.  ", "Can act as a Stack, a Queue, or both simultaneously."),
        ("Python standard.  ", "Implemented natively as collections.deque in Python, optimized for fast O(1) appends and pops at both ends."),
    ]
    add_bullets(s, Inches(0.9), Inches(2.0), Inches(5.8), Inches(2.8),
                items, size=15, line_spacing=1.35)
                
    # Python code example
    code = (
        "from collections import deque\n"
        "\n"
        "dq = deque()          # create deque\n"
        "dq.append(20)         # rear append: [20]\n"
        "dq.appendleft(10)     # front append: [10, 20]\n"
        "dq.append(30)         # rear append: [10, 20, 30]\n"
        "\n"
        "dq.popleft()          # front remove: returns 10\n"
        "dq.pop()              # rear remove: returns 30"
    )
    add_code_block(s, Inches(0.9), Inches(4.9), Inches(5.8), Inches(2.1), code, size=13)
    
    # Right panel: Deque visualization
    add_rect(s, Inches(7.1), Inches(1.45), Inches(5.6), Inches(5.5), LIGHT_GREY)
    add_text(s, Inches(7.1), Inches(1.7), Inches(5.6), Inches(0.4),
             "Deque Layout", size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
             
    # Draw deque lane
    add_rect(s, Inches(7.3), Inches(3.2), Inches(5.2), Inches(0.08), DARK_GREY)
    add_rect(s, Inches(7.3), Inches(4.8), Inches(5.2), Inches(0.08), DARK_GREY)
    
    # Items
    dq_items = ["10", "20", "30"]
    dq_x = Inches(7.9)
    for el in dq_items:
        add_rrect(s, dq_x, Inches(3.5), Inches(1.0), Inches(1.0), WHITE, line=NAVY)
        add_text(s, dq_x, Inches(3.5), Inches(1.0), Inches(1.0), el,
                 size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        dq_x += Inches(1.2)
        
    # Arrows on both sides
    add_arrow(s, Inches(7.1), Inches(3.7), Inches(7.7), Inches(3.7))
    add_arrow(s, Inches(7.7), Inches(4.3), Inches(7.1), Inches(4.3), color=GOLD)
    add_text(s, Inches(6.9), Inches(4.6), Inches(1.4), Inches(0.35),
             "Front appends/pops", size=10, italic=True, color=NAVY)
             
    add_arrow(s, Inches(12.7), Inches(3.7), Inches(12.1), Inches(3.7))
    add_arrow(s, Inches(12.1), Inches(4.3), Inches(12.7), Inches(4.3), color=GOLD)
    add_text(s, Inches(11.5), Inches(4.6), Inches(1.4), Inches(0.35),
             "Rear appends/pops", size=10, italic=True, color=NAVY)


def slide_applications(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "7.  Real-World Applications",
               "How Stacks and Queues power critical computer systems", n, total)
               
    apps = [
        ("The Program Call Stack", NAVY,
         "Manages function execution frames, recursion stack space (costing O(N) auxiliary space), and local variables in runtime environments."),
        ("Undo/Redo History", TEAL,
         "Text editors use two stacks (Undo stack and Redo stack) to track document modifications in reverse order."),
        ("Breadth-First Search (BFS)", GOLD,
         "Graph traversals use a queue to track visited nodes in a level-order fashion (shortest path, web crawlers)."),
        ("Asynchronous Buffers", RED,
         "Data sent between asymmetric processes (network socket packets, printer spoolers, keyboard event loops)."),
        ("Round-Robin Scheduling", GREEN,
         "OS CPU schedulers cycle executing processes by rotating a circular queue of threads."),
        ("Backtracking (DFS)", PURPLE,
         "DFS uses a stack (explicit or recursion call stack) to backtrack. Dijkstra's Shunting-yard uses an operator stack for math parsing."),
    ]
    
    cw = Inches(4.0); ch = Inches(2.4); gx = Inches(0.25); gy = Inches(0.2)
    x0 = Inches(0.6); y0 = Inches(1.55)
    for i, (title, color, body) in enumerate(apps):
        col = i % 3; row = i // 3
        cx = x0 + (cw + gx) * col
        cy = y0 + (ch + gy) * row
        add_rect(s, cx, cy, cw, Inches(0.55), color)
        add_text(s, cx, cy, cw, Inches(0.55), title,
                 size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, cx, cy + Inches(0.55), cw, ch - Inches(0.55), LIGHT_GREY)
        add_text(s, cx + Inches(0.18), cy + Inches(0.72),
                 cw - Inches(0.35), ch - Inches(0.8),
                 body, size=12, color=DARK_GREY)


def slide_practice(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "8.  Practice Problems",
               "Crucial algorithmic challenges to build your mastery", n, total)
               
    problems = [
        ("Valid Parentheses", TEAL,
         "LeetCode 20 (Easy)\n"
         "Given a string of brackets '{}()[]', check if closed in correct order.\n"
         "Use a Stack to match open-to-close brackets."),
        ("Implement Queue using Stacks", GOLD,
         "LeetCode 232 (Medium)\n"
         "FIFO queue using Stack1 (enqueue) and Stack2 (dequeue). Transfer elements if empty.\n"
         "Amortized O(1) dequeue: elements are pushed/popped a constant number of times."),
        ("Implement Stack using Queues", NAVY,
         "LeetCode 225 (Easy)\n"
         "Build a LIFO stack API using only FIFO queues.\n"
         "Rotate elements within a single queue on push/pop."),
        ("Min Stack", GREEN,
         "LeetCode 155 (Medium)\n"
         "Design a stack that supports get_min() in O(1) time.\n"
         "Uses a parallel auxiliary stack to track the minimum values."),
        ("Evaluate Reverse Polish Notation", PURPLE,
         "LeetCode 150 (Medium)\n"
         "Evaluate math strings like ['2', '1', '+', '3', '*'] -> 9.\n"
         "Push operands; pop for operations. (Related: Dijkstra's Shunting-yard)."),
        ("Next Greater Element", RED,
         "LeetCode 496 / 739 (Medium)\n"
         "Find the next larger element for each array index in O(N) time.\n"
         "Solved using a Monotonic Stack to store elements in decreasing order."),
    ]
    
    cw = Inches(4.0); ch = Inches(2.35); gx = Inches(0.25); gy = Inches(0.2)
    x0 = Inches(0.6); y0 = Inches(1.55)
    for i, (title, color, body) in enumerate(problems):
        col = i % 3; row = i // 3
        cx = x0 + (cw + gx) * col
        cy = y0 + (ch + gy) * row
        add_rect(s, cx, cy, cw, Inches(0.52), color)
        add_text(s, cx, cy, cw, Inches(0.52), title,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, cx, cy + Inches(0.52), cw, ch - Inches(0.52), LIGHT_GREY)
        add_text(s, cx + Inches(0.18), cy + Inches(0.68),
                 cw - Inches(0.35), ch - Inches(0.75),
                 body, size=12, color=DARK_GREY)


def slide_summary(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Summary — Core ADT Concepts",
               "Key takeaways to review before leaving", n, total)
               
    items = [
        "Abstract Data Types define the behaviors (API) and operations, decoupling from physical memory representations.",
        "A Stack is LIFO (Last-In-First-Out). Operations push, pop, and peek occur at the top of the stack.",
        "A Queue is FIFO (First-In-First-Out). Elements are added at the rear and removed from the front.",
        "Stack implements O(1) easily using a Singly Linked List (at head) or dynamic array (at end).",
        "Queue requires dual pointers (front and rear) in a Linked List for O(1) performance.",
        "Queue in arrays is implemented as a Circular Buffer to avoid O(N) element shifting.",
    ]
    add_bullets(s, Inches(0.9), Inches(1.55), Inches(12.1), Inches(5.55),
                items, size=18, line_spacing=1.45, bullet_color=GOLD)


def slide_references(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "References & Further Reading",
               "Resources for deeper study", n, total)
               
    refs = [
        ("Knuth, D. E.  ", "The Art of Computer Programming, Vol. 1. Fundamental Algorithms, Section 2.2. Linear Lists (Stacks and Queues)."),
        ("CLRS 4th Ed.  ", "Cormen et al. Introduction to Algorithms. Chapter 10: Elementary Data Structures (Stacks and Queues)."),
        ("Python Collections Documentation.  ", "Official reference for collections.deque, dynamic lists, and optimization mechanics."),
        ("LeetCode Challenges.  ", "Problems: 20 (Parentheses), 155 (Min Stack), 225 & 232 (Cross Implementations), 239 (Monotonic Deque)."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.55), Inches(12.1), Inches(5.6),
                refs, size=16, line_spacing=1.4)


def slide_thanks():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(5.6), SLIDE_W, Inches(0.14), GOLD)
    add_rect(s, 0, Inches(5.8), SLIDE_W, Inches(0.05), TEAL)
    
    add_text(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.5),
             "Questions?", size=80, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.2), Inches(11.5), Inches(0.7),
             "Bring your traces, your arrays, and your pointer puzzles.",
             size=22, italic=True, color=GOLD)
    add_text(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.5),
             "Remember:  🥞  stacks push/pop at the top; 🚶 queues enter rear, leave front.",
             size=16, color=LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.4),
             "Next topic  ·  Trees and Graphs — hierarchical and networked relations.",
             size=15, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.4),
             "Data Structures & Algorithms  ·  Week 03 Supplement  ·  Stacks & Queues",
             size=11, color=LIGHT_GREY, italic=True)


# ---------- Build Sequence ----------

content_slides = [
    slide_outline,
    slide_adt_concept,
    slide_stack_concept,
    slide_stack_ops,
    slide_stack_linked_list,
    slide_stack_array,
    slide_queue_concept,
    slide_queue_ops,
    slide_queue_linked_list,
    slide_queue_circular_array,
    slide_queue_circular_trace,
    slide_deque,
    slide_applications,
    slide_practice,
    slide_summary,
    slide_references,
]

TOTAL = 1 + len(content_slides) + 1  # Title + Content + Thanks

slide_title()
for idx, fn in enumerate(content_slides, start=2):
    fn(idx, TOTAL)
slide_thanks()

# Save PPTX
OUT = os.path.join(os.path.dirname(__file__), "stacks_queues_dsa.pptx")
prs.save(OUT)
print(f"✓ Wrote {OUT} ({len(prs.slides)} slides)")
