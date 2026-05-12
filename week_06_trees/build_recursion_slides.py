"""
Generate a university-grade PowerPoint deck for Week 06: Recursion Algorithms.

Output: recursion_algorithms.pptx (16:9 widescreen)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# -------- Theme --------
NAVY        = RGBColor(0x0B, 0x2A, 0x4A)
TEAL        = RGBColor(0x12, 0x7A, 0x8A)
GOLD        = RGBColor(0xE0, 0xA8, 0x1F)
LIGHT_GREY  = RGBColor(0xF2, 0xF4, 0xF7)
DARK_GREY   = RGBColor(0x33, 0x3D, 0x4A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xC0, 0x39, 0x2B)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
PURPLE      = RGBColor(0x6A, 0x1B, 0x9A)
CODE_BG     = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG     = RGBColor(0xE6, 0xE6, 0xE6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------- Generic helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=18, color=DARK_GREY,
                bullet_color=TEAL, line_spacing=1.25, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        rb = p.add_run(); rb.text = "▸  "
        rb.font.name = font; rb.font.size = Pt(size)
        rb.font.bold = True; rb.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run(); r1.text = head
            r1.font.name = font; r1.font.size = Pt(size)
            r1.font.bold = True; r1.font.color.rgb = NAVY
            r2 = p.add_run(); r2.text = tail
            r2.font.name = font; r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = item
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None, page_num=None, total=None):
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    add_rect(slide, Inches(0.35), 0, SLIDE_W - Inches(0.35), Inches(0.08), GOLD)
    add_text(slide, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.78), Inches(11.5), Inches(0.45),
                 subtitle, size=15, italic=True, color=TEAL)
    add_rect(slide, Inches(0.6), Inches(1.22), Inches(2.0), Inches(0.04), GOLD)
    add_text(slide, Inches(0.6), SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
             "Data Structures & Algorithms · Recursion Algorithms",
             size=10, color=DARK_GREY, italic=True)
    if page_num is not None and total is not None:
        add_text(slide, SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.45),
                 Inches(1.2), Inches(0.3),
                 f"Slide {page_num} / {total}", size=10, color=DARK_GREY,
                 align=PP_ALIGN.RIGHT)


def add_code_block(slide, x, y, w, h, code, size=14):
    add_rect(slide, x, y, w, h, CODE_BG)
    add_rect(slide, x, y, Inches(0.08), h, GOLD)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.25), h - Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.line_spacing = 1.15
        r = p.add_run(); r.text = ln if ln else " "
        r.font.name = "Consolas"; r.font.size = Pt(size)
        r.font.color.rgb = CODE_FG

# ---------- Slides ----------
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(5.5), SLIDE_W, Inches(0.12), GOLD)
    add_rect(s, 0, Inches(5.7), SLIDE_W, Inches(0.04), TEAL)
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.9), Inches(1.0),
                              Inches(3.4), Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb = GOLD
    chip.line.fill.background(); chip.shadow.inherit = False
    add_text(s, Inches(0.9), Inches(1.0), Inches(3.4), Inches(0.55),
             "PRE-REQUISITE", size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(1.4),
             "Recursive Algorithms", size=58, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.7),
             "The Foundation for Tree Traversals and Divide & Conquer",
             size=24, color=GOLD, italic=True)
    add_text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.6),
             "A University-Grade Lecture in Data Structures & Algorithms",
             size=18, color=LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(6.0), Inches(8), Inches(0.4),
             "Course   ·   CSC: Data Structures & Algorithms",
             size=14, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.4), Inches(8), Inches(0.4),
             "Prerequisite for Trees & Divide/Conquer",
             size=14, color=WHITE)

def slide_outline(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Lecture Outline",
               "Understanding functions that call themselves", n, total)
    items = [
        ("1. What is Recursion?  ", "Definition and real-world intuition."),
        ("2. The Two Rules.  ", "Base Case and Recursive Step."),
        ("3. The Call Stack.  ", "How the computer keeps track of deferred operations."),
        ("4. Example: Factorial.  ", "A classic math example broken down line-by-line."),
        ("5. Example: Fibonacci.  ", "Branching recursion and why overlapping subproblems matter."),
        ("6. Recursion vs Iteration.  ", "When to use which: memory overhead vs elegance."),
        ("7. Why Recursion for Trees?  ", "How hierarchical data structures naturally demand recursive logic."),
        ("8. Common Pitfalls.  ", "Stack Overflow and missing base cases.")
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(12), Inches(5.5),
                items, size=18, line_spacing=1.32)

def slide_what_is_recursion(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "1. What is Recursion?",
               "Solving a problem by solving smaller instances of the same problem", n, total)
    add_text(s, Inches(0.9), Inches(1.5), Inches(6), Inches(0.5),
             "Definition", size=20, bold=True, color=NAVY)
    add_text(s, Inches(0.9), Inches(2.05), Inches(6), Inches(1.5),
             "In computer science, recursion occurs when a function calls itself directly or indirectly. "
             "It allows us to break down complex problems into smaller, manageable subproblems of the exact same structure.",
             size=15, color=DARK_GREY)
    
    add_text(s, Inches(0.9), Inches(3.8), Inches(6), Inches(0.5),
             "Real-world Intuition", size=18, bold=True, color=TEAL)
    items = [
        "A dictionary definition that refers to other words in the dictionary.",
        "Russian Matryoshka dolls (a doll inside a doll).",
        "Looking between two facing mirrors."
    ]
    add_bullets(s, Inches(0.9), Inches(4.3), Inches(6.0), Inches(2.4),
                items, size=15, line_spacing=1.3)
    
    # Right side code visualization
    add_rect(s, Inches(7.5), Inches(1.5), Inches(5.4), Inches(5.3), LIGHT_GREY)
    add_text(s, Inches(7.5), Inches(1.55), Inches(5.4), Inches(0.4),
             "The Structure of a Recursive Function", size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    
    code = (
        "def recursive_function(parameters):\n"
        "    if base_condition:\n"
        "        # BASE CASE\n"
        "        return simple_result\n"
        "    else:\n"
        "        # RECURSIVE STEP\n"
        "        modified_params = change(parameters)\n"
        "        return recursive_function(modified_params)"
    )
    add_code_block(s, Inches(7.7), Inches(2.2), Inches(5.0), Inches(3.0), code, size=14)

def slide_two_rules(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "2. The Two Essential Rules",
               "Every valid recursive algorithm must satisfy these two conditions", n, total)
    
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
             "Rule #1: The Base Case", size=22, bold=True, color=RED)
    add_text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.0),
             "There must be at least one condition where the function stops calling itself and returns a value. "
             "Without a base case, the function will recurse infinitely, leading to a Stack Overflow.",
             size=16, color=DARK_GREY)
             
    add_text(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.5),
             "Rule #2: The Recursive Step (making progress)", size=22, bold=True, color=NAVY)
    add_text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(1.0),
             "Each time the function calls itself, it must alter the parameters so that they move closer to the base case. "
             "If the parameters do not change, or move away from the base case, the recursion will never terminate.",
             size=16, color=DARK_GREY)

def slide_call_stack(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "3. The Call Stack",
               "How the computer remembers where it was", n, total)
    
    add_text(s, Inches(0.9), Inches(1.5), Inches(6.5), Inches(5.0),
             "When a function calls itself, the computer pauses the current execution and pushes the context onto the Call Stack.\n\n"
             "• A stack operates Last-In, First-Out (LIFO).\n"
             "• Each call gets its own memory frame (variables and state).\n"
             "• When the base case is reached, the stack starts popping off, returning values back up the chain.\n\n"
             "This explains why recursion can consume O(N) memory even if there are no explicit data structures.",
             size=16, color=DARK_GREY)
             
    # Right side: stack visualization
    add_rect(s, Inches(8.0), Inches(1.5), Inches(4.5), Inches(5.0), LIGHT_GREY)
    add_text(s, Inches(8.0), Inches(1.6), Inches(4.5), Inches(0.4),
             "Call Stack Example", size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
             
    add_rect(s, Inches(8.5), Inches(4.5), Inches(3.5), Inches(0.6), GOLD)
    add_text(s, Inches(8.5), Inches(4.5), Inches(3.5), Inches(0.6), "func(1) -> BASE CASE", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    
    add_rect(s, Inches(8.5), Inches(3.8), Inches(3.5), Inches(0.6), TEAL)
    add_text(s, Inches(8.5), Inches(3.8), Inches(3.5), Inches(0.6), "func(2) -> waiting...", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, color=WHITE)
    
    add_rect(s, Inches(8.5), Inches(3.1), Inches(3.5), Inches(0.6), NAVY)
    add_text(s, Inches(8.5), Inches(3.1), Inches(3.5), Inches(0.6), "func(3) -> waiting...", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, color=WHITE)

def slide_factorial(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "4. Classic Example: Factorial",
               "N! = N * (N-1)!", n, total)
               
    code = (
        "def factorial(n):\n"
        "    # Base Case\n"
        "    if n == 1 or n == 0:\n"
        "        return 1\n"
        "        \n"
        "    # Recursive Step\n"
        "    return n * factorial(n - 1)"
    )
    add_code_block(s, Inches(0.9), Inches(1.5), Inches(5.5), Inches(3.5), code, size=15)
    
    add_text(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5), "Execution of factorial(4):", bold=True, color=NAVY)
    
    trace = [
        "factorial(4) returns 4 * factorial(3)",
        "factorial(3) returns 3 * factorial(2)",
        "factorial(2) returns 2 * factorial(1)",
        "factorial(1) returns 1 (Base Case!)"
    ]
    add_bullets(s, Inches(7.0), Inches(2.0), Inches(6.0), Inches(2.0), trace, size=14)
    
    add_text(s, Inches(7.0), Inches(4.2), Inches(6.0), Inches(2.0), 
             "Unwinding the stack:\n"
             "• 2 * 1 = 2\n"
             "• 3 * 2 = 6\n"
             "• 4 * 6 = 24  → Final Result", size=15, bold=True, color=GREEN)

def slide_fibonacci(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "5. Branching Recursion: Fibonacci",
               "Multiple recursive calls per function", n, total)
               
    code = (
        "def fib(n):\n"
        "    if n <= 1:\n"
        "        return n\n"
        "    return fib(n-1) + fib(n-2)"
    )
    add_code_block(s, Inches(0.9), Inches(1.5), Inches(5.0), Inches(2.5), code, size=15)
    
    add_text(s, Inches(0.9), Inches(4.2), Inches(5.0), Inches(2.0),
             "Time Complexity: O(2^N)\n"
             "Space Complexity: O(N) due to call stack depth.", size=14, bold=True, color=RED)
             
    # Tree visualization on the right
    add_text(s, Inches(6.5), Inches(1.5), Inches(6.0), Inches(0.5), "Call Tree for fib(4):", bold=True, color=NAVY)
    add_rect(s, Inches(6.5), Inches(2.0), Inches(6.0), Inches(4.5), LIGHT_GREY)
    add_text(s, Inches(6.5), Inches(2.5), Inches(6.0), Inches(3.5), 
             "            fib(4)\n"
             "           /      \\\n"
             "      fib(3)      fib(2)\n"
             "      /    \\      /    \\\n"
             "  fib(2) fib(1) fib(1) fib(0)\n"
             "  /   \\\n"
             "fib(1) fib(0)", size=16, font="Consolas", align=PP_ALIGN.CENTER)

def slide_vs_iteration(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "6. Recursion vs Iteration",
               "When to use which?", n, total)
               
    add_rect(s, Inches(0.9), Inches(1.5), Inches(5.5), Inches(4.5), LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(1.6), Inches(5.5), Inches(0.5), "Iterative (Loops)", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    iter_bullets = [
        "Uses explicit loops (for, while).",
        "O(1) auxiliary space (usually).",
        "Faster execution (no function call overhead).",
        "Can be harder to read for complex data structures like trees."
    ]
    add_bullets(s, Inches(1.1), Inches(2.2), Inches(5.1), Inches(3.5), iter_bullets, size=15)
    
    add_rect(s, Inches(6.9), Inches(1.5), Inches(5.5), Inches(4.5), LIGHT_GREY)
    add_text(s, Inches(6.9), Inches(1.6), Inches(5.5), Inches(0.5), "Recursive", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rec_bullets = [
        "Calls itself.",
        "O(N) auxiliary space for call stack.",
        "Slower due to stack setup/teardown.",
        "Leads to extremely elegant, concise code for hierarchical data."
    ]
    add_bullets(s, Inches(7.1), Inches(2.2), Inches(5.1), Inches(3.5), rec_bullets, size=15)

def slide_why_trees(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "7. Why Recursion for Trees?",
               "The shape of the data matches the shape of the algorithm", n, total)
               
    add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.0),
             "A tree is inherently recursive by definition:\n"
             "A Tree is either empty, or a Root node connected to left and right Subtrees.",
             size=20, bold=True, color=NAVY)
             
    code = (
        "# Traversing a tree recursively is beautifully simple:\n"
        "def dfs(node):\n"
        "    if node is None:         # Base Case\n"
        "        return\n"
        "    \n"
        "    print(node.val)          # Process current\n"
        "    dfs(node.left)           # Recurse left\n"
        "    dfs(node.right)          # Recurse right\n"
    )
    add_code_block(s, Inches(0.9), Inches(3.0), Inches(7.5), Inches(3.5), code, size=16)
    
    add_text(s, Inches(8.8), Inches(3.5), Inches(3.5), Inches(2.0),
             "Attempting to do this iteratively requires managing an explicit Stack manually, which complicates the code significantly.", size=16, color=DARK_GREY)

def slide_summary(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Summary & Pitfalls",
               "Key takeaways before we jump into Trees", n, total)
               
    items = [
        ("Base Cases are mandatory. ", "Without them, you get a RecursionError (Stack Overflow)."),
        ("Trust the recursion. ", "Assume the recursive call will correctly solve the subproblem, and use its result."),
        ("Watch out for overlapping subproblems. ", "Like in Fibonacci, standard recursion can be exponentially slow unless optimized (e.g., Memoization)."),
        ("Recursion is the key to Trees. ", "You must be comfortable with recursive thinking to master Binary Trees and BSTs.")
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.0), items, size=18, line_spacing=1.5)

# --- Build ---
slides = [
    slide_title,
    slide_outline,
    slide_what_is_recursion,
    slide_two_rules,
    slide_call_stack,
    slide_factorial,
    slide_fibonacci,
    slide_vs_iteration,
    slide_why_trees,
    slide_summary
]

for i, f in enumerate(slides):
    if f == slide_title:
        f()
    else:
        f(i, len(slides)-1)

prs.save("recursion_algorithms.pptx")
print("Saved recursion_algorithms.pptx")
