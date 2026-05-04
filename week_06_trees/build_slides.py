"""
Generate a university-grade PowerPoint deck for Week 06: Trees & BSTs.

Output: trees_dsa_final.pptx (16:9 widescreen)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# -------- Theme (matches Week 03 deck) --------
NAVY        = RGBColor(0x0B, 0x2A, 0x4A)
TEAL        = RGBColor(0x12, 0x7A, 0x8A)
GOLD        = RGBColor(0xE0, 0xA8, 0x1F)
LIGHT_GREY  = RGBColor(0xF2, 0xF4, 0xF7)
DARK_GREY   = RGBColor(0x33, 0x3D, 0x4A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
RED         = RGBColor(0xC0, 0x39, 0x2B)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
PURPLE      = RGBColor(0x6A, 0x1B, 0x9A)
CODE_BG     = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG     = RGBColor(0xE6, 0xE6, 0xE6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------- Generic helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=18, color=DARK_GREY,
                bullet_color=TEAL, line_spacing=1.25, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        rb = p.add_run(); rb.text = "▸  "
        rb.font.name = font; rb.font.size = Pt(size)
        rb.font.bold = True; rb.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run(); r1.text = head
            r1.font.name = font; r1.font.size = Pt(size)
            r1.font.bold = True; r1.font.color.rgb = NAVY
            r2 = p.add_run(); r2.text = tail
            r2.font.name = font; r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = item
            r.font.name = font; r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None, page_num=None, total=None):
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, NAVY)
    add_rect(slide, Inches(0.35), 0, SLIDE_W - Inches(0.35), Inches(0.08), GOLD)
    add_text(slide, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.78), Inches(11.5), Inches(0.45),
                 subtitle, size=15, italic=True, color=TEAL)
    add_rect(slide, Inches(0.6), Inches(1.22), Inches(2.0), Inches(0.04), GOLD)
    add_text(slide, Inches(0.6), SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
             "Data Structures & Algorithms · Week 06 · Trees & BSTs",
             size=10, color=DARK_GREY, italic=True)
    if page_num is not None and total is not None:
        add_text(slide, SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.45),
                 Inches(1.2), Inches(0.3),
                 f"Slide {page_num} / {total}", size=10, color=DARK_GREY,
                 align=PP_ALIGN.RIGHT)


def add_code_block(slide, x, y, w, h, code, size=14):
    add_rect(slide, x, y, w, h, CODE_BG)
    add_rect(slide, x, y, Inches(0.08), h, GOLD)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.25), h - Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.line_spacing = 1.15
        r = p.add_run(); r.text = ln if ln else " "
        r.font.name = "Consolas"; r.font.size = Pt(size)
        r.font.color.rgb = CODE_FG


def add_arrow(slide, x1, y1, x2, y2, color=TEAL, weight=2.0, head=True):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if head:
        ln = conn.line._get_or_add_ln()
        etree.SubElement(ln, qn("a:tailEnd"),
                         {"type": "triangle", "w": "med", "len": "med"})
    return conn


# ---------- Tree-drawing helpers ----------
def add_tree_node(slide, cx, cy, value, *, r=Inches(0.4), fill=WHITE,
                  border=NAVY, text_color=NAVY, size=16, highlight=False):
    """Draw a circular tree node centered at (cx, cy)."""
    if highlight:
        fill = GOLD
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 cx - r, cy - r, r * 2, r * 2)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(2)
    shp.shadow.inherit = False
    add_text(slide, cx - r, cy - r, r * 2, r * 2,
             str(value), size=size, bold=True, color=text_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_edge(slide, x1, y1, x2, y2, color=DARK_GREY, weight=1.75):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    return conn


def draw_bst(slide, root_cx, root_cy, layout, *, level_h=Inches(0.95),
             node_r=Inches(0.35), highlight_path=None, font_size=14,
             leaf_color=None):
    """
    Draw a small BST.
    layout: dict mapping value -> (col_offset_in_units, level)
    where col_offset is float column index from root, level is depth (root=0).
    Edges are inferred from BST structure of the values themselves.
    Simpler: pass edges explicitly via `edges` param.
    """
    pass  # placeholder; we'll use draw_tree below


def draw_tree(slide, positions, edges, *, node_r=Inches(0.35),
              fill=WHITE, border=NAVY, text_color=NAVY,
              highlight_nodes=None, highlight_edges=None,
              font_size=14, edge_color=DARK_GREY):
    """
    positions: dict node_label -> (cx, cy) in EMU/Inches
    edges: list of (parent_label, child_label)
    """
    highlight_nodes = highlight_nodes or set()
    highlight_edges = highlight_edges or set()
    # edges first so nodes overlay
    for (a, b) in edges:
        x1, y1 = positions[a]; x2, y2 = positions[b]
        col = GOLD if (a, b) in highlight_edges else edge_color
        weight = 2.5 if (a, b) in highlight_edges else 1.75
        add_edge(slide, x1, y1 + node_r, x2, y2 - node_r,
                 color=col, weight=weight)
    for label, (cx, cy) in positions.items():
        hl = label in highlight_nodes
        add_tree_node(slide, cx, cy, label, r=node_r,
                      fill=fill, border=border,
                      text_color=text_color, size=font_size,
                      highlight=hl)


# ---------- Slides ----------
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(5.5), SLIDE_W, Inches(0.12), GOLD)
    add_rect(s, 0, Inches(5.7), SLIDE_W, Inches(0.04), TEAL)
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.9), Inches(1.0),
                              Inches(2.4), Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb = GOLD
    chip.line.fill.background(); chip.shadow.inherit = False
    add_text(s, Inches(0.9), Inches(1.0), Inches(2.4), Inches(0.55),
             "WEEK 06", size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(1.4),
             "Trees & Binary Search Trees", size=58, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.7),
             "Hierarchies, Recursion, and Logarithmic Search",
             size=24, color=GOLD, italic=True)
    add_text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.6),
             "A University-Grade Lecture in Data Structures & Algorithms",
             size=18, color=LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(6.0), Inches(8), Inches(0.4),
             "Course   ·   CSC: Data Structures & Algorithms",
             size=14, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.4), Inches(8), Inches(0.4),
             "Module   ·   Week 06 of 12",
             size=14, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.8), Inches(8), Inches(0.4),
             "Format   ·   Lecture · Demonstration · Practical Lab",
             size=14, color=WHITE)


def slide_outline(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Lecture Outline",
               "From flat lists to branching hierarchies", n, total)
    items = [
        ("1. From linear to hierarchical.  ", "Why we need trees."),
        ("2. Tree terminology.  ", "Root, parent, child, leaf, depth, height."),
        ("3. Binary Trees.  ", "At most two children — and what that buys us."),
        ("4. Binary Search Trees.  ", "The ordering invariant that enables O(log N) lookup."),
        ("5. Insertion & Search.  ", "Walking the tree under the BST property."),
        ("6. DFS Traversals.  ", "Pre-order, In-order, Post-order — and what each is for."),
        ("7. BFS / Level-order.  ", "Visiting tier by tier with a queue."),
        ("8. Lowest Common Ancestor.  ", "A canonical interview question."),
        ("9. Balanced vs Degenerate.  ", "Why a sorted insert ruins your tree."),
        ("10. Pitfalls, applications, and the lab."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(12), Inches(5.5),
                items, size=18, line_spacing=1.32)


def slide_objectives(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Learning Outcomes",
               "By the end of this lecture you will be able to…", n, total)
    objs = [
        "Define the canonical tree vocabulary and identify each component on a diagram.",
        "Implement a TreeNode and a Binary Search Tree with iterative insertion.",
        "State the BST invariant and explain why it implies O(log N) search on balanced inputs.",
        "Implement Pre-, In-, and Post-order DFS traversals recursively and reason about which to use when.",
        "Describe Level-order (BFS) traversal and implement it with a queue.",
        "Use the BST property to compute the Lowest Common Ancestor of two nodes in O(h).",
        "Recognise degenerate trees and describe the motivation for self-balancing variants (AVL, Red-Black).",
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(12), Inches(5.5),
                objs, size=17, line_spacing=1.32)


def slide_motivation(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "1.  From Linear to Hierarchical",
               "Where a single chain stops being enough", n, total)
    add_text(s, Inches(0.9), Inches(1.5), Inches(6), Inches(0.5),
             "Linked lists are 1-D chains.", size=20, bold=True, color=NAVY)
    add_text(s, Inches(0.9), Inches(2.05), Inches(6), Inches(2.5),
             "But the world is hierarchical: a filesystem is folders inside "
             "folders; an HTML page is elements inside elements; a chess "
             "engine reasons about moves that branch into more moves.\n\n"
             "Whenever an item logically owns several other items, you have "
             "a tree.",
             size=15, color=DARK_GREY)
    add_text(s, Inches(0.9), Inches(4.4), Inches(6), Inches(0.5),
             "Why the shape matters", size=18, bold=True, color=TEAL)
    add_bullets(s, Inches(0.9), Inches(4.85), Inches(6.0), Inches(2.4),
                [
                    "A chain of N items is searched in O(N).",
                    "A balanced tree of N items is searched in O(log N).",
                    "log₂(1,000,000) ≈ 20 — twenty steps to find one in a million.",
                ],
                size=15, line_spacing=1.3)
    # right side: quick tree
    add_rect(s, Inches(7.5), Inches(1.5), Inches(5.4), Inches(5.3), LIGHT_GREY)
    add_text(s, Inches(7.5), Inches(1.55), Inches(5.4), Inches(0.4),
             "A small tree", size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    cx0 = Inches(10.2); top = Inches(2.3); dx = Inches(1.2); dy = Inches(1.05)
    pos = {
        "10": (cx0,            top),
        "5":  (cx0 - dx,       top + dy),
        "15": (cx0 + dx,       top + dy),
        "2":  (cx0 - dx*1.6,   top + 2*dy),
        "7":  (cx0 - dx*0.4,   top + 2*dy),
        "12": (cx0 + dx*0.5,   top + 2*dy),
        "20": (cx0 + dx*1.6,   top + 2*dy),
    }
    edges = [("10","5"),("10","15"),("5","2"),("5","7"),
             ("15","12"),("15","20")]
    draw_tree(s, pos, edges, font_size=12, node_r=Inches(0.3))
    add_text(s, Inches(7.5), Inches(6.2), Inches(5.4), Inches(0.5),
             "Same data — but ordered for O(log N) search.",
             size=12, italic=True, color=DARK_GREY,
             align=PP_ALIGN.CENTER)


def slide_terminology(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "2.  Tree Terminology",
               "The vocabulary of every tree question you will ever face",
               n, total)
    items = [
        ("Root.  ", "The single topmost node; the entry point."),
        ("Parent / Child.  ", "Direct hierarchical relationship along an edge."),
        ("Sibling.  ", "Two nodes that share a parent."),
        ("Leaf.  ", "A node with no children; the end of a branch."),
        ("Edge.  ", "The connection between a parent and a child."),
        ("Path.  ", "The sequence of edges from one node to another."),
        ("Depth of a node.  ", "Number of edges from the root to that node."),
        ("Height of a tree.  ", "Length of the longest root-to-leaf path."),
        ("Subtree.  ", "Any node together with all of its descendants."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(6.5), Inches(5.5),
                items, size=15, line_spacing=1.25)
    # right diagram with annotations
    add_rect(s, Inches(7.7), Inches(1.5), Inches(5.2), Inches(5.4), LIGHT_GREY)
    cx0 = Inches(10.3); top = Inches(2.0); dx = Inches(0.95); dy = Inches(1.0)
    pos = {
        "A": (cx0,            top),
        "B": (cx0 - dx,       top + dy),
        "C": (cx0 + dx,       top + dy),
        "D": (cx0 - dx*1.6,   top + 2*dy),
        "E": (cx0 - dx*0.4,   top + 2*dy),
        "F": (cx0 + dx,       top + 2*dy),
        "G": (cx0 - dx*0.4,   top + 3*dy),
    }
    edges = [("A","B"),("A","C"),("B","D"),("B","E"),("C","F"),("E","G")]
    draw_tree(s, pos, edges, font_size=12, node_r=Inches(0.3))
    # annotations
    add_text(s, Inches(7.9), Inches(2.0), Inches(1.6), Inches(0.3),
             "← Root", size=11, bold=True, color=RED)
    add_text(s, Inches(11.1), Inches(2.95), Inches(1.5), Inches(0.3),
             "Parent of F", size=10, italic=True, color=DARK_GREY)
    add_text(s, Inches(7.7), Inches(5.0), Inches(2.0), Inches(0.3),
             "Leaf →", size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(7.9), Inches(6.4), Inches(5), Inches(0.3),
             "Height of this tree = 3 (longest path A→B→E→G)",
             size=11, italic=True, color=NAVY)


def slide_binary_tree(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "3.  Binary Trees",
               "Each node has at most two children — left and right",
               n, total)
    add_text(s, Inches(0.9), Inches(1.5), Inches(6.0), Inches(0.5),
             "Definition", size=20, bold=True, color=NAVY)
    add_text(s, Inches(0.9), Inches(2.0), Inches(6.0), Inches(2.0),
             "A binary tree is a tree where every node has zero, one, or two "
             "children, distinguished as left and right.\n\n"
             "Binary trees are the workhorse of CS: BSTs, heaps, expression "
             "trees, Huffman codes, and decision trees are all binary.",
             size=14, color=DARK_GREY)
    add_text(s, Inches(0.9), Inches(4.2), Inches(6.0), Inches(0.5),
             "Special shapes", size=18, bold=True, color=TEAL)
    items = [
        ("Full.  ", "Every node has either 0 or 2 children."),
        ("Complete.  ", "Every level full except possibly the last, filled left-to-right."),
        ("Perfect.  ", "All internal nodes have 2 children, all leaves at the same depth."),
        ("Balanced.  ", "Heights of left and right subtrees differ by at most a constant."),
    ]
    add_bullets(s, Inches(0.9), Inches(4.65), Inches(6.0), Inches(2.5),
                items, size=14, line_spacing=1.25)
    # right: TreeNode code
    add_text(s, Inches(7.4), Inches(1.5), Inches(5.5), Inches(0.5),
             "The atom — TreeNode", size=18, bold=True, color=NAVY)
    code = (
        "class TreeNode:\n"
        "    \"\"\"A single node of a binary tree.\"\"\"\n"
        "\n"
        "    __slots__ = (\"val\", \"left\", \"right\")\n"
        "\n"
        "    def __init__(self, val=0):\n"
        "        self.val   = val\n"
        "        self.left  = None\n"
        "        self.right = None"
    )
    add_code_block(s, Inches(7.4), Inches(2.0), Inches(5.5), Inches(3.4),
                   code, size=14)
    add_text(s, Inches(7.4), Inches(5.5), Inches(5.5), Inches(0.5),
             "Mini binary tree", size=14, bold=True, color=NAVY)
    cx0 = Inches(10.1); top = Inches(6.0); dx = Inches(0.7); dy = Inches(0.6)
    pos = {"R": (cx0, top),
           "L": (cx0-dx, top+dy),
           "Rt":(cx0+dx, top+dy)}
    edges = [("R","L"),("R","Rt")]
    draw_tree(s, pos, edges, font_size=12, node_r=Inches(0.25))


def slide_bst_property(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "4.  Binary Search Trees",
               "The invariant that turns a tree into a sorted index",
               n, total)
    # left: invariant
    add_text(s, Inches(0.9), Inches(1.5), Inches(6.0), Inches(0.5),
             "The BST invariant", size=20, bold=True, color=NAVY)
    add_rect(s, Inches(0.9), Inches(2.05), Inches(6.0), Inches(2.0),
             LIGHT_GREY)
    add_text(s, Inches(1.05), Inches(2.15), Inches(5.7), Inches(1.85),
             "For every node N:\n"
             "   • every value in the left subtree of N  <  N.val\n"
             "   • every value in the right subtree of N >  N.val",
             size=14, color=NAVY, font="Consolas")
    add_text(s, Inches(0.9), Inches(4.25), Inches(6.0), Inches(0.5),
             "Why this matters", size=18, bold=True, color=TEAL)
    items = [
        "Searching halves the candidate set at every step.",
        "An in-order DFS yields the values in sorted order — for free.",
        "Insertion preserves the invariant by always going one side.",
        "When balanced, height = ⌈log₂ N⌉ + 1 — hence O(log N) operations.",
    ]
    add_bullets(s, Inches(0.9), Inches(4.7), Inches(6.0), Inches(2.5),
                items, size=14, line_spacing=1.25)
    # right: example BST
    add_rect(s, Inches(7.4), Inches(1.5), Inches(5.5), Inches(5.4),
             LIGHT_GREY)
    add_text(s, Inches(7.4), Inches(1.55), Inches(5.5), Inches(0.4),
             "Example", size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    cx0 = Inches(10.15); top = Inches(2.3); dx = Inches(1.0); dy = Inches(1.05)
    pos = {
        "10": (cx0,           top),
        "5":  (cx0-dx,        top+dy),
        "15": (cx0+dx,        top+dy),
        "2":  (cx0-dx*1.55,   top+2*dy),
        "7":  (cx0-dx*0.45,   top+2*dy),
        "12": (cx0+dx*0.45,   top+2*dy),
        "20": (cx0+dx*1.55,   top+2*dy),
    }
    edges = [("10","5"),("10","15"),("5","2"),("5","7"),
             ("15","12"),("15","20")]
    draw_tree(s, pos, edges, font_size=12, node_r=Inches(0.32))
    add_text(s, Inches(7.5), Inches(6.0), Inches(5.3), Inches(0.4),
             "Each left subtree is < parent.   Each right subtree is > parent.",
             size=11, italic=True, color=DARK_GREY,
             align=PP_ALIGN.CENTER)


def slide_search_logn(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "5.  Why O(log N)?",
               "Each comparison eliminates half the remaining tree",
               n, total)
    add_text(s, Inches(0.9), Inches(1.5), Inches(11), Inches(0.5),
             "Searching for 7 in the BST below — only 3 comparisons.",
             size=17, color=DARK_GREY)
    # tree with highlighted path
    cx0 = Inches(7.0); top = Inches(2.5); dx = Inches(1.4); dy = Inches(1.2)
    pos = {
        "10": (cx0,           top),
        "5":  (cx0-dx,        top+dy),
        "15": (cx0+dx,        top+dy),
        "2":  (cx0-dx*1.6,    top+2*dy),
        "7":  (cx0-dx*0.4,    top+2*dy),
        "12": (cx0+dx*0.4,    top+2*dy),
        "20": (cx0+dx*1.6,    top+2*dy),
    }
    edges = [("10","5"),("10","15"),("5","2"),("5","7"),
             ("15","12"),("15","20")]
    highlight = {"10","5","7"}
    h_edges = {("10","5"),("5","7")}
    draw_tree(s, pos, edges, font_size=14, node_r=Inches(0.4),
              highlight_nodes=highlight, highlight_edges=h_edges)
    # narration
    add_rect(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.3),
             LIGHT_GREY)
    add_text(s, Inches(1.1), Inches(5.8), Inches(11), Inches(0.5),
             "Step-by-step", size=15, bold=True, color=NAVY)
    add_text(s, Inches(1.1), Inches(6.2), Inches(11), Inches(0.8),
             "1) 7 < 10  →  go left.    "
             "2) 7 > 5   →  go right.    "
             "3) Match: return node 7.   "
             "Three steps, not seven.",
             size=14, color=DARK_GREY)


def slide_insert_walk(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "6.  BST Insertion — Walkthrough",
               "Insert 8 into the tree below",
               n, total)
    # before
    add_text(s, Inches(0.9), Inches(1.5), Inches(11), Inches(0.4),
             "Step 1 — Start at root. 8 < 10  →  go left.",
             size=15, bold=True, color=NAVY)
    cx0 = Inches(3.5); top = Inches(2.4); dx = Inches(1.0); dy = Inches(0.85)
    pos = {
        "10": (cx0,           top),
        "5":  (cx0-dx,        top+dy),
        "15": (cx0+dx,        top+dy),
        "2":  (cx0-dx*1.55,   top+2*dy),
        "7":  (cx0-dx*0.45,   top+2*dy),
        "12": (cx0+dx*0.45,   top+2*dy),
        "20": (cx0+dx*1.55,   top+2*dy),
    }
    edges = [("10","5"),("10","15"),("5","2"),("5","7"),
             ("15","12"),("15","20")]
    draw_tree(s, pos, edges, font_size=12, node_r=Inches(0.3),
              highlight_nodes={"10"})
    # middle
    add_text(s, Inches(0.9), Inches(4.5), Inches(11), Inches(0.4),
             "Step 2 — At 5: 8 > 5  →  go right.   "
             "Step 3 — At 7: 8 > 7 and 7.right is None → insert.",
             size=15, bold=True, color=NAVY)
    # after tree on right
    cx1 = Inches(10.0); top1 = Inches(2.4)
    pos2 = {
        "10": (cx1,            top1),
        "5":  (cx1-dx,         top1+dy),
        "15": (cx1+dx,         top1+dy),
        "2":  (cx1-dx*1.55,    top1+2*dy),
        "7":  (cx1-dx*0.45,    top1+2*dy),
        "12": (cx1+dx*0.45,    top1+2*dy),
        "20": (cx1+dx*1.55,    top1+2*dy),
        "8":  (cx1-dx*0.0,     top1+3*dy),
    }
    edges2 = edges + [("7","8")]
    draw_tree(s, pos2, edges2, font_size=12, node_r=Inches(0.3),
              highlight_nodes={"8"}, highlight_edges={("7","8")})
    add_text(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.5),
             "Each comparison decides one direction. Insertion cost = path length = O(h).",
             size=14, italic=True, color=TEAL)


def slide_insert_code(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "7.  BST Insertion — Code",
               "Iterative insertion is robust and stack-safe",
               n, total)
    code = (
        "class BST:\n"
        "    def __init__(self):\n"
        "        self.root = None\n"
        "\n"
        "    def insert(self, val):\n"
        "        if self.root is None:\n"
        "            self.root = TreeNode(val)\n"
        "            return\n"
        "\n"
        "        current = self.root\n"
        "        while True:\n"
        "            if val < current.val:\n"
        "                if current.left is None:\n"
        "                    current.left = TreeNode(val)\n"
        "                    return\n"
        "                current = current.left\n"
        "            else:\n"
        "                if current.right is None:\n"
        "                    current.right = TreeNode(val)\n"
        "                    return\n"
        "                current = current.right"
    )
    add_code_block(s, Inches(0.9), Inches(1.45), Inches(7.7), Inches(5.6),
                   code, size=13)
    add_text(s, Inches(8.9), Inches(1.45), Inches(4.0), Inches(0.5),
             "Notes", size=18, bold=True, color=NAVY)
    notes = [
        "Iterative form avoids recursion depth limits on large trees.",
        "Empty-tree branch creates the root and returns immediately.",
        "Equal values go right by convention — choose and document.",
        "Time O(h);  for balanced trees O(log N);  worst-case O(N).",
    ]
    add_bullets(s, Inches(8.9), Inches(2.0), Inches(4.1), Inches(4.5),
                notes, size=13, line_spacing=1.3)


def slide_traversal_overview(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "8.  Tree Traversals — Overview",
               "Two strategies, four orders, one tree",
               n, total)
    cards = [
        ("Pre-Order (DFS)", NAVY,
         ["Visit Node → Left → Right.",
          "Use: deep-copy a tree; serialise structure.",
          "Implemented with recursion or an explicit stack."]),
        ("In-Order (DFS)", TEAL,
         ["Visit Left → Node → Right.",
          "On a BST: produces values in sorted ascending order.",
          "Use: validate a BST; output sorted view."]),
        ("Post-Order (DFS)", GOLD,
         ["Visit Left → Right → Node.",
          "Use: free / delete a tree; evaluate expression trees.",
          "Children are processed before their parent."]),
        ("Level-Order (BFS)", PURPLE,
         ["Visit tier by tier from the root downward.",
          "Implemented with a FIFO queue, not recursion.",
          "Use: shortest path in unweighted trees, level summaries."]),
    ]
    qw = Inches(5.95); qh = Inches(2.6); gx = Inches(0.25); gy = Inches(0.2)
    x0 = Inches(0.6); y0 = Inches(1.5)
    for i, (title, color, bullets) in enumerate(cards):
        col = i % 2; row = i // 2
        cx = x0 + (qw + gx) * col
        cy = y0 + (qh + gy) * row
        add_rect(s, cx, cy, qw, Inches(0.55), color)
        add_text(s, cx, cy, qw, Inches(0.55), title,
                 size=17, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, cx, cy + Inches(0.55), qw, qh - Inches(0.55), LIGHT_GREY)
        add_bullets(s, cx + Inches(0.25), cy + Inches(0.7),
                    qw - Inches(0.5), qh - Inches(0.7),
                    bullets, size=13, line_spacing=1.3,
                    bullet_color=color)


def slide_traversal_outputs(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "9.  Same Tree, Four Outputs",
               "What each traversal actually produces", n, total)
    # tree on left
    cx0 = Inches(3.4); top = Inches(2.0); dx = Inches(1.0); dy = Inches(1.0)
    pos = {
        "10": (cx0,            top),
        "5":  (cx0-dx,         top+dy),
        "15": (cx0+dx,         top+dy),
        "2":  (cx0-dx*1.55,    top+2*dy),
        "7":  (cx0-dx*0.45,    top+2*dy),
        "12": (cx0+dx*0.45,    top+2*dy),
        "20": (cx0+dx*1.55,    top+2*dy),
    }
    edges = [("10","5"),("10","15"),("5","2"),("5","7"),
             ("15","12"),("15","20")]
    draw_tree(s, pos, edges, font_size=14, node_r=Inches(0.32))
    # right column: outputs
    add_text(s, Inches(7.0), Inches(1.7), Inches(6), Inches(0.45),
             "Pre-Order  (Node, Left, Right)",
             size=15, bold=True, color=NAVY)
    add_text(s, Inches(7.0), Inches(2.1), Inches(6), Inches(0.45),
             "10 · 5 · 2 · 7 · 15 · 12 · 20",
             size=14, color=DARK_GREY, font="Consolas")
    add_text(s, Inches(7.0), Inches(2.8), Inches(6), Inches(0.45),
             "In-Order  (Left, Node, Right)",
             size=15, bold=True, color=TEAL)
    add_text(s, Inches(7.0), Inches(3.2), Inches(6), Inches(0.45),
             "2 · 5 · 7 · 10 · 12 · 15 · 20    ← sorted!",
             size=14, color=GREEN, font="Consolas")
    add_text(s, Inches(7.0), Inches(3.9), Inches(6), Inches(0.45),
             "Post-Order  (Left, Right, Node)",
             size=15, bold=True, color=GOLD)
    add_text(s, Inches(7.0), Inches(4.3), Inches(6), Inches(0.45),
             "2 · 7 · 5 · 12 · 20 · 15 · 10",
             size=14, color=DARK_GREY, font="Consolas")
    add_text(s, Inches(7.0), Inches(5.0), Inches(6), Inches(0.45),
             "Level-Order  (BFS)",
             size=15, bold=True, color=PURPLE)
    add_text(s, Inches(7.0), Inches(5.4), Inches(6), Inches(0.45),
             "10 · 5 · 15 · 2 · 7 · 12 · 20",
             size=14, color=DARK_GREY, font="Consolas")
    add_rect(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.7),
             LIGHT_GREY)
    add_text(s, Inches(1.1), Inches(6.45), Inches(11), Inches(0.5),
             "In-order on a BST is the cheapest sort you'll ever write.",
             size=14, italic=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)


def slide_traversal_code(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "10.  Traversal Implementations",
               "Three lines change between Pre-, In-, and Post-order",
               n, total)
    code = (
        "def inorder_traversal(root):\n"
        "    res = []\n"
        "    def dfs(node):\n"
        "        if node is None: return\n"
        "        dfs(node.left)\n"
        "        res.append(node.val)     # ← position decides the order\n"
        "        dfs(node.right)\n"
        "    dfs(root)\n"
        "    return res\n"
        "\n"
        "def preorder_traversal(root):\n"
        "    res = []\n"
        "    def dfs(node):\n"
        "        if node is None: return\n"
        "        res.append(node.val)     # node first\n"
        "        dfs(node.left)\n"
        "        dfs(node.right)\n"
        "    dfs(root)\n"
        "    return res"
    )
    add_code_block(s, Inches(0.9), Inches(1.45), Inches(7.7), Inches(5.6),
                   code, size=13)
    # right
    add_text(s, Inches(8.9), Inches(1.45), Inches(4.1), Inches(0.5),
             "Level-order (BFS)", size=17, bold=True, color=NAVY)
    bfs_code = (
        "from collections import deque\n"
        "\n"
        "def level_order(root):\n"
        "    if root is None: return []\n"
        "    res, q = [], deque([root])\n"
        "    while q:\n"
        "        n = q.popleft()\n"
        "        res.append(n.val)\n"
        "        if n.left:  q.append(n.left)\n"
        "        if n.right: q.append(n.right)\n"
        "    return res"
    )
    add_code_block(s, Inches(8.9), Inches(2.0), Inches(4.1), Inches(4.5),
                   bfs_code, size=11)


def slide_lca(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "11.  Lowest Common Ancestor (BST)",
               "The first node where the paths to p and q diverge",
               n, total)
    # tree
    cx0 = Inches(3.4); top = Inches(2.0); dx = Inches(1.0); dy = Inches(1.0)
    pos = {
        "10": (cx0,            top),
        "5":  (cx0-dx,         top+dy),
        "15": (cx0+dx,         top+dy),
        "2":  (cx0-dx*1.55,    top+2*dy),
        "7":  (cx0-dx*0.45,    top+2*dy),
        "12": (cx0+dx*0.45,    top+2*dy),
        "20": (cx0+dx*1.55,    top+2*dy),
    }
    edges = [("10","5"),("10","15"),("5","2"),("5","7"),
             ("15","12"),("15","20")]
    draw_tree(s, pos, edges, font_size=14, node_r=Inches(0.32),
              highlight_nodes={"5"},
              highlight_edges={("10","5")})
    add_text(s, Inches(0.9), Inches(5.4), Inches(6), Inches(0.5),
             "LCA(2, 7) = 5", size=18, bold=True, color=NAVY)
    add_text(s, Inches(0.9), Inches(5.85), Inches(6), Inches(1.0),
             "Both 2 and 7 are smaller than 10 → go left.\n"
             "Now at 5: 2 < 5 and 7 > 5 → split point. Stop.",
             size=14, color=DARK_GREY)
    # code
    add_text(s, Inches(7.2), Inches(1.5), Inches(5.7), Inches(0.5),
             "Iterative LCA (uses BST property)",
             size=16, bold=True, color=NAVY)
    code = (
        "def lowest_common_ancestor(root, p, q):\n"
        "    current = root\n"
        "    while current:\n"
        "        if p.val < current.val and q.val < current.val:\n"
        "            current = current.left\n"
        "        elif p.val > current.val and q.val > current.val:\n"
        "            current = current.right\n"
        "        else:\n"
        "            return current   # split point"
    )
    add_code_block(s, Inches(7.2), Inches(2.05), Inches(5.7), Inches(3.3),
                   code, size=12)
    add_rect(s, Inches(7.2), Inches(5.5), Inches(5.7), Inches(1.5),
             LIGHT_GREY)
    add_text(s, Inches(7.4), Inches(5.6), Inches(5.4), Inches(0.4),
             "Complexity",
             size=14, bold=True, color=NAVY)
    add_text(s, Inches(7.4), Inches(5.95), Inches(5.4), Inches(1.0),
             "Time  O(h)   ·   Space  O(1).\n"
             "On a balanced BST: O(log N) and no recursion stack.",
             size=12, color=DARK_GREY)


def slide_balanced_vs_degenerate(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "12.  Balanced vs Degenerate",
               "Insertion order silently destroys your O(log N) guarantee",
               n, total)
    # left: balanced
    add_text(s, Inches(0.9), Inches(1.5), Inches(5.8), Inches(0.45),
             "Balanced — height ≈ log₂ N",
             size=16, bold=True, color=GREEN)
    cx0 = Inches(3.5); top = Inches(2.2); dx = Inches(0.95); dy = Inches(0.95)
    pos = {
        "10": (cx0,            top),
        "5":  (cx0-dx,         top+dy),
        "15": (cx0+dx,         top+dy),
        "2":  (cx0-dx*1.55,    top+2*dy),
        "7":  (cx0-dx*0.45,    top+2*dy),
        "12": (cx0+dx*0.45,    top+2*dy),
        "20": (cx0+dx*1.55,    top+2*dy),
    }
    edges = [("10","5"),("10","15"),("5","2"),("5","7"),
             ("15","12"),("15","20")]
    draw_tree(s, pos, edges, font_size=12, node_r=Inches(0.3))
    # right: degenerate
    add_text(s, Inches(7.5), Inches(1.5), Inches(5.4), Inches(0.45),
             "Degenerate — sorted insert collapses to a chain",
             size=16, bold=True, color=RED)
    cx2 = Inches(10.4); top2 = Inches(2.2); dx2 = Inches(0.55); dy2 = Inches(0.7)
    pos2 = {
        "1": (cx2,             top2),
        "2": (cx2+dx2,         top2+dy2),
        "3": (cx2+2*dx2,       top2+2*dy2),
        "4": (cx2+3*dx2,       top2+3*dy2),
        "5": (cx2+4*dx2,       top2+4*dy2),
    }
    edges2 = [("1","2"),("2","3"),("3","4"),("4","5")]
    draw_tree(s, pos2, edges2, font_size=12, node_r=Inches(0.25))
    # caption
    add_rect(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.2),
             LIGHT_GREY)
    add_text(s, Inches(1.1), Inches(6.0), Inches(11), Inches(0.5),
             "What went wrong?", size=15, bold=True, color=NAVY)
    add_text(s, Inches(1.1), Inches(6.4), Inches(11), Inches(0.7),
             "Inserting [1,2,3,4,5] into a plain BST yields a right-only "
             "chain. Search becomes O(N). Self-balancing trees (AVL, "
             "Red-Black, B-trees) restore O(log N) by re-shaping after "
             "insertions.",
             size=13, color=DARK_GREY)


def slide_complexity_table(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "13.  Complexity Reference",
               "h = tree height. Balanced ⇒ h = log N. Degenerate ⇒ h = N.",
               n, total)
    rows = [
        ("Operation",      "Avg / Balanced", "Worst / Degenerate"),
        ("Search",         "O(log N)",       "O(N)"),
        ("Insert",         "O(log N)",       "O(N)"),
        ("Delete",         "O(log N)",       "O(N)"),
        ("Min / Max",      "O(log N)",       "O(N)"),
        ("In-order walk",  "O(N)",           "O(N)"),
        ("Space",          "O(N)",           "O(N)"),
        ("Recursion depth","O(log N)",       "O(N)  ← stack-overflow risk"),
    ]
    x0 = Inches(1.6); y0 = Inches(1.65)
    col_w = [Inches(4.5), Inches(3.0), Inches(3.5)]
    rh = Inches(0.55)
    for r, row in enumerate(rows):
        cx = x0
        for c, txt in enumerate(row):
            fill = NAVY if r == 0 else (LIGHT_GREY if r % 2 == 1 else WHITE)
            color = WHITE if r == 0 else DARK_GREY
            bold = (r == 0) or (c == 0)
            add_rect(s, cx, y0 + rh * r, col_w[c], rh, fill,
                     line=RGBColor(0xC0, 0xC4, 0xCC))
            add_text(s, cx + Inches(0.15), y0 + rh * r,
                     col_w[c] - Inches(0.2), rh,
                     txt, size=14, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[c]
    add_text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.4),
             "Self-balancing trees (AVL, Red-Black) guarantee the average column even on adversarial input.",
             size=13, italic=True, color=TEAL)


def slide_pitfalls(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "14.  Common Pitfalls",
               "Where students lose marks — and where production code breaks",
               n, total)
    items = [
        ("Forgetting the empty tree.  ",
         "Calling root.val on None raises AttributeError; always guard root is None first."),
        ("Confusing the BST property.  ",
         "It applies recursively: every value in the left subtree, not just the immediate child."),
        ("Equal values inconsistency.  ",
         "Pick a side (left or right) for duplicates and apply it everywhere; document it."),
        ("Recursive depth on degenerate trees.  ",
         "A chain of 10⁴ nodes can blow Python's default 1000-frame stack. Prefer iteration."),
        ("Mutating during traversal.  ",
         "Inserting or deleting while walking the tree corrupts the iteration. Snapshot first."),
        ("Wrong traversal for the job.  ",
         "Use post-order to delete; in-order to print sorted; pre-order to copy structure."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.5), Inches(12), Inches(5.5),
                items, size=15, line_spacing=1.3, bullet_color=RED)


def slide_applications(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "15.  Trees in the Real World",
               "Where the structures from this lecture run in production",
               n, total)
    cards = [
        ("Filesystems",      NAVY,
         "Directories form a tree; ext4 and NTFS use B-trees for indexing."),
        ("Databases",        TEAL,
         "Indexes are B+ trees — every SQL WHERE you've ever written touches one."),
        ("HTML / DOM",       GOLD,
         "Browsers parse pages into a tree the JS engine traverses to render."),
        ("Compilers",        RED,
         "Source code becomes an Abstract Syntax Tree before optimisation and codegen."),
        ("Game AI",          GREEN,
         "Decision trees and Monte-Carlo Tree Search drive chess, Go, and modern game AI."),
        ("Networking",       PURPLE,
         "IP routing tables use radix / trie trees for longest-prefix match."),
    ]
    cw = Inches(4.0); ch = Inches(2.4); gx = Inches(0.2); gy = Inches(0.2)
    x0 = Inches(0.6); y0 = Inches(1.55)
    for i, (title, color, body) in enumerate(cards):
        col = i % 3; row = i // 3
        cx = x0 + (cw + gx) * col
        cy = y0 + (ch + gy) * row
        add_rect(s, cx, cy, cw, Inches(0.55), color)
        add_text(s, cx, cy, cw, Inches(0.55), title,
                 size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, cx, cy + Inches(0.55), cw, ch - Inches(0.55), LIGHT_GREY)
        add_text(s, cx + Inches(0.2), cy + Inches(0.7),
                 cw - Inches(0.4), ch - Inches(0.7),
                 body, size=14, color=DARK_GREY)


def slide_lab(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "16.  Practical Lab — practice.py",
               "Three exercises. Build, traverse, reason.",
               n, total)
    items = [
        ("Exercise 1.  ",
         "Implement BST.insert() iteratively, preserving the BST invariant."),
        ("Exercise 2.  ",
         "Implement inorder_traversal() and preorder_traversal() recursively."),
        ("Exercise 3.  ",
         "Implement lowest_common_ancestor(root, p, q) using the BST property in O(h)."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.55), Inches(12), Inches(2.5),
                items, size=18, line_spacing=1.4)
    add_text(s, Inches(0.9), Inches(4.2), Inches(12), Inches(0.5),
             "Submission expectations",
             size=18, bold=True, color=NAVY)
    rules = [
        "Test the empty tree, single-node tree, and a degenerate (sorted-insert) tree.",
        "Verify that inorder_traversal returns a sorted list.",
        "Annotate each function with its time and space complexity.",
        "Avoid recursion deeper than necessary; prefer iteration where natural.",
    ]
    add_bullets(s, Inches(0.9), Inches(4.65), Inches(12), Inches(2.5),
                rules, size=15, line_spacing=1.3, bullet_color=GOLD)


def slide_summary(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "17.  Summary",
               "Six sentences worth memorising", n, total)
    items = [
        "A tree is a hierarchy of nodes connected by edges; one root, no cycles.",
        "A binary tree has at most two children per node — left and right.",
        "A BST imposes left < parent < right recursively, enabling O(log N) search on balanced inputs.",
        "Pre-order, In-order, and Post-order are DFS traversals distinguished by where the node visit happens relative to the recursive calls.",
        "Level-order is BFS over a queue — used for tier-by-tier algorithms.",
        "Insertion order matters: a sorted insert produces a degenerate tree; balancing trees fix this.",
    ]
    add_bullets(s, Inches(0.9), Inches(1.55), Inches(12), Inches(5.5),
                items, size=17, line_spacing=1.4, bullet_color=GOLD)


def slide_references(n, total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "18.  Further Reading",
               "Trusted sources for deeper study", n, total)
    refs = [
        ("CLRS — Introduction to Algorithms (4th ed.).  ",
         "Cormen, Leiserson, Rivest, Stein. Chapter 12 — Binary Search Trees."),
        ("The Art of Computer Programming, Vol. 3.  ",
         "Donald E. Knuth. Section 6.2 — Searching by Tree Insertion."),
        ("Algorithms (4th ed.).  ",
         "Sedgewick & Wayne. Princeton University. Chapter 3.2 — BSTs."),
        ("Adelson-Velsky & Landis (1962).  ",
         "An algorithm for the organization of information. Origin of AVL trees."),
        ("LeetCode practice set.  ",
         "94 (Inorder), 144 (Preorder), 102 (Level-order), 235 (LCA-BST), 98 (Validate BST)."),
    ]
    add_bullets(s, Inches(0.9), Inches(1.55), Inches(12), Inches(5.5),
                refs, size=16, line_spacing=1.4)


def slide_thanks():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, 0, Inches(5.5), SLIDE_W, Inches(0.12), GOLD)
    add_text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.5),
             "Questions?", size=80, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.7),
             "Bring your traces, your edge cases, and your degenerate trees.",
             size=22, italic=True, color=GOLD)
    add_text(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.4),
             "Next week  ·  Heaps & Priority Queues — the BST's specialised cousin.",
             size=16, color=LIGHT_GREY)
    add_text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4),
             "Data Structures & Algorithms  ·  Week 06  ·  Trees & BSTs",
             size=12, color=LIGHT_GREY, italic=True)


# -------- Build ----------
slide_title()
content = [
    slide_outline,
    slide_objectives,
    slide_motivation,
    slide_terminology,
    slide_binary_tree,
    slide_bst_property,
    slide_search_logn,
    slide_insert_walk,
    slide_insert_code,
    slide_traversal_overview,
    slide_traversal_outputs,
    slide_traversal_code,
    slide_lca,
    slide_balanced_vs_degenerate,
    slide_complexity_table,
    slide_pitfalls,
    slide_applications,
    slide_lab,
    slide_summary,
    slide_references,
]
TOTAL = 1 + len(content) + 1
for i, fn in enumerate(content, start=2):
    fn(i, TOTAL)
slide_thanks()

out = "/Users/silasgah/Downloads/enquire.ai/ds/week_06_trees/trees_dsa_final.pptx"
prs.save(out)
print(f"Wrote {out}  ({len(prs.slides)} slides)")
